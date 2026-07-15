#!/usr/bin/env python3
"""
Extract one slide from a source deck into a single-slide asset.pptx.

Usage:
    python scripts/extract_slide.py <source.pptx> <slide_index> <target_dir>
"""

from __future__ import annotations

import shutil
import sys
from pathlib import Path

from pptx import Presentation


def _delete_slide(prs, index: int) -> None:
    slide_id = prs.slides._sldIdLst[index]
    rId = slide_id.rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[index]


def extract_slide(source_pptx: str, slide_index: int, target_dir: str) -> Path:
    source_path = Path(source_pptx)
    target_path = Path(target_dir)
    target_path.mkdir(parents=True, exist_ok=True)

    prs = Presentation(str(source_path))
    if not (0 <= slide_index < len(prs.slides)):
        raise ValueError(
            f"Slide index {slide_index} out of range for {source_path} ({len(prs.slides)} slides)"
        )

    # Remove all slides except the desired one, from the end toward the front
    # so indices stay stable.
    for i in range(len(prs.slides) - 1, -1, -1):
        if i != slide_index:
            _delete_slide(prs, i)

    out_pptx = target_path / "asset.pptx"
    prs.save(str(out_pptx))
    return out_pptx


def main() -> None:
    if len(sys.argv) != 4:
        print(__doc__)
        sys.exit(1)
    source_pptx, slide_index, target_dir = sys.argv[1:4]
    out = extract_slide(source_pptx, int(slide_index), target_dir)
    print(f"Saved single-slide asset to {out}")


if __name__ == "__main__":
    main()
