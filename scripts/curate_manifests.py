#!/usr/bin/env python3
"""
Curate the seven newly ingested assets and clear placeholder text from
unbound shapes so certification passes without leakage errors.
"""

from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


_PLACEHOLDER_LEAKAGE_VALUES = {"text", "title", "subtitle", "placeholder", "lorem ipsum", "xx", "…"}


def _load(path: Path) -> dict:
    return json.loads(path.read_text(encoding="utf-8"))


def _save(path: Path, data: dict) -> None:
    path.write_text(json.dumps(data, indent=2), encoding="utf-8")


def _ph(id: str, role: str, kind: str, cardinality: str, shape_name: str | None,
        native_idx: int | None, required: bool, max_chars: int, max_lines: int) -> dict:
    return {
        "id": id,
        "role": role,
        "kind": kind,
        "cardinality": cardinality,
        "required": required,
        "content_schema": {},
        "constraints": {"max_chars": max_chars, "max_lines": max_lines},
        "binding": {
            "native_placeholder_idx": native_idx,
            "shape_name": shape_name,
        },
    }


def _iter_leaf_shapes(shapes):
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_leaf_shapes(shape.shapes)
        else:
            yield shape


def _clear_unbound_leakage(pptx_path: Path, bound_names: set[str]) -> None:
    prs = Presentation(str(pptx_path))
    slide = prs.slides[0]
    for shape in _iter_leaf_shapes(slide.shapes):
        if not getattr(shape, "has_text_frame", False):
            continue
        name = getattr(shape, "name", "")
        if name in bound_names:
            continue
        text = " ".join((shape.text or "").strip().lower().split())
        if text in _PLACEHOLDER_LEAKAGE_VALUES:
            shape.text_frame.clear()
    prs.save(str(pptx_path))


def _bound_names_from_manifest(manifest: dict) -> set[str]:
    names: set[str] = set()
    repeating = manifest.get("repeating")
    count = repeating["count"] if repeating else None
    for p in manifest["placeholders"]:
        shape_name = p["binding"]["shape_name"]
        if not shape_name:
            continue
        if p["cardinality"] == "N" and repeating and count:
            token = repeating.get("index_token", "{N}")
            names.update(shape_name.replace(token, str(i)) for i in range(1, count + 1))
        else:
            names.add(shape_name)
    return names


def _update_meta(m: dict, **kwargs) -> None:
    m.update(kwargs)


def _curate_benefits():
    path = Path("presentation_assets/business_benefits/BENEFITS-6FACTOR-001/asset.json")
    m = _load(path)
    _update_meta(
        m,
        purpose="Six-factor benefits layout with a label and short description per factor.",
        audience_tags=["board", "executive", "leadership", "program_team"],
        style_tags=["structured", "list", "board-ready", "clean"],
        recommended_for=["business benefits", "value drivers", "success factors", "benefits case"],
        avoid_for=["data-heavy dashboards", "comparisons", "process flows", "risk assessment"],
        density=6,
        density_range=[4, 6],
        fits_content_kinds=["business_benefits", "benefits", "value_drivers", "success_factors"],
        message_type="benefits",
        information_shape="list",
    )
    m["placeholders"] = [
        _ph("title", "title", "title", "1", None, 0, True, 120, 2),
        _ph("factor_label", "factor_name", "body", "N", "Factor{N}Label", None, False, 40, 1),
        _ph("factor_body", "factor_description", "body", "N", "Factor{N}Body", None, False, 200, 3),
    ]
    m["repeating"] = {
        "group_template": "Factor{N}Group",
        "placeholders_per_group": ["factor_label", "factor_body"],
        "index_token": "{N}",
        "count": 6,
    }
    _save(path, m)
    _clear_unbound_leakage(path.parent / "asset.pptx", _bound_names_from_manifest(m))
    print("[benefits] curated")


def _curate_risk():
    path = Path("presentation_assets/risk/RISK-ASSESSMENT-001/asset.json")
    m = _load(path)
    _update_meta(
        m,
        purpose="Risk assessment matrix with risk title, description, assessment and confidence per row.",
        audience_tags=["board", "executive", "leadership", "program_team", "risk_committee"],
        style_tags=["structured", "matrix", "board-ready", "analytical"],
        recommended_for=["risk assessment", "risk matrix", "implementation risks", "mitigation"],
        avoid_for=["simple lists", "timeline", "process flows", "kpi dashboards"],
        density=4,
        density_range=[3, 4],
        fits_content_kinds=["risk", "assessment", "matrix"],
        message_type="risk",
        information_shape="matrix",
    )
    m["placeholders"] = [
        _ph("title", "title", "title", "1", None, 0, True, 120, 2),
        _ph("header_findings", "header_findings", "body", "1", "HeaderFindings", None, True, 35, 1),
        _ph("header_assessment", "header_assessment", "body", "1", "HeaderAssessment", None, True, 35, 1),
        _ph("header_confidence", "header_confidence", "body", "1", "HeaderConfidence", None, True, 35, 1),
        _ph("risk_title", "risk_title", "body", "N", "Risk{N}Title", None, False, 50, 2),
        _ph("risk_description", "risk_description", "body", "N", "Risk{N}Description", None, False, 220, 3),
        _ph("risk_assessment", "risk_assessment", "body", "N", "Risk{N}Assessment", None, False, 30, 1),
        _ph("risk_confidence", "risk_confidence", "body", "N", "Risk{N}Confidence", None, False, 30, 1),
    ]
    m["repeating"] = {
        "group_template": "Risk{N}Group",
        "placeholders_per_group": ["risk_title", "risk_description", "risk_assessment", "risk_confidence"],
        "index_token": "{N}",
        "count": 4,
    }
    _save(path, m)
    _clear_unbound_leakage(path.parent / "asset.pptx", _bound_names_from_manifest(m))
    print("[risk] curated")


def _curate_opportunity():
    path = Path("presentation_assets/opportunity_matrix/OPPORTUNITY-MATRIX-3SEGMENT-001/asset.json")
    m = _load(path)
    _update_meta(
        m,
        purpose="Three-segment opportunity matrix with category panels and supporting item lines.",
        audience_tags=["board", "executive", "leadership", "strategy_team"],
        style_tags=["structured", "matrix", "strategic", "board-ready"],
        recommended_for=["opportunity matrix", "market opportunity", "growth levers", "strategic opportunities"],
        avoid_for=["risk assessment", "process flows", "kpi dashboards", "timeline"],
        density=20,
        density_range=[8, 20],
        fits_content_kinds=["opportunity_matrix", "matrix", "opportunities", "growth"],
        message_type="opportunity",
        information_shape="matrix",
    )
    placeholders = [
        _ph("title", "title", "title", "1", None, 0, True, 120, 2),
        _ph("category_1", "category", "body", "1", "Category1", None, True, 120, 3),
        _ph("category_2", "category", "body", "1", "Category2", None, True, 120, 3),
        _ph("category_3", "category", "body", "1", "Category3", None, True, 120, 3),
        _ph("category_indicator_1", "category_indicator", "body", "1", "Category1Indicator", None, True, 5, 1),
        _ph("category_indicator_2", "category_indicator", "body", "1", "Category2Indicator", None, True, 5, 1),
        _ph("category_indicator_3", "category_indicator", "body", "1", "Category3Indicator", None, True, 5, 1),
    ]
    for i in range(1, 13):
        placeholders.append(_ph(f"item_{i}_label", "item_label", "body", "1", f"Item{i}Label", None, False, 80, 2))
    for i in range(1, 21):
        placeholders.append(_ph(f"item_{i}_body", "item_body", "body", "1", f"Item{i}Body", None, False, 120, 2))
    m["placeholders"] = placeholders
    m["repeating"] = None
    _save(path, m)
    _clear_unbound_leakage(path.parent / "asset.pptx", _bound_names_from_manifest(m))
    print("[opportunity] curated")


def _curate_roadmap():
    path = Path("presentation_assets/roadmap/ROADMAP-3PHASE-4WORKSTREAM-001/asset.json")
    m = _load(path)
    _update_meta(
        m,
        purpose="Three-phase roadmap with four workstreams, durations and phase results.",
        audience_tags=["board", "executive", "leadership", "program_team"],
        style_tags=["structured", "roadmap", "phased", "board-ready"],
        recommended_for=["roadmap", "implementation roadmap", "phase plan", "workstream plan"],
        avoid_for=["simple lists", "data-heavy dashboards", "risk matrix", "process flows"],
        density=5,
        density_range=[3, 5],
        fits_content_kinds=["roadmap", "phases", "workstreams", "milestones"],
        message_type="roadmap",
        information_shape="sequence",
    )
    placeholders = [
        _ph("title", "title", "title", "1", None, 0, True, 120, 2),
        _ph("phase_1_title", "phase_title", "body", "1", "Phase1Title", None, True, 90, 2),
        _ph("phase_2_title", "phase_title", "body", "1", "Phase2Title", None, True, 90, 2),
        _ph("phase_3_title", "phase_title", "body", "1", "Phase3Title", None, True, 90, 2),
        _ph("phase_1_roman", "phase_roman", "body", "1", "Phase1Roman", None, True, 10, 1),
        _ph("phase_2_roman", "phase_roman", "body", "1", "Phase2Roman", None, True, 10, 1),
        _ph("phase_3_roman", "phase_roman", "body", "1", "Phase3Roman", None, True, 10, 1),
        _ph("phase_1_duration", "phase_duration", "body", "1", "Phase1Duration", None, False, 30, 1),
        _ph("phase_2_duration", "phase_duration", "body", "1", "Phase2Duration", None, False, 30, 1),
        _ph("phase_1_result", "phase_result", "body", "1", "Phase1Result", None, False, 40, 1),
        _ph("phase_2_result", "phase_result", "body", "1", "Phase2Result", None, False, 40, 1),
        _ph("phase_3_result", "phase_result", "body", "1", "Phase3Result", None, False, 40, 1),
    ]
    for i in range(1, 6):
        placeholders.append(_ph(f"phase_block_{i}", "phase_block", "body", "1", f"PhaseBlock{i}", None, False, 300, 6))
    for i in range(1, 5):
        placeholders.append(_ph(f"workstream_{i}_label", "workstream_label", "body", "1", f"Workstream{i}Label", None, True, 60, 2))
        placeholders.append(_ph(f"workstream_{i}_number", "workstream_number", "metric", "1", f"Workstream{i}Number", None, True, 10, 1))
    m["placeholders"] = placeholders
    m["repeating"] = None
    _save(path, m)
    _clear_unbound_leakage(path.parent / "asset.pptx", _bound_names_from_manifest(m))
    print("[roadmap] curated")


def _curate_process():
    path = Path("presentation_assets/process/PROCESS-MATURITY-SLIDER-001/asset.json")
    m = _load(path)
    _update_meta(
        m,
        purpose="Process maturity slider showing current state, ambition and capability gaps.",
        audience_tags=["board", "executive", "leadership", "operations", "program_team"],
        style_tags=["structured", "process", "maturity", "analytical"],
        recommended_for=["process maturity", "capability assessment", "as-is to-be", "maturity model"],
        avoid_for=["simple lists", "kpi dashboards", "roadmap", "risk matrix"],
        density=10,
        density_range=[4, 10],
        fits_content_kinds=["process", "maturity", "capability", "assessment"],
        message_type="process",
        information_shape="comparison",
    )
    placeholders = [
        _ph("title", "title", "title", "1", None, 0, True, 120, 2),
    ]
    for i in range(1, 5):
        placeholders.append(_ph(f"top_block_{i}", "top_block", "body", "1", f"TopBlock{i}", None, False, 250, 4))
    for i in range(1, 3):
        placeholders.append(_ph(f"bottom_block_{i}", "bottom_block", "body", "1", f"BottomBlock{i}", None, False, 250, 4))
    for i in range(1, 11):
        placeholders.append(_ph(f"slider_{i}", "slider", "body", "1", f"Slider{i}", None, False, 60, 2))
    for i in range(1, 3):
        placeholders.append(_ph(f"timeline_{i}", "timeline", "body", "1", f"Timeline{i}", None, False, 40, 1))
    m["placeholders"] = placeholders
    m["repeating"] = None
    _save(path, m)
    _clear_unbound_leakage(path.parent / "asset.pptx", _bound_names_from_manifest(m))
    print("[process] curated")


def _curate_nextsteps():
    path = Path("presentation_assets/next_steps/NEXTSTEPS-REGISTER-001/asset.json")
    m = _load(path)
    _update_meta(
        m,
        purpose="Next steps register with priority, instrument, action, timing and owner.",
        audience_tags=["board", "executive", "leadership", "program_team", "steering_committee"],
        style_tags=["structured", "register", "action-oriented", "board-ready"],
        recommended_for=["next steps", "action register", "initiative tracker", "implementation plan"],
        avoid_for=["data-heavy dashboards", "process flows", "risk matrix", "kpi dashboards"],
        density=8,
        density_range=[4, 8],
        fits_content_kinds=["next_steps", "actions", "register", "tracker"],
        message_type="next_steps",
        information_shape="list",
    )
    placeholders = [
        _ph("title", "title", "title", "1", None, 0, True, 120, 2),
        _ph("header_nr", "header_nr", "body", "1", "HeaderNr", None, True, 15, 1),
        _ph("header_priority", "header_priority", "body", "1", "HeaderPriority", None, True, 25, 1),
        _ph("header_instrument", "header_instrument", "body", "1", "HeaderInstrument", None, True, 25, 1),
        _ph("header_next_step", "header_next_step", "body", "1", "HeaderNextStep", None, True, 25, 1),
        _ph("header_when", "header_when", "body", "1", "HeaderWhen", None, True, 25, 1),
        _ph("header_who", "header_who", "body", "1", "HeaderWho", None, True, 25, 1),
    ]
    for i in range(1, 5):
        placeholders.append(_ph(f"row_{i}_nr", "row_nr", "metric", "1", f"Row{i}Nr", None, True, 10, 1))
        placeholders.append(_ph(f"row_{i}_priority", "row_priority", "body", "1", f"Row{i}Priority", None, True, 60, 3))
        placeholders.append(_ph(f"row_{i}_instrument", "row_instrument", "body", "1", f"Row{i}Instrument", None, True, 80, 3))
    for i in range(1, 9):
        placeholders.append(_ph(f"row_{i}_next_step", "row_next_step", "body", "1", f"Row{i}NextStep", None, False, 160, 3))
        placeholders.append(_ph(f"row_{i}_when", "row_when", "body", "1", f"Row{i}When", None, False, 60, 2))
        placeholders.append(_ph(f"row_{i}_who", "row_who", "body", "1", f"Row{i}Who", None, False, 60, 2))
    m["placeholders"] = placeholders
    m["repeating"] = None
    _save(path, m)
    _clear_unbound_leakage(path.parent / "asset.pptx", _bound_names_from_manifest(m))
    print("[next_steps] curated")


def _curate_kpi():
    path = Path("presentation_assets/kpi/KPI-BAR-INDICATOR-001/asset.json")
    m = _load(path)
    _update_meta(
        m,
        purpose="Bar-indicator KPI slide with categories, values and supporting labels.",
        audience_tags=["board", "executive", "leadership", "operations"],
        style_tags=["data-driven", "kpi", "indicator", "board-ready"],
        recommended_for=["kpi", "metrics", "performance indicators", "dashboard"],
        avoid_for=["narrative lists", "process flows", "risk matrix", "roadmap"],
        density=6,
        density_range=[4, 6],
        fits_content_kinds=["kpi", "metrics", "indicators"],
        message_type="kpi",
        information_shape="comparison",
    )
    placeholders = [
        _ph("title", "title", "title", "1", None, 0, True, 120, 2),
        _ph("category", "kpi_category", "body", "N", "Category{N}", None, True, 40, 1),
        _ph("value", "kpi_value", "metric", "N", "Value{N}", None, True, 30, 1),
        _ph("right_label", "kpi_label", "body", "N", "RightLabel{N}", None, False, 30, 1),
    ]
    m["placeholders"] = placeholders
    m["repeating"] = {
        "group_template": "KPI{N}Group",
        "placeholders_per_group": ["category", "value", "right_label"],
        "index_token": "{N}",
        "count": 6,
    }
    _save(path, m)
    _clear_unbound_leakage(path.parent / "asset.pptx", _bound_names_from_manifest(m))
    print("[kpi] curated")


def main():
    _curate_benefits()
    _curate_risk()
    _curate_opportunity()
    _curate_roadmap()
    _curate_process()
    _curate_nextsteps()
    _curate_kpi()


if __name__ == "__main__":
    main()
