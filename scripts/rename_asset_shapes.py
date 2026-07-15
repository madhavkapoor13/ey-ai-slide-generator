#!/usr/bin/env python3
"""
Rename content-bearing shapes in the newly extracted assets so each has a
unique, meaningful name. Decorative shapes are left untouched.

Run after extraction and before re-ingesting:
    . venv/bin/activate && python scripts/rename_asset_shapes.py
"""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

_PROJECT_ROOT = Path(__file__).resolve().parents[1]
if str(_PROJECT_ROOT) not in sys.path:
    sys.path.insert(0, str(_PROJECT_ROOT))


def _to_inches(emu: int) -> float:
    return emu / 914400


def _collect_leaf_shapes(slide):
    """Return all leaf shapes recursively."""
    out = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            out.extend(_collect_leaf_shapes(shape))
        else:
            out.append(shape)
    return out


def _is_native_title(shape) -> bool:
    return bool(getattr(shape, "is_placeholder", False) and getattr(shape, "placeholder_format", None) is not None and shape.placeholder_format.idx == 0)


def _text(shape) -> str:
    return shape.text_frame.text if shape.has_text_frame else ""


def _rename_benefits():
    pptx = Path("presentation_assets/business_benefits/BENEFITS-6FACTOR-001/asset.pptx")
    prs = Presentation(str(pptx))
    shapes = _collect_leaf_shapes(prs.slides[0])

    # Labels: left column AUTO_SHAPE, x<2, width 2.5-3.0, has text, not placeholder
    labels = [s for s in shapes
              if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
              and not _is_native_title(s)
              and _to_inches(s.left) < 1.5
              and 2.5 <= _to_inches(s.width) <= 3.5
              and s.has_text_frame]
    labels.sort(key=lambda s: s.top)
    for i, s in enumerate(labels, 1):
        s.name = f"Factor{i}Label"

    # Bodies: right column AUTO_SHAPE, x≈3.3-3.7, width 8.5-9.7
    bodies = [s for s in shapes
              if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
              and not _is_native_title(s)
              and 3.0 <= _to_inches(s.left) <= 4.0
              and 8.0 <= _to_inches(s.width) <= 10.0
              and s.has_text_frame]
    bodies.sort(key=lambda s: s.top)
    for i, s in enumerate(bodies, 1):
        s.name = f"Factor{i}Body"

    prs.save(str(pptx))
    print(f"[benefits] renamed {len(labels)} labels and {len(bodies)} bodies")


def _rename_risk():
    pptx = Path("presentation_assets/risk/RISK-ASSESSMENT-001/asset.pptx")
    prs = Presentation(str(pptx))
    shapes = _collect_leaf_shapes(prs.slides[0])

    # Header row at y≈1.45
    headers = [s for s in shapes
               if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
               and not _is_native_title(s)
               and 1.35 <= _to_inches(s.top) <= 1.6
               and s.has_text_frame
               and _text(s).strip()]
    headers.sort(key=lambda s: s.left)
    names = ["HeaderFindings", "HeaderAssessment", "HeaderConfidence"]
    for s, name in zip(headers, names):
        s.name = name

    # Row y centers
    row_ys = [1.97, 3.17, 4.37, 5.57]
    for i, y in enumerate(row_ys, 1):
        row = [s for s in shapes
               if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
               and not _is_native_title(s)
               and abs(_to_inches(s.top) - y) < 0.2
               and s.has_text_frame]
        row.sort(key=lambda s: s.left)
        # Expected order: title (x~0.66, width~1.64), description (x~2.37, width~7.71),
        # assessment (x~10.15, width~1.22), confidence (x~11.44, width~1.22).
        # Ignore small indicator shapes (width < 0.8).
        mapping = {}
        for s in row:
            x = _to_inches(s.left)
            w = _to_inches(s.width)
            if x < 1.5 and w > 1.0:
                mapping.setdefault("title", []).append(s)
            elif x < 5 and w > 3.0:
                mapping.setdefault("description", []).append(s)
            elif 9.5 <= x < 10.8 and w > 0.8:
                mapping.setdefault("assessment", []).append(s)
            elif x >= 10.8 and w > 0.8:
                mapping.setdefault("confidence", []).append(s)
        for key, name in [("title", f"Risk{i}Title"), ("description", f"Risk{i}Description"),
                          ("assessment", f"Risk{i}Assessment"), ("confidence", f"Risk{i}Confidence")]:
            for s in mapping.get(key, []):
                s.name = name

    # Legend footnotes at bottom y≈6.8-7.1, TEXT_BOX
    footnotes = [s for s in shapes
                 if s.shape_type == MSO_SHAPE_TYPE.TEXT_BOX
                 and _to_inches(s.top) > 6.7
                 and _text(s).strip()]
    footnotes.sort(key=lambda s: s.left)
    for i, s in enumerate(footnotes, 1):
        s.name = f"Legend{i}"

    prs.save(str(pptx))
    print(f"[risk] renamed headers/rows/legend")


def _rename_opportunity():
    pptx = Path("presentation_assets/opportunity_matrix/OPPORTUNITY-MATRIX-3SEGMENT-001/asset.pptx")
    prs = Presentation(str(pptx))
    shapes = _collect_leaf_shapes(prs.slides[0])

    # Left big category boxes
    cats = [s for s in shapes
            if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
            and not _is_native_title(s)
            and _to_inches(s.left) < 2.0
            and _to_inches(s.width) > 1.5
            and _text(s).strip()]
    cats.sort(key=lambda s: s.top)
    for i, s in enumerate(cats, 1):
        s.name = f"Category{i}"

    # Category indicators A/B/C (small squares at x≈0.74)
    indicators = [s for s in shapes
                  if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
                  and _to_inches(s.left) < 1.5
                  and _text(s).strip() in ("A", "B", "C")]
    indicators.sort(key=lambda s: s.top)
    for i, s in enumerate(indicators, 1):
        s.name = f"Category{i}Indicator"

    # Right-side item labels at x≈3.04 (not the category column)
    labels = [s for s in shapes
              if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
              and 2.5 <= _to_inches(s.left) <= 4.0
              and _text(s).strip()
              and _text(s).strip() != "Text"]
    labels.sort(key=lambda s: s.top)
    for i, s in enumerate(labels, 1):
        s.name = f"Item{i}Label"

    # Right-side item bodies at x≈5.63 (wide text placeholders under each label)
    bodies = [s for s in shapes
              if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
              and 5.0 <= _to_inches(s.left) <= 6.5
              and _to_inches(s.width) > 5.0]
    bodies.sort(key=lambda s: s.top)
    for i, s in enumerate(bodies, 1):
        s.name = f"Item{i}Body"

    # Numbered ellipses (small circles with digits)
    ellipses = [s for s in shapes
                if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
                and 5.5 <= _to_inches(s.left) <= 6.5
                and 0.15 <= _to_inches(s.width) <= 0.3
                and _text(s).strip().isdigit()]
    ellipses.sort(key=lambda s: s.top)
    for i, s in enumerate(ellipses, 1):
        s.name = f"Item{i}Number"

    prs.save(str(pptx))
    print(f"[opportunity] renamed {len(cats)} categories, {len(labels)} labels")


def _rename_roadmap():
    pptx = Path("presentation_assets/roadmap/ROADMAP-3PHASE-4WORKSTREAM-001/asset.pptx")
    prs = Presentation(str(pptx))
    shapes = _collect_leaf_shapes(prs.slides[0])

    # Phase headers at top y≈1.76, x>2, wide-ish
    phases = [s for s in shapes
              if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
              and not _is_native_title(s)
              and 1.6 <= _to_inches(s.top) <= 2.0
              and _to_inches(s.left) > 2.0
              and _to_inches(s.width) > 2.0
              and _text(s).strip()]
    phases.sort(key=lambda s: s.left)
    for i, s in enumerate(phases, 1):
        s.name = f"Phase{i}Title"

    # Phase Roman numerals I/II/III (small circles at x≈2.45, 5.95, 10.05)
    romans = [s for s in shapes
              if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
              and _text(s).strip() in ("I", "II", "III")]
    romans.sort(key=lambda s: s.left)
    for i, s in enumerate(romans, 1):
        s.name = f"Phase{i}Roman"

    # Durations (~x weeks)
    durations = [s for s in shapes if s.has_text_frame and "week" in _text(s).lower()]
    durations.sort(key=lambda s: s.left)
    for i, s in enumerate(durations, 1):
        s.name = f"Phase{i}Duration"

    # Phase body blocks (large rectangles under phases, text "xx" or empty)
    bodies = [s for s in shapes
              if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
              and not _is_native_title(s)
              and _to_inches(s.top) > 2.0
              and _to_inches(s.width) > 2.0
              and _to_inches(s.height) > 0.5
              and _to_inches(s.left) > 2.0]
    bodies.sort(key=lambda s: (s.left, s.top))
    for i, s in enumerate(bodies, 1):
        s.name = f"PhaseBlock{i}"

    # Workstream labels on far left (x<1.7, y>2, not roman numerals)
    streams = [s for s in shapes
               if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
               and not _is_native_title(s)
               and _to_inches(s.left) < 1.7
               and _to_inches(s.top) > 2.0
               and _text(s).strip()
               and _text(s).strip() not in ("I", "II", "III", "1", "2", "3", "4")]
    streams.sort(key=lambda s: s.top)
    for i, s in enumerate(streams, 1):
        s.name = f"Workstream{i}Label"

    # Bubble numbers 1-4
    bubbles = [s for s in shapes
               if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
               and _to_inches(s.left) < 1.7
               and _text(s).strip() in ("1", "2", "3", "4")]
    bubbles.sort(key=lambda s: s.top)
    for i, s in enumerate(bubbles, 1):
        s.name = f"Workstream{i}Number"

    # Result labels at bottom
    results = [s for s in shapes
               if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
               and _text(s).strip() == "Result"]
    results.sort(key=lambda s: s.left)
    for i, s in enumerate(results, 1):
        s.name = f"Phase{i}Result"

    prs.save(str(pptx))
    print(f"[roadmap] renamed phases, durations, blocks, workstreams")


def _rename_process():
    pptx = Path("presentation_assets/process/PROCESS-MATURITY-SLIDER-001/asset.pptx")
    prs = Presentation(str(pptx))
    shapes = _collect_leaf_shapes(prs.slides[0])

    # Title stays native.
    # Top-left cluster: text blocks around x 1.4-4.5, y 1.5-2.2
    top = [s for s in shapes
           if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
           and not _is_native_title(s)
           and 1.4 <= _to_inches(s.left) <= 5.0
           and 1.4 <= _to_inches(s.top) <= 2.2
           and s.has_text_frame]
    top.sort(key=lambda s: s.top)
    for i, s in enumerate(top, 1):
        s.name = f"TopBlock{i}"

    # Bottom-left cluster: y 3.0-4.0
    bottom = [s for s in shapes
              if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
              and not _is_native_title(s)
              and 1.4 <= _to_inches(s.left) <= 5.0
              and 3.0 <= _to_inches(s.top) <= 4.2
              and s.has_text_frame]
    bottom.sort(key=lambda s: s.top)
    for i, s in enumerate(bottom, 1):
        s.name = f"BottomBlock{i}"

    # Right-side maturity sliders: textboxes + limited/high labels at x>9
    sliders = [s for s in shapes
               if not _is_native_title(s)
               and _to_inches(s.left) > 9.0
               and s.has_text_frame
               and _text(s).strip()]
    sliders.sort(key=lambda s: (s.top, s.left))
    for i, s in enumerate(sliders, 1):
        s.name = f"Slider{i}"

    # Timeline labels at bottom y>6.8
    timeline = [s for s in shapes
                if not _is_native_title(s)
                and _to_inches(s.top) > 6.7
                and s.has_text_frame
                and _text(s).strip()]
    timeline.sort(key=lambda s: s.left)
    for i, s in enumerate(timeline, 1):
        s.name = f"Timeline{i}"

    prs.save(str(pptx))
    print(f"[process] renamed blocks/sliders/timeline")


def _rename_nextsteps():
    pptx = Path("presentation_assets/next_steps/NEXTSTEPS-REGISTER-001/asset.pptx")
    prs = Presentation(str(pptx))
    shapes = _collect_leaf_shapes(prs.slides[0])

    # Header labels at y≈1.45, x sorted
    headers = [s for s in shapes
               if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
               and not _is_native_title(s)
               and 1.35 <= _to_inches(s.top) <= 1.6
               and _text(s).strip()]
    headers.sort(key=lambda s: s.left)
    header_names = ["HeaderNr", "HeaderPriority", "HeaderInstrument", "HeaderNextStep", "HeaderWhen", "HeaderWho"]
    for s, name in zip(headers, header_names):
        s.name = name

    # Row number cells (1-4) at x≈0.66, tall
    for s in shapes:
        if (s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
                and _to_inches(s.left) < 1.0
                and _to_inches(s.height) > 0.5
                and _text(s).strip() in ("1", "2", "3", "4")):
            s.name = f"Row{_text(s).strip()}Nr"

    # Priority column (x≈1.07, tall boxes, below header)
    priorities = [s for s in shapes
                  if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
                  and 0.9 <= _to_inches(s.left) <= 1.5
                  and _to_inches(s.top) > 1.6
                  and _to_inches(s.height) > 0.8
                  and not _text(s).strip().isdigit()]
    priorities.sort(key=lambda s: s.top)
    for i, s in enumerate(priorities, 1):
        s.name = f"Row{i}Priority"

    # Instrument column (x≈2.08, tall boxes, below header)
    instruments = [s for s in shapes
                   if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
                   and 1.8 <= _to_inches(s.left) <= 2.5
                   and _to_inches(s.top) > 1.6
                   and _to_inches(s.height) > 0.8]
    instruments.sort(key=lambda s: s.top)
    for i, s in enumerate(instruments, 1):
        s.name = f"Row{i}Instrument"

    # Next step column (x≈3.5, wide, short rows, below header)
    steps = [s for s in shapes
             if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
             and 3.0 <= _to_inches(s.left) <= 4.2
             and _to_inches(s.top) > 1.6
             and _to_inches(s.width) > 3
             and _to_inches(s.height) < 0.8]
    steps.sort(key=lambda s: s.top)
    for i, s in enumerate(steps, 1):
        s.name = f"Row{i}NextStep"

    # When column (x≈9.68, below header)
    whens = [s for s in shapes
             if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
             and 9.0 <= _to_inches(s.left) <= 10.5
             and _to_inches(s.top) > 1.6
             and _to_inches(s.height) < 0.8]
    whens.sort(key=lambda s: s.top)
    for i, s in enumerate(whens, 1):
        s.name = f"Row{i}When"

    # Who column (x≈11.4, below header)
    whos = [s for s in shapes
            if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
            and 11.0 <= _to_inches(s.left) <= 12.0
            and _to_inches(s.top) > 1.6
            and _to_inches(s.height) < 0.8]
    whos.sort(key=lambda s: s.top)
    for i, s in enumerate(whos, 1):
        s.name = f"Row{i}Who"

    prs.save(str(pptx))
    print(f"[next_steps] renamed headers/rows")


def _rename_kpi():
    pptx = Path("presentation_assets/kpi/KPI-BAR-INDICATOR-001/asset.pptx")
    prs = Presentation(str(pptx))
    shapes = _collect_leaf_shapes(prs.slides[0])

    # 6 category labels on far left
    cats = [s for s in shapes
            if not _is_native_title(s)
            and _to_inches(s.left) < 1.5
            and _text(s).strip().startswith("Category")]
    cats.sort(key=lambda s: s.top)
    for i, s in enumerate(cats, 1):
        s.name = f"Category{i}"

    # 6 bar value labels (Rechteck 73-78 area)
    values = [s for s in shapes
              if not _is_native_title(s)
              and 1.5 <= _to_inches(s.left) <= 3.0
              and 1.5 <= _to_inches(s.top) <= 7.0
              and _text(s).strip()]
    values.sort(key=lambda s: s.top)
    for i, s in enumerate(values, 1):
        s.name = f"Value{i}"

    # 4 right-side labels
    rights = [s for s in shapes
              if not _is_native_title(s)
              and 7.0 <= _to_inches(s.left) <= 8.0
              and _text(s).strip()]
    rights.sort(key=lambda s: s.top)
    for i, s in enumerate(rights, 1):
        s.name = f"RightLabel{i}"

    prs.save(str(pptx))
    print(f"[kpi] renamed {len(cats)} categories, {len(values)} values, {len(rights)} right labels")


def main():
    _rename_benefits()
    _rename_risk()
    _rename_opportunity()
    _rename_roadmap()
    _rename_process()
    _rename_nextsteps()
    _rename_kpi()


if __name__ == "__main__":
    main()
