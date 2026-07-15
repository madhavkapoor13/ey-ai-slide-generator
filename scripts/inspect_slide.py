#!/usr/bin/env python3
"""Inspect a specific slide in a PPTX and print shape summary."""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu


def _to_inches(emu: int) -> float:
    return round(emu / 914400, 2)


def _collect(shapes, group_path, out):
    for shape in shapes:
        name = shape.name or ""
        stype = str(shape.shape_type)
        try:
            stype = shape.shape_type.name
        except AttributeError:
            pass

        if stype == "GROUP":
            _collect(shape.shapes, group_path + [name], out)
            continue

        left = _to_inches(shape.left)
        top = _to_inches(shape.top)
        width = _to_inches(shape.width) if hasattr(shape, "width") else 0
        height = _to_inches(shape.height) if hasattr(shape, "height") else 0
        text = ""
        if shape.has_text_frame:
            text = shape.text_frame.text or ""
        is_placeholder = getattr(shape, "is_placeholder", False)
        ph_idx = None
        if is_placeholder:
            try:
                ph_idx = int(shape.placeholder_format.idx)
            except Exception:
                pass
        visible = True
        try:
            visible = not shape.element.get("hidden") == "1"
        except Exception:
            pass

        out.append({
            "name": name,
            "type": stype,
            "group": " / ".join(group_path),
            "left": left,
            "top": top,
            "width": width,
            "height": height,
            "text": text,
            "placeholder": is_placeholder,
            "ph_idx": ph_idx,
            "visible": visible,
        })


def inspect_slide(pptx_path: str, slide_index: int):
    prs = Presentation(pptx_path)
    if not (0 <= slide_index < len(prs.slides)):
        print(f"Slide index {slide_index} out of range (0-{len(prs.slides)-1})")
        sys.exit(1)
    slide = prs.slides[slide_index]
    out = []
    _collect(slide.shapes, [], out)
    print(f"Slide {slide_index}: {len(out)} leaf shapes")
    print(f"{'idx':<4} {'name':<35} {'type':<18} {'left':>6} {'top':>6} {'w':>6} {'h':>6} {'ph':>4} {'vis':>4} text")
    for i, s in enumerate(out):
        text_preview = s['text'].replace('\n', ' ')[:60]
        print(f"{i:<4} {s['name']:<35} {s['type']:<18} {s['left']:>6} {s['top']:>6} {s['width']:>6} {s['height']:>6} {str(s['ph_idx']):>4} {str(s['visible'])[0]:>4} {text_preview}")


if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("Usage: python inspect_slide.py <pptx> <slide_index>")
        sys.exit(1)
    inspect_slide(sys.argv[1], int(sys.argv[2]))
