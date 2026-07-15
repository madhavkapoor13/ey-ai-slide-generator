"""
tests/_asset_factory.py
=======================
Synthetic asset builders for Presentation Asset Library tests.

Built at test time (no committed binary fixtures). Produces two kinds of
artifacts into a caller-supplied temp directory:

- ``write_asset_pptx(path, ...)`` — a one-slide .pptx with:
    * a native title placeholder,
    * a named custom shape "SubtitleShape",
    * three named phase shapes "Phase1Label", "Phase2Label",
      "Phase3Label" (a repeating group).

- ``write_asset_manifest(asset_dir, manifest_dict)`` — writes an
  ``asset.json`` next to an ``asset.pptx``.

- ``write_full_asset(root, family, asset_id, ...)`` — builds the full
  ``<family>/<asset_id>/{asset.pptx,asset.json}`` tree in one call.

Designed for ``unittest``: callers use ``tempfile.TemporaryDirectory()``
in ``setUp`` and pass the temp path here. Hermetic and offline.
"""

from __future__ import annotations

import json
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt

from schemas.presentation_asset import (
    AssetManifest,
    AssetPlaceholder,
    PlaceholderBinding,
    PlaceholderKind,
    RepeatingGroup,
)


def write_asset_pptx(
    path: str | Path,
    *,
    title: str = "Sample Title",
    subtitle: str = "Sample Subtitle",
    phases: list[str] | None = None,
    extra_named_shapes: list[tuple[str, str]] | None = None,
) -> Path:
    """
    Build a synthetic one-slide .pptx at ``path`` (or path/asset.pptx).

    The slide contains:
      - a native title placeholder (filled with ``title``),
      - a named custom textbox "SubtitleShape" (filled with ``subtitle``),
      - one named textbox per phase label
        ("Phase1Label", "Phase2Label", ...) filled with the phase text,
      - any ``extra_named_shapes`` as ``(name, text)`` textboxes (for
        kind-heuristic / multi-field-repeating tests).

    Returns the path to the written .pptx file.
    """
    if phases is None:
        phases = ["Phase 1", "Phase 2", "Phase 3"]

    pptx_path = Path(path)
    pptx_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    blank = prs.slide_layouts[5]
    slide = prs.slides.add_slide(blank)

    title_ph = slide.shapes.title
    if title_ph is not None:
        title_ph.text = title

    sub_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(6), Inches(1))
    sub_box.name = "SubtitleShape"
    sub_box.text_frame.text = subtitle

    for i, phase in enumerate(phases, start=1):
        box = slide.shapes.add_textbox(
            Inches(1 + i), Inches(3), Inches(1.5), Inches(0.6)
        )
        box.name = f"Phase{i}Label"
        box.text_frame.text = phase

    for i, (shape_name, shape_text) in enumerate(extra_named_shapes or []):
        box = slide.shapes.add_textbox(
            Inches(1 + i), Inches(4.5), Inches(1.5), Inches(0.6)
        )
        box.name = shape_name
        box.text_frame.text = shape_text

    prs.save(str(pptx_path))
    return pptx_path


def _placeholder(
    pid: str,
    role: str,
    kind: PlaceholderKind,
    *,
    binding=None,
    required: bool = True,
    cardinality: str = "1",
) -> AssetPlaceholder:
    return AssetPlaceholder(
        id=pid,
        role=role,
        kind=kind,
        required=required,
        cardinality=cardinality,
        binding=binding or PlaceholderBinding(shape_name=pid),
    )


def build_roadmap_manifest(
    asset_id: str = "ROADMAP-3PHASE-001",
    *,
    family: str = "roadmap",
    density: int = 3,
    density_range: list[int] | None = None,
    audience_tags: list[str] | None = None,
    style_tags: list[str] | None = None,
    recommended_for: list[str] | None = None,
    avoid_for: list[str] | None = None,
    purpose: str = "Sequenced transformation roadmap with phased milestones.",
    family_aliases: list[str] | None = None,
) -> AssetManifest:
    """Build a valid 3-phase roadmap AssetManifest for tests."""
    placeholders = [
        _placeholder("title", "title", PlaceholderKind.TITLE, binding=PlaceholderBinding(native_placeholder_idx=0)),
        _placeholder("subtitle", "subtitle", PlaceholderKind.BODY, binding=PlaceholderBinding(shape_name="SubtitleShape")),
        _placeholder(
            f"phase_label",
            "phase",
            PlaceholderKind.BODY,
            binding=PlaceholderBinding(shape_name="Phase{N}Label"),
            cardinality="N",
            required=False,
        ),
    ]
    return AssetManifest(
        asset_id=asset_id,
        family=family,
        family_aliases=family_aliases or ["phased plan", "migration path"],
        purpose=purpose,
        audience_tags=audience_tags or ["board", "executive"],
        style_tags=style_tags or ["minimal", "modern"],
        recommended_for=recommended_for or ["Transformation Roadmap", "Strategy"],
        avoid_for=avoid_for or ["Detailed Process Documentation", "Technical Architecture"],
        density=density,
        density_range=density_range or [3, 6],
        fits_content_kinds=["phases", "milestones", "steps"],
        supports_images=False,
        placeholders=placeholders,
        repeating=RepeatingGroup(
            group_template="Phase{N}Label",
            placeholders_per_group=["phase_label"],
            index_token="{N}",
            count=density,
        ),
    )


def write_asset_manifest(asset_dir: str | Path, manifest: AssetManifest | dict[str, Any]) -> Path:
    """Write ``asset.json`` into ``asset_dir``. Accepts a model or a dict."""
    asset_dir = Path(asset_dir)
    asset_dir.mkdir(parents=True, exist_ok=True)
    manifest_path = asset_dir / "asset.json"
    if isinstance(manifest, AssetManifest):
        data = manifest.model_dump(mode="json")
    else:
        data = manifest
    with manifest_path.open("w", encoding="utf-8") as f:
        json.dump(data, f, indent=2)
    return manifest_path


def write_full_asset(
    root: str | Path,
    family: str,
    asset_id: str,
    *,
    manifest: AssetManifest | None = None,
    write_pptx: bool = True,
) -> Path:
    """
    Build a full ``<family>/<asset_id>/`` asset directory under ``root``.

    Writes ``asset.json`` (and ``asset.pptx`` unless ``write_pptx=False``)
    and returns the asset directory path.
    """
    asset_dir = Path(root) / family / asset_id
    asset_dir.mkdir(parents=True, exist_ok=True)
    if write_pptx:
        write_asset_pptx(asset_dir / "asset.pptx")
    if manifest is None:
        manifest = build_roadmap_manifest(asset_id=asset_id, family=family)
    write_asset_manifest(asset_dir, manifest)
    return asset_dir