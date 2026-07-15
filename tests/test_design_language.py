import unittest
from unittest.mock import MagicMock, patch

from pptx import Presentation
from pptx.util import Inches

from backend.design_system import design_language
from backend.design_system.theme_loader import clear_cache, get_current_theme, set_active_theme
from ppt_renderer.components.component_dispatcher import render as dispatch_render
from ppt_renderer.components.executive_card_renderer import render as render_executive_card
from schemas.executive_card import ExecutiveCardContent
from schemas.layout import ComponentSpecification


class DesignLanguageEngineTests(unittest.TestCase):
    """Tests for the Sprint I5 Design Language Extraction Engine."""

    def setUp(self):
        design_language.clear_cache()
        clear_cache()
        set_active_theme("ey_blue")

    def tearDown(self):
        design_language.clear_cache()
        clear_cache()
        set_active_theme("ey_blue")

    def test_loads_default_spacing_rules(self):
        spacing = design_language.get_card_spacing()
        self.assertEqual(spacing["padding"], 0.12)
        self.assertEqual(spacing["gap"], 0.10)
        self.assertEqual(spacing["border_width"], 0.011111)

    def test_loads_alignment_rules(self):
        self.assertEqual(design_language.get_alignment("card_grid"), "distributed")
        self.assertEqual(design_language.get_alignment("timeline"), "distributed")
        self.assertEqual(design_language.get_alignment("matrix"), "stretch")
        self.assertEqual(design_language.get_alignment("label"), "center")

    def test_loads_hierarchy_rules(self):
        title = design_language.get_typography_hierarchy("executive_title")
        self.assertEqual(title["font"], "title")
        self.assertEqual(title["size"], 24.0)
        self.assertTrue(title["bold"])

        body = design_language.get_typography_hierarchy("body")
        self.assertEqual(body["font"], "body")
        self.assertEqual(body["size"], 10.5)
        self.assertFalse(body["bold"])

    def test_loads_emphasis_rules(self):
        primary = design_language.get_emphasis("primary")
        self.assertTrue(primary["bold"])
        self.assertEqual(primary["color_role"], "ink")

        muted = design_language.get_emphasis("muted")
        self.assertFalse(muted["bold"])
        self.assertEqual(muted["color_role"], "grey")

    def test_pattern_specific_override_infographic(self):
        # IG-02 overrides timeline connector_width in infographics.json.
        value = design_language.get(
            "timeline", "connector_width", pattern_id="IG-02", category="spacing"
        )
        self.assertEqual(value, 0.020000)

    def test_visual_patterns_design_metadata_is_ingested(self):
        # CL-01 design_metadata in creative_patterns.json is readable.
        value = design_language.get(
            "card", "padding", pattern_id="CL-01", category="spacing"
        )
        self.assertEqual(value, 0.12)

    def test_missing_rule_returns_default(self):
        self.assertEqual(
            design_language.get("nonexistent", "key", default="fallback"),
            "fallback",
        )
        self.assertEqual(
            design_language.get("nonexistent", default={}),
            {},
        )

    def test_missing_rule_returns_none_when_no_default(self):
        self.assertIsNone(design_language.get("nonexistent", "key"))

    def test_generic_lookup_without_category(self):
        # "card" context exists in spacing.json.
        self.assertEqual(design_language.get("card", "padding"), 0.12)
        # "primary" context exists in emphasis.json.
        self.assertEqual(design_language.get("primary", "color_role"), "ink")

    def test_get_header_and_footer_spacing(self):
        header = design_language.get_header_spacing()
        self.assertEqual(header["height"], 0.15)
        self.assertEqual(header["margin_bottom"], 0.03)

        footer = design_language.get_footer_spacing()
        self.assertEqual(footer["height"], 0.10)
        self.assertEqual(footer["margin_top"], 0.03)

    def test_get_whitespace(self):
        whitespace = design_language.get_whitespace()
        self.assertEqual(whitespace["min_component_gap"], 0.05)
        self.assertEqual(whitespace["preferred_component_gap"], 0.10)


class DesignLanguageRendererIntegrationTests(unittest.TestCase):
    """Tests that component renderers consume design language rules."""

    def setUp(self):
        design_language.clear_cache()
        clear_cache()
        set_active_theme("ey_blue")
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

    def tearDown(self):
        design_language.clear_cache()
        clear_cache()
        set_active_theme("ey_blue")

    def _make_card_spec(self, placeholder: str) -> ComponentSpecification:
        return ComponentSpecification(
            component_id="ec",
            type="executive_card",
            placeholder=placeholder,
            x=0.1,
            y=0.2,
            width=0.3,
            height=0.4,
        )

    def _content(self, card: ExecutiveCardContent) -> dict:
        return {"cards": [card.model_dump(mode="json")]}

    def test_dispatcher_passes_layout_context(self):
        """The dispatcher forwards layout_context to the renderer module."""
        mock_renderer = MagicMock()
        with patch.dict(
            "ppt_renderer.components.component_dispatcher._DISPATCH_TABLE",
            {"text": mock_renderer},
            clear=False,
        ):
            spec = ComponentSpecification(
                component_id="txt",
                type="text",
                placeholder="content",
                x=0.1,
                y=0.2,
                width=0.3,
                height=0.4,
            )
            dispatch_render(spec, self.prs, self.slide, {}, layout_context={"pattern_id": "CL-01"})

        mock_renderer.render.assert_called_once()
        kwargs = mock_renderer.render.call_args.kwargs
        self.assertEqual(kwargs["layout_context"]["pattern_id"], "CL-01")

    def test_text_renderer_uses_design_language_alignment(self):
        """Label components are center-aligned according to design language."""
        from pptx.enum.text import PP_ALIGN

        spec = ComponentSpecification(
            component_id="lbl",
            type="label",
            placeholder="content",
            x=0.1,
            y=0.2,
            width=0.3,
            height=0.1,
        )
        dispatch_render(spec, self.prs, self.slide, {"text": "Label"})

        textbox = [s for s in self.slide.shapes if s.has_text_frame and s.text_frame.text == "Label"][0]
        self.assertEqual(textbox.text_frame.paragraphs[0].alignment, PP_ALIGN.CENTER)


if __name__ == "__main__":
    unittest.main()
