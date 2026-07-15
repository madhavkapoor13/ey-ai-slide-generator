from __future__ import annotations

import unittest

from pptx import Presentation
from pptx.util import Inches

from ppt_renderer.components import component_dispatcher
from schemas.layout import ComponentSpecification


class ConnectorRendererTests(unittest.TestCase):
    def test_connector_component_renders(self):
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        component = ComponentSpecification(
            component_id="conn_1",
            type="connector",
            x=0.2,
            y=0.4,
            width=0.3,
            height=0.02,
            placeholder="conn_1",
            constraints=["arrow"],
        )

        component_dispatcher.render(component, prs, slide, {})

        self.assertEqual(len(slide.shapes), 1)
        self.assertIsNotNone(slide.shapes[0].line)

    def test_connector_renders_arrowhead_in_ooxml(self):
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        component = ComponentSpecification(
            component_id="conn_1",
            type="connector",
            x=0.2,
            y=0.4,
            width=0.3,
            height=0.02,
            placeholder="conn_1",
            constraints=["arrow"],
        )

        component_dispatcher.render(component, prs, slide, {})

        connector = slide.shapes[0]
        spPr = connector._element.find(
            "{http://schemas.openxmlformats.org/presentationml/2006/main}spPr"
        )
        self.assertIsNotNone(spPr)
        ln = spPr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}ln")
        self.assertIsNotNone(ln)
        tail = ln.find("{http://schemas.openxmlformats.org/drawingml/2006/main}tailEnd")
        self.assertIsNotNone(tail)
        self.assertEqual(tail.get("type"), "triangle")


if __name__ == "__main__":
    unittest.main()
