"""
tests/test_asset_inspector.py
=============================
Tests for the Asset Inspector: shape analysis, repeating detection, and
``AssetManifest`` draft generation.

Synthetic .pptx files are built at test time via ``_asset_factory`` and
raw python-pptx. Hermetic; offline; no real assets required.
"""

import os
import tempfile
import unittest
from pathlib import Path

from backend.presentation_assets import asset_inspector
from schemas.presentation_asset import PlaceholderKind
from _asset_factory import write_asset_pptx


class AssetInspectorTests(unittest.TestCase):

    def setUp(self):
        self._tmp = tempfile.TemporaryDirectory()
        self.root = self._tmp.name

    def tearDown(self):
        self._tmp.cleanup()

    def _pptx(self, name: str = "asset.pptx", **kwargs) -> Path:
        path = os.path.join(self.root, name)
        return write_asset_pptx(path, **kwargs)

    def test_inspect_emits_title_placeholder_with_native_binding(self):
        pptx_path = self._pptx()
        report = asset_inspector.inspect(pptx_path)

        title = next(
            p for p in report.proposed_manifest.placeholders if p.id == "title"
        )
        self.assertEqual(title.kind, PlaceholderKind.TITLE)
        self.assertEqual(title.role, "title")
        self.assertEqual(title.cardinality, "1")
        self.assertIsNotNone(title.binding.native_placeholder_idx)

    def test_inspect_emits_named_shape_placeholder(self):
        pptx_path = self._pptx()
        report = asset_inspector.inspect(pptx_path)

        subtitle = next(
            p for p in report.proposed_manifest.placeholders if p.id == "subtitle_shape"
        )
        self.assertEqual(subtitle.kind, PlaceholderKind.BODY)
        self.assertIsNotNone(subtitle.binding.shape_name)

    def test_inspect_detects_repeating_three_phases(self):
        pptx_path = self._pptx()
        report = asset_inspector.inspect(pptx_path)

        self.assertEqual(len(report.repeating), 1)
        det = report.repeating[0]
        self.assertEqual(det.stripped_id, "phase_label")
        self.assertEqual(det.shape_name_template, "Phase{N}Label")
        self.assertEqual(det.indices, [1, 2, 3])

    def test_inspect_collapses_repeating_into_single_template_placeholder(self):
        pptx_path = self._pptx()
        report = asset_inspector.inspect(pptx_path)
        manifest = report.proposed_manifest

        phase_placeholders = [p for p in manifest.placeholders if p.cardinality == "N"]
        self.assertEqual(len(phase_placeholders), 1)
        ph = phase_placeholders[0]
        self.assertEqual(ph.id, "phase_label")
        self.assertEqual(ph.binding.shape_name, "Phase{N}Label")
        self.assertEqual(ph.cardinality, "N")

    def test_inspect_emits_repeating_group_metadata(self):
        pptx_path = self._pptx()
        report = asset_inspector.inspect(pptx_path)
        rep = report.proposed_manifest.repeating
        self.assertIsNotNone(rep)
        self.assertEqual(rep.count, 3)
        self.assertEqual(rep.group_template, "Phase{N}Label")
        self.assertEqual(rep.placeholders_per_group, ["phase_label"])

    def test_inspect_density_matches_repeating_count(self):
        pptx_path = self._pptx()
        report = asset_inspector.inspect(pptx_path)
        self.assertEqual(report.proposed_manifest.density, 3)
        self.assertEqual(report.proposed_manifest.density_range, [3, 3])

    def test_inspect_multi_field_repeating(self):
        pptx_path = self._pptx(
            extra_named_shapes=[
                ("Phase1Owner", "Alice"),
                ("Phase2Owner", "Bob"),
                ("Phase3Owner", "Carol"),
            ],
        )
        report = asset_inspector.inspect(pptx_path)
        repeating_ids = {r.stripped_id for r in report.repeating}
        self.assertEqual(repeating_ids, {"phase_label", "phase_owner"})

        manifest_ids = {p.id for p in report.proposed_manifest.placeholders if p.cardinality == "N"}
        self.assertEqual(manifest_ids, {"phase_label", "phase_owner"})

        rep = report.proposed_manifest.repeating
        self.assertEqual(rep.count, 3)
        self.assertEqual(rep.placeholders_per_group, ["phase_label", "phase_owner"])

    def test_inspect_non_sequential_indices_not_repeating(self):
        # Only Phase1 + Phase3 (no Phase2 → not a 1..N sequence)
        path = os.path.join(self.root, "non-seq.pptx")
        from pptx import Presentation
        from pptx.util import Inches
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        for n in [1, 3]:
            box = slide.shapes.add_textbox(Inches(n), Inches(3), Inches(1.5), Inches(0.6))
            box.name = f"Phase{n}Label"
            box.text_frame.text = f"Phase {n}"
        prs.save(path)

        report = asset_inspector.inspect(path)
        self.assertEqual(report.repeating, [])
        self.assertIsNone(report.proposed_manifest.repeating)

    def test_inspect_kind_heuristics(self):
        pptx_path = self._pptx(
            extra_named_shapes=[
                ("KpiValue", "42"),
                ("CostMetric", "$1.2M"),
                ("GrowthTrend", "12%"),
                ("QDate", "Q1 2025"),
                ("IconChevron", ""),
                ("JourneyNode", ""),
            ],
        )
        report = asset_inspector.inspect(pptx_path)
        by_id = {p.id: p for p in report.proposed_manifest.placeholders}

        self.assertEqual(by_id["kpi_value"].kind, PlaceholderKind.METRIC)
        self.assertEqual(by_id["cost_metric"].kind, PlaceholderKind.CURRENCY)
        self.assertEqual(by_id["growth_trend"].kind, PlaceholderKind.PERCENTAGE)
        self.assertEqual(by_id["q_date"].kind, PlaceholderKind.DATE)
        # Empty text → fall back to shape-name hint
        self.assertEqual(by_id["icon_chevron"].kind, PlaceholderKind.CHEVRON)
        self.assertEqual(by_id["journey_node"].kind, PlaceholderKind.TIMELINE_NODE)

    def test_inspect_supports_images_false_no_pictures(self):
        pptx_path = self._pptx()
        report = asset_inspector.inspect(pptx_path)
        self.assertFalse(report.supports_images)
        self.assertFalse(report.proposed_manifest.supports_images)

    def test_inspect_asset_id_from_asset_pptx_filename(self):
        named_path = os.path.join(self.root, "ROADMAP-3PHASE-001.pptx")
        write_asset_pptx(named_path)
        report = asset_inspector.inspect(named_path)
        self.assertEqual(report.proposed_manifest.asset_id, "ROADMAP-3PHASE-001")

    def test_inspect_asset_id_from_parent_dir_when_file_is_asset_pptx(self):
        asset_dir = os.path.join(self.root, "ROADMAP-XYZ-002")
        os.makedirs(asset_dir, exist_ok=True)
        path = os.path.join(asset_dir, "asset.pptx")
        write_asset_pptx(path)
        report = asset_inspector.inspect(path)
        self.assertEqual(report.proposed_manifest.asset_id, "ROADMAP-XYZ-002")

    def test_inspect_family_inferred_from_known_parent(self):
        family_dir = os.path.join(self.root, "roadmap", "ROADMAP-3PHASE-001")
        os.makedirs(family_dir, exist_ok=True)
        path = os.path.join(family_dir, "asset.pptx")
        write_asset_pptx(path)
        report = asset_inspector.inspect(path)
        self.assertEqual(report.proposed_manifest.family, "roadmap")

    def test_inspect_family_blank_when_parent_unknown(self):
        # path under a non-family dir → family remains blank for human
        path = os.path.join(self.root, "asset.pptx")
        write_asset_pptx(path)
        report = asset_inspector.inspect(path)
        self.assertEqual(report.proposed_manifest.family, "")

    def test_inspect_family_override_arg(self):
        path = os.path.join(self.root, "asset.pptx")
        write_asset_pptx(path)
        report = asset_inspector.inspect(path, family="timeline")
        self.assertEqual(report.proposed_manifest.family, "timeline")

    def test_inspect_asset_id_override_arg(self):
        path = os.path.join(self.root, "asset.pptx")
        write_asset_pptx(path)
        report = asset_inspector.inspect(path, asset_id="EXEC-001")
        self.assertEqual(report.proposed_manifest.asset_id, "EXEC-001")

    def test_inspect_authored_metadata_left_blank(self):
        pptx_path = self._pptx()
        report = asset_inspector.inspect(pptx_path)
        m = report.proposed_manifest
        self.assertEqual(m.purpose, "")
        self.assertEqual(m.audience_tags, [])
        self.assertEqual(m.style_tags, [])
        self.assertEqual(m.recommended_for, [])
        self.assertEqual(m.avoid_for, [])
        self.assertEqual(m.family_aliases, [])
        self.assertEqual(m.fits_content_kinds, [])

    def test_inspect_skip_decorative_shapes_without_text(self):
        # No text frame ⇒ not enumerated as a fillable placeholder.
        # All factory shapes have text frames; basic smoke test confirms
        # the proposed manifest has only fillable slots.
        pptx_path = self._pptx()
        report = asset_inspector.inspect(pptx_path)
        for ph in report.proposed_manifest.placeholders:
            self.assertTrue(
                ph.binding.native_placeholder_idx is not None
                or ph.binding.shape_name is not None
            )


if __name__ == "__main__":
    unittest.main()