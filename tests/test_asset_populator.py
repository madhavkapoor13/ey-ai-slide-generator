import base64
import io
import tempfile
import unittest
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

from backend.presentation_assets.asset_loader import enumerate_shapes, open_for_population
from backend.presentation_assets.asset_populator import (
    _BLIP_TAG,
    _R_NS,
    _copy_slide,
    PopulatorError,
    populate_asset_slide,
    populate_slide,
)
from backend.presentation_assets.asset_registry import clear_cache
from schemas.presentation_asset import (
    AssetManifest,
    AssetPlaceholder,
    PlaceholderBinding,
    PlaceholderKind,
    RepeatingGroup,
)
from schemas.slide_spec import SlideSpec
from tests._asset_factory import build_roadmap_manifest, write_full_asset


def _shape_text(slide, name: str) -> str:
    """Return the current text of the shape named ``name`` on ``slide``."""
    for shape in slide.shapes:
        if shape.name == name:
            return shape.text_frame.text if shape.has_text_frame else ""
        if shape.shape_type == 6:  # GROUP
            for child in shape.shapes:
                if child.name == name:
                    return child.text_frame.text if child.has_text_frame else ""
    return ""


def _shape_by_name(slide, name: str):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
        if shape.shape_type == 6:  # GROUP
            for child in shape.shapes:
                if child.name == name:
                    return child
    return None


class AssetPopulatorTests(unittest.TestCase):
    """Tests for Sprint E — Asset Populator."""

    def setUp(self):
        self.temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(self.temp_dir.cleanup)
        self.assets_root = Path(self.temp_dir.name) / "presentation_assets"
        write_full_asset(self.assets_root, "roadmap", "ROADMAP-3PHASE-001")
        # Ensure the default registry cache is cleared so the temp assets can be
        # discovered when tests run in arbitrary order.
        clear_cache()

    def test_populate_slide_sets_title_and_subtitle(self):
        _, slide = open_for_population(
            "ROADMAP-3PHASE-001", assets_dir=self.assets_root
        )
        manifest = build_roadmap_manifest(asset_id="ROADMAP-3PHASE-001")
        content = {
            "title": "Transformation Roadmap",
            "subtitle": "TestCo Operations — Three-Phase Journey",
        }
        populate_slide(slide, content, manifest)

        self.assertEqual(_shape_text(slide, "Title 1"), "Transformation Roadmap")
        self.assertEqual(_shape_text(slide, "SubtitleShape"), "TestCo Operations — Three-Phase Journey")

    def test_populate_slide_preserves_title_font_size_and_bold(self):
        _, slide = open_for_population(
            "ROADMAP-3PHASE-001", assets_dir=self.assets_root
        )
        title_shape = _shape_by_name(slide, "Title 1")
        self.assertIsNotNone(title_shape)
        original_run = title_shape.text_frame.paragraphs[0].runs[0]
        original_run.font.size = Pt(40)
        original_run.font.bold = True

        manifest = build_roadmap_manifest(asset_id="ROADMAP-3PHASE-001")
        content = {
            "title": "Transformation Roadmap",
            "subtitle": "TestCo Operations — Three-Phase Journey",
        }
        populate_slide(slide, content, manifest)

        populated_run = title_shape.text_frame.paragraphs[0].runs[0]
        self.assertEqual(populated_run.text, "Transformation Roadmap")
        self.assertEqual(populated_run.font.size, Pt(40))
        self.assertTrue(populated_run.font.bold)

    def test_populate_slide_expands_repeating_group(self):
        _, slide = open_for_population(
            "ROADMAP-3PHASE-001", assets_dir=self.assets_root
        )
        manifest = build_roadmap_manifest(asset_id="ROADMAP-3PHASE-001")
        content = {
            "title": "Roadmap",
            "subtitle": "Subtitle",
            "phase_label": ["Discover", "Design", "Deliver"],
        }
        populate_slide(slide, content, manifest)

        self.assertEqual(_shape_text(slide, "Phase1Label"), "Discover")
        self.assertEqual(_shape_text(slide, "Phase2Label"), "Design")
        self.assertEqual(_shape_text(slide, "Phase3Label"), "Deliver")

    def test_populate_slide_clears_unused_repeating_shapes(self):
        _, slide = open_for_population(
            "ROADMAP-3PHASE-001", assets_dir=self.assets_root
        )
        manifest = build_roadmap_manifest(asset_id="ROADMAP-3PHASE-001")
        content = {
            "title": "Roadmap",
            "subtitle": "Subtitle",
            "phase_label": ["Discover", "Design"],
        }
        populate_slide(slide, content, manifest)

        self.assertEqual(_shape_text(slide, "Phase1Label"), "Discover")
        self.assertEqual(_shape_text(slide, "Phase2Label"), "Design")
        self.assertEqual(_shape_text(slide, "Phase3Label"), "")

    def test_populate_asset_slide_copies_into_target_presentation(self):
        target_prs = Presentation()
        manifest = build_roadmap_manifest(asset_id="ROADMAP-3PHASE-001")
        slide_spec = SlideSpec(
            slide_type="operating_model",
            raw_spec={
                "title": "Roadmap",
                "subtitle": "Subtitle",
                "phase_label": ["Q1", "Q2", "Q3"],
            },
            version="2.0",
            asset_id="ROADMAP-3PHASE-001",
        )

        new_slide = populate_asset_slide(
            target_prs, slide_spec, manifest, assets_dir=self.assets_root
        )

        self.assertEqual(len(target_prs.slides), 1)
        self.assertIs(new_slide, target_prs.slides[0])
        self.assertEqual(_shape_text(new_slide, "Title 1"), "Roadmap")
        self.assertEqual(_shape_text(new_slide, "Phase2Label"), "Q2")

    def test_copy_slide_preserves_embedded_picture_relationships(self):
        source_prs = Presentation()
        source_slide = source_prs.slides.add_slide(source_prs.slide_layouts[6])
        png_blob = base64.b64decode(
            "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR4nGP4z8DwHwAFgwJ/l0u1"
            "NwAAAABJRU5ErkJggg=="
        )
        source_slide.shapes.add_picture(
            io.BytesIO(png_blob),
            Inches(1),
            Inches(1),
            width=Inches(0.25),
            height=Inches(0.25),
        )
        target_prs = Presentation()

        new_slide, _ = _copy_slide(source_slide, target_prs)

        embed_attr = f"{{{_R_NS}}}embed"
        rel_ids = [
            blip.get(embed_attr)
            for blip in new_slide.element.iter(_BLIP_TAG)
            if blip.get(embed_attr)
        ]
        self.assertEqual(len(rel_ids), 1)
        self.assertEqual(new_slide.part.related_part(rel_ids[0]).blob, png_blob)

    def test_populate_slide_looks_up_manifest_when_not_supplied(self):
        target_prs = Presentation()
        slide_spec = SlideSpec(
            slide_type="operating_model",
            raw_spec={
                "title": "Roadmap",
                "subtitle": "Subtitle",
                "phase_label": ["Alpha", "Beta", "Gamma"],
            },
            version="2.0",
            asset_id="ROADMAP-3PHASE-001",
        )

        new_slide = populate_asset_slide(
            target_prs, slide_spec, assets_dir=self.assets_root
        )

        self.assertEqual(_shape_text(new_slide, "Title 1"), "Roadmap")
        self.assertEqual(_shape_text(new_slide, "Phase1Label"), "Alpha")

    def test_populate_slide_coerces_structured_dict_content(self):
        manifest = AssetManifest(
            asset_id="STRUCTURED-001",
            family="roadmap",
            purpose="Structured phases.",
            density=2,
            density_range=[2, 4],
            placeholders=[
                AssetPlaceholder(
                    id="phase_label",
                    role="phase",
                    kind=PlaceholderKind.BODY,
                    cardinality="N",
                    binding=PlaceholderBinding(shape_name="Phase{N}Label"),
                    content_schema={"phase": "string", "owner": "string?"},
                )
            ],
            repeating=RepeatingGroup(
                group_template="Phase{N}Group",
                placeholders_per_group=["phase_label"],
                index_token="{N}",
                count=2,
            ),
        )
        _, slide = open_for_population(
            "ROADMAP-3PHASE-001", assets_dir=self.assets_root
        )
        content = {
            "phase_label": [
                {"phase": "Discover", "owner": "Alice"},
                {"phase": "Design", "owner": "Bob"},
            ],
        }
        populate_slide(slide, content, manifest)

        self.assertEqual(_shape_text(slide, "Phase1Label"), "Discover")
        self.assertEqual(_shape_text(slide, "Phase2Label"), "Design")

    def test_populate_slide_missing_shape_does_not_crash(self):
        manifest = AssetManifest(
            asset_id="MISSING-001",
            family="roadmap",
            purpose="Missing shape test.",
            density=1,
            density_range=[1, 1],
            placeholders=[
                AssetPlaceholder(
                    id="ghost",
                    role="ghost",
                    kind=PlaceholderKind.BODY,
                    binding=PlaceholderBinding(shape_name="GhostShape"),
                )
            ],
        )
        _, slide = open_for_population(
            "ROADMAP-3PHASE-001", assets_dir=self.assets_root
        )
        content = {"ghost": "Boo"}
        # Should not raise.
        populate_slide(slide, content, manifest)

    def test_populate_asset_slide_raises_when_asset_id_missing(self):
        target_prs = Presentation()
        slide_spec = SlideSpec(
            slide_type="operating_model",
            raw_spec={},
            version="2.0",
            asset_id="UNKNOWN-ASSET-999",
        )
        with self.assertRaises(PopulatorError):
            populate_asset_slide(target_prs, slide_spec)

    def test_populate_asset_slide_raises_when_manifest_missing(self):
        target_prs = Presentation()
        slide_spec = SlideSpec(
            slide_type="operating_model",
            raw_spec={},
            version="2.0",
            asset_id="UNKNOWN-ASSET-999",
        )
        with self.assertRaises(PopulatorError):
            populate_asset_slide(target_prs, slide_spec, manifest=None)


class RealAssetIngestionTests(unittest.TestCase):
    """Smoke tests for the pilot assets ingested from the EYP template deck."""

    def setUp(self):
        from backend.presentation_assets.asset_registry import clear_cache, load_assets

        clear_cache()
        self.manifests = load_assets()

    def test_pilot_assets_are_registered(self):
        for asset_id in ["TIMELINE-6STEP-001", "PROCESS-7STEP-001", "LIST-6ITEM-001"]:
            self.assertIn(asset_id, self.manifests, f"{asset_id} should be registered")

    def test_pilot_assets_populate_without_error(self):
        from backend.presentation_assets.asset_loader import open_for_population

        for asset_id in ["TIMELINE-6STEP-001", "PROCESS-7STEP-001", "LIST-6ITEM-001"]:
            manifest = self.manifests[asset_id]
            _, slide = open_for_population(asset_id)
            content = {"title": f"{asset_id} Title"}
            for placeholder in manifest.placeholders:
                if placeholder.cardinality == "N":
                    count = manifest.density
                    if placeholder.content_schema:
                        content[placeholder.id] = [
                            {k: f"{asset_id} {i + 1}" for k in placeholder.content_schema}
                            for i in range(count)
                        ]
                    else:
                        content[placeholder.id] = [f"{asset_id} {i + 1}" for i in range(count)]
            # Should not raise.
            populate_slide(slide, content, manifest)

    def test_investment_asset_copy_preserves_svg_icon_relationships(self):
        _, source_slide = open_for_population("INVESTMENT-CASE-SUMMARY-001")
        target_prs = Presentation()

        new_slide, _ = _copy_slide(source_slide, target_prs)

        embed_attr = f"{{{_R_NS}}}embed"
        image_refs = []
        for el in new_slide.element.iter():
            rel_id = el.get(embed_attr)
            if not rel_id:
                continue
            related_part = new_slide.part.related_part(rel_id)
            if getattr(related_part, "content_type", "").startswith("image/"):
                image_refs.append((el.tag, related_part.partname, related_part.content_type))

        self.assertEqual(len(image_refs), 4)
        self.assertTrue(all(content_type == "image/svg+xml" for _, _, content_type in image_refs))

    def test_kpi_scorecard_copy_preserves_editable_table_shapes(self):
        _, source_slide = open_for_population("KPI-SCORECARD-TABLE-001")
        target_prs = Presentation()

        new_slide, _ = _copy_slide(source_slide, target_prs)

        self.assertEqual(_shape_text(new_slide, "TextBox 4"), "KPI")
        self.assertIn("Procurement Cycle Time", _shape_text(new_slide, "TextBox 27"))
        self.assertEqual(_shape_text(new_slide, "TextBox 31"), "COO")

    def test_kpi_scorecard_table_cells_are_manifest_fillable(self):
        manifest = self.manifests["KPI-SCORECARD-TABLE-001"]
        _, source_slide = open_for_population("KPI-SCORECARD-TABLE-001")
        target_prs = Presentation()
        new_slide, _ = _copy_slide(source_slide, target_prs)

        populate_slide(
            new_slide,
            {
                "title": "HSBC Finance AI KPI Scorecard",
                "subtitle": "Finance transformation controls and adoption metrics.",
                "kpi_name_1": "Close Cycle Time",
                "baseline_1": "8 days",
                "target_1": "5 days",
                "current_1": "6 days",
                "owner_1": "Controllership",
                "cadence_1": "Monthly",
                "status_1": "Improving",
                "comment_1": "AI reconciliations reducing manual work.",
            },
            manifest,
        )

        self.assertEqual(_shape_text(new_slide, "TextBox 27"), "Close Cycle Time")
        self.assertEqual(_shape_text(new_slide, "TextBox 31"), "Controllership")


if __name__ == "__main__":
    unittest.main()
