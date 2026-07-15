"""
tests/test_asset_loader.py
============================
Tests for the Asset Loader (python-pptx reader).

Synthetic .pptx files are built at test time in a temp dir via
``_asset_factory``. Verifies shape enumeration (native placeholder +
named custom shapes) and the population-open path resolution.
"""

import os
import tempfile
import unittest
from pathlib import Path

from backend.presentation_assets import asset_loader, asset_registry
from _asset_factory import write_asset_pptx, write_full_asset


class AssetLoaderTests(unittest.TestCase):

    def setUp(self):
        self._tmp = tempfile.TemporaryDirectory()
        self.root = self._tmp.name
        asset_registry.clear_cache()

    def tearDown(self):
        asset_registry.clear_cache()
        self._tmp.cleanup()

    def test_enumerate_shapes_finds_named_shapes_and_placeholder(self):
        pptx_path = write_asset_pptx(
            os.path.join(self.root, "sample.pptx"),
            title="Board Update",
            subtitle="FY26 Priorities",
            phases=["Phase 1", "Phase 2", "Phase 3"],
        )
        shapes = asset_loader.enumerate_shapes(pptx_path)

        names = {s.name for s in shapes}
        self.assertIn("SubtitleShape", names)
        self.assertIn("Phase1Label", names)
        self.assertIn("Phase2Label", names)
        self.assertIn("Phase3Label", names)

    def test_enumerate_shapes_reports_text_frames_and_current_text(self):
        pptx_path = write_asset_pptx(
            os.path.join(self.root, "sample.pptx"),
            subtitle="My Subtitle Text",
        )
        shapes = asset_loader.enumerate_shapes(pptx_path)

        subtitle = next(s for s in shapes if s.name == "SubtitleShape")
        self.assertTrue(subtitle.has_text_frame)
        self.assertEqual(subtitle.current_text, "My Subtitle Text")
        self.assertFalse(subtitle.is_placeholder)
        self.assertIsNone(subtitle.placeholder_idx)

    def test_enumerate_shapes_detects_native_placeholder(self):
        pptx_path = write_asset_pptx(os.path.join(self.root, "sample.pptx"), title="Hello")
        shapes = asset_loader.enumerate_shapes(pptx_path)

        placeholders = [s for s in shapes if s.is_placeholder]
        self.assertGreaterEqual(len(placeholders), 1, "should detect the native title placeholder")

        title_ph = next(s for s in shapes if s.is_placeholder)
        self.assertIsNotNone(title_ph.placeholder_idx)
        self.assertEqual(title_ph.current_text, "Hello")

    def test_enumerate_shapes_empty_presentation(self):
        prs_path = os.path.join(self.root, "empty.pptx")
        from pptx import Presentation
        Presentation().save(prs_path)
        shapes = asset_loader.enumerate_shapes(prs_path)
        self.assertEqual(shapes, [])

    def test_open_for_population_returns_slide(self):
        write_full_asset(self.root, "roadmap", "ROADMAP-3PHASE-001")

        prs, slide = asset_loader.open_for_population(
            "ROADMAP-3PHASE-001", assets_dir=self.root
        )
        self.assertEqual(len(prs.slides), 1)
        self.assertEqual(slide, prs.slides[0])

    def test_open_for_population_missing_asset(self):
        with self.assertRaises(FileNotFoundError):
            asset_loader.open_for_population("NOPE", assets_dir=self.root)

    def test_open_for_population_missing_pptx(self):
        asset_dir = write_full_asset(self.root, "roadmap", "ROADMAP-NO-PPTX")
        os.remove(os.path.join(asset_dir, "asset.pptx"))

        with self.assertRaises(FileNotFoundError):
            asset_loader.open_for_population("ROADMAP-NO-PPTX", assets_dir=self.root)


if __name__ == "__main__":
    unittest.main()