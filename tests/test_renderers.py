import os
import tempfile
import unittest

from pptx import Presentation

from ppt_renderer.operating_model_renderer import OperatingModelRenderer
from ppt_renderer.renderer import ProcessFlowRenderer


class RendererTests(unittest.TestCase):
    def _operating_model_spec(self, title: str = "Current State") -> dict:
        return {
            "title": title,
            "subtitle": "Toyota Procurement Operating Model",
            "description": "Current-state operating model for Procure-to-Pay.",
            "executive_summary": (
                "The Procure-to-Pay operating model coordinates core Procurement workflows for Toyota. "
                "Primary operational challenges arise from fragmented handoffs, inconsistent controls, and limited cross-functional visibility."
            ),
            "summary": {
                "headline": "Procure-to-Pay",
                "description": (
                    "The Procure-to-Pay operating model coordinates core Procurement workflows for Toyota. "
                    "Primary operational challenges arise from fragmented handoffs, inconsistent controls, and limited cross-functional visibility."
                ),
                "metrics": [],
            },
            "stages": [
                {"number": 1, "title": "Stage 1", "label": "Stage 1", "activities": ["Activity one", "Activity two"]},
                {"number": 2, "title": "Stage 2", "label": "Stage 2", "activities": ["Activity three", "Activity four"]},
                {"number": 3, "title": "Stage 3", "label": "Stage 3", "activities": ["Activity five", "Activity six"]},
                {"number": 4, "title": "Stage 4", "label": "Stage 4", "activities": ["Activity seven", "Activity eight"]},
                {"number": 5, "title": "Stage 5", "label": "Stage 5", "activities": ["Activity nine", "Activity ten"]},
                {"number": 6, "title": "Stage 6", "label": "Stage 6", "activities": ["Activity eleven", "Activity twelve"]},
            ],
            "pain_points": [
                {"stage": "Stage 1", "text": "Fragmented stage one ownership delays decision support."},
                {"stage": "Stage 2", "text": "Fragmented stage two ownership delays decision support."},
                {"stage": "Stage 3", "text": "Fragmented stage three ownership delays decision support."},
                {"stage": "Stage 4", "text": "Fragmented stage four ownership delays decision support."},
                {"stage": "Stage 5", "text": "Fragmented stage five ownership delays decision support."},
                {"stage": "Stage 6", "text": "Fragmented stage six ownership delays decision support."},
            ],
            "risks": [
                {"stage": 1, "text": "Fragmented stage one ownership delays decision support."},
                {"stage": 2, "text": "Fragmented stage two ownership delays decision support."},
                {"stage": 3, "text": "Fragmented stage three ownership delays decision support."},
                {"stage": 4, "text": "Fragmented stage four ownership delays decision support."},
                {"stage": 5, "text": "Fragmented stage five ownership delays decision support."},
                {"stage": 6, "text": "Fragmented stage six ownership delays decision support."},
            ],
            "metadata": {"company": "Toyota", "industry": "Automotive", "process": "Procure-to-Pay"},
        }

    def _process_flow_spec(self, title: str = "Process Flow") -> dict:
        return {
            "title": title,
            "subtitle": "Toyota Procurement Process Flow",
            "description": "Procure-to-Pay process flow.",
            "nodes": [
                {"id": "req", "label": "Requisition"},
                {"id": "po", "label": "Purchase Order"},
                {"id": "pay", "label": "Payment"},
            ],
            "connections": [
                {"source": "req", "target": "po"},
                {"source": "po", "target": "pay"},
            ],
            "pain_points": [
                {"stage": "Requisition", "text": "Manual requisition entry creates delays."},
                {"stage": "Payment", "text": "Late payments damage supplier relationships."},
            ],
        }

    def test_operating_model_renderer_creates_single_slide(self):
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = os.path.join(tmpdir, "om.pptx")
            renderer = OperatingModelRenderer()
            renderer.render(self._operating_model_spec(), output_path)

            self.assertTrue(os.path.exists(output_path))
            prs = Presentation(output_path)
            self.assertEqual(len(prs.slides), 1)

    def test_process_flow_renderer_creates_single_slide(self):
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = os.path.join(tmpdir, "pf.pptx")
            renderer = ProcessFlowRenderer()
            renderer.render(self._process_flow_spec(), output_path)

            self.assertTrue(os.path.exists(output_path))
            prs = Presentation(output_path)
            self.assertEqual(len(prs.slides), 1)

    def test_operating_model_renderer_appends_to_existing_presentation(self):
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = os.path.join(tmpdir, "multi.pptx")
            renderer = OperatingModelRenderer()
            shared_prs = Presentation()

            returned_prs = renderer.render(
                self._operating_model_spec("Slide 1"), output_path, presentation=shared_prs
            )
            self.assertEqual(len(returned_prs.slides), 1)

            returned_prs = renderer.render(
                self._operating_model_spec("Slide 2"), output_path, presentation=shared_prs
            )
            self.assertEqual(len(returned_prs.slides), 2)
            self.assertIs(returned_prs, shared_prs)

            # No save occurred during append calls.
            self.assertFalse(os.path.exists(output_path))

            returned_prs.save(output_path)
            self.assertTrue(os.path.exists(output_path))
            prs = Presentation(output_path)
            self.assertEqual(len(prs.slides), 2)

    def test_process_flow_renderer_appends_to_existing_presentation(self):
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = os.path.join(tmpdir, "multi.pptx")
            renderer = ProcessFlowRenderer()
            shared_prs = Presentation()

            returned_prs = renderer.render(
                self._process_flow_spec("Slide 1"), output_path, presentation=shared_prs
            )
            self.assertEqual(len(returned_prs.slides), 1)

            returned_prs = renderer.render(
                self._process_flow_spec("Slide 2"), output_path, presentation=shared_prs
            )
            self.assertEqual(len(returned_prs.slides), 2)
            self.assertIs(returned_prs, shared_prs)

            returned_prs.save(output_path)
            prs = Presentation(output_path)
            self.assertEqual(len(prs.slides), 2)

    def test_mixed_renderers_append_to_shared_presentation(self):
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = os.path.join(tmpdir, "mixed.pptx")
            om_renderer = OperatingModelRenderer()
            pf_renderer = ProcessFlowRenderer()
            shared_prs = Presentation()

            om_renderer.render(
                self._operating_model_spec("Operating Model"), output_path, presentation=shared_prs
            )
            pf_renderer.render(
                self._process_flow_spec("Process Flow"), output_path, presentation=shared_prs
            )

            shared_prs.save(output_path)
            prs = Presentation(output_path)
            self.assertEqual(len(prs.slides), 2)


if __name__ == "__main__":
    unittest.main()
