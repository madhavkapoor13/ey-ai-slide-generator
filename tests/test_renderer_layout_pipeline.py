import os
import tempfile
import unittest

from pptx import Presentation
from pptx.util import Inches

from backend.layout_engine.layout_engine import generate_layout
from ppt_renderer.components import component_dispatcher
from ppt_renderer.components.card_renderer import render as render_card
from ppt_renderer.components.coordinates import convert_bounds, convert_height, convert_width, convert_x, convert_y
from ppt_renderer.components.header_renderer import render_header
from ppt_renderer.components.matrix_renderer import render as render_matrix
from ppt_renderer.components.timeline_renderer import render as render_timeline
from ppt_renderer.operating_model_renderer import OperatingModelRenderer
from ppt_renderer.renderer import ProcessFlowRenderer
from schemas.layout import (
    BodySpecification,
    ComponentSpecification,
    FooterSpecification,
    HeaderSpecification,
    LayoutSpecification,
)
from schemas.visual import VisualPatternSelection


def _slide_spec():
    return {
        "title": "Test Slide",
        "subtitle": "Subtitle",
        "description": "Description",
    }


class RendererLayoutPipelineTests(unittest.TestCase):

    def test_legacy_process_flow_renderer_still_works(self):
        renderer = ProcessFlowRenderer()
        spec = {
            "title": "Procure-to-Pay",
            "subtitle": "Current state",
            "nodes": [
                {"id": "1", "label": "Receive"},
                {"id": "2", "label": "Validate"},
                {"id": "3", "label": "Approve"},
            ],
            "connections": [{"from": "1", "to": "2"}, {"from": "2", "to": "3"}],
            "pain_points": [],
        }

        with tempfile.TemporaryDirectory() as tmp:
            path = os.path.join(tmp, "legacy.pptx")
            prs = renderer.render(spec, output_path=path)

            self.assertTrue(os.path.exists(path))
            self.assertEqual(len(prs.slides), 1)
            self.assertGreater(len(prs.slides[0].shapes), 0)

    def test_legacy_operating_model_renderer_still_works(self):
        renderer = OperatingModelRenderer()
        spec = {
            "title": "Operating Model",
            "subtitle": "Current State",
            "summary": {"headline": "Summary", "description": "Desc", "metrics": []},
            "stages": [
                {"number": 1, "title": "Stage 1", "activities": ["A", "B"]},
                {"number": 2, "title": "Stage 2", "activities": ["C", "D"]},
            ],
            "risks": [{"stage": 1, "text": "Risk"}],
        }

        with tempfile.TemporaryDirectory() as tmp:
            path = os.path.join(tmp, "legacy_om.pptx")
            prs = renderer.render(spec, output_path=path)

            self.assertTrue(os.path.exists(path))
            self.assertEqual(len(prs.slides), 1)
            self.assertGreater(len(prs.slides[0].shapes), 0)

    def test_layout_specification_rendering_works(self):
        selection = VisualPatternSelection(
            pattern_id="CL-01",
            category="creative_listing",
            confidence=0.9,
            reasoning="Four insight cards",
        )
        layout = generate_layout(selection)

        renderer = ProcessFlowRenderer()
        content = {
            "title": "Trends",
            "subtitle": "Megatrends",
            "cards": [
                {"title": "Digital", "description": "..."},
                {"title": "Regulation", "description": "..."},
                {"title": "Talent", "description": "..."},
                {"title": "Sustainability", "description": "..."},
            ],
        }

        with tempfile.TemporaryDirectory() as tmp:
            path = os.path.join(tmp, "layout.pptx")
            prs = renderer.render(content, output_path=path, layout_spec=layout)

            self.assertTrue(os.path.exists(path))
            self.assertEqual(len(prs.slides), 1)
            # Header, divider, 4 cards, footer rule + 3 footer textboxes = 10 shapes.
            self.assertGreaterEqual(len(prs.slides[0].shapes), 8)

    def test_component_dispatcher_selects_correct_renderer(self):
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        content = {
            "cards": [{"title": "Card Title", "description": "Card body"}],
            "steps": [{"label": "Step 1"}],
            "cells": [{"value": "Cell"}],
        }

        component_dispatcher.render(
            ComponentSpecification(
                component_id="c1",
                type="card",
                x=0.1,
                y=0.2,
                width=0.2,
                height=0.2,
                placeholder="card_1",
            ),
            prs,
            slide,
            content,
        )
        component_dispatcher.render(
            ComponentSpecification(
                component_id="n1", type="node", x=0.4, y=0.2, width=0.1, height=0.1, placeholder="step_1"
            ),
            prs,
            slide,
            content,
        )
        component_dispatcher.render(
            ComponentSpecification(
                component_id="m1", type="cell", x=0.6, y=0.2, width=0.1, height=0.1, placeholder="cell_1"
            ),
            prs,
            slide,
            content,
        )

        self.assertGreaterEqual(len(slide.shapes), 3)

    def test_normalized_coordinates_converted_correctly(self):
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        component = ComponentSpecification(
            component_id="box",
            type="card",
            x=0.5,
            y=0.25,
            width=0.4,
            height=0.2,
            placeholder="card_1",
        )

        left, top, width, height = convert_bounds(component, prs)

        self.assertAlmostEqual(left.inches, 13.333 * 0.5, places=3)
        self.assertAlmostEqual(top.inches, 7.5 * 0.25, places=3)
        self.assertAlmostEqual(width.inches, 13.333 * 0.4, places=3)
        self.assertAlmostEqual(height.inches, 7.5 * 0.2, places=3)

    def test_header_rendering(self):
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        header = HeaderSpecification(
            height=0.15,
            title_area={"x": 0.05, "y": 0.04, "width": 0.9, "height": 0.06},
            subtitle_area={"x": 0.05, "y": 0.10, "width": 0.9, "height": 0.04},
        )
        content = {"title": "Header Title", "subtitle": "Header Subtitle"}

        render_header(header, prs, slide, content)

        texts = [shape.text_frame.text for shape in slide.shapes if shape.has_text_frame]
        self.assertIn("Header Title", texts)
        self.assertIn("Header Subtitle", texts)

    def test_header_subtitle_and_description_do_not_overlap(self):
        """Regression: subtitle and description must render at distinct y
        positions, never on top of each other. Pre-fix they shared the
        ``subtitle_area`` and overlapped on every slide."""
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        header = HeaderSpecification(
            height=0.15,
            title_area={"x": 0.05, "y": 0.04, "width": 0.9, "height": 0.06},
            subtitle_area={"x": 0.05, "y": 0.10, "width": 0.9, "height": 0.04},
        )
        content = {"title": "T", "subtitle": "S", "description": "D"}

        render_header(header, prs, slide, content)

        text_boxes = {}
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text
            if text in ("S", "D"):
                text_boxes[text] = (shape.top, shape.height)
        self.assertIn("S", text_boxes)
        self.assertIn("D", text_boxes)
        sub_top, sub_h = text_boxes["S"]
        desc_top, _ = text_boxes["D"]
        # Description starts strictly below the end of the subtitle.
        self.assertGreater(desc_top, sub_top + sub_h)

    def test_card_rendering(self):
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        component = ComponentSpecification(
            component_id="card_1",
            type="card",
            x=0.1,
            y=0.2,
            width=0.2,
            height=0.2,
            placeholder="card_1",
        )
        content = {"cards": [{"title": "Insight 1", "description": "Description"}]}

        render_card(component, prs, slide, content)

        texts = [shape.text_frame.text for shape in slide.shapes if shape.has_text_frame]
        self.assertTrue(any("Insight 1" in text for text in texts))

    def test_timeline_rendering(self):
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        component = ComponentSpecification(
            component_id="step_1",
            type="node",
            x=0.1,
            y=0.3,
            width=0.15,
            height=0.15,
            placeholder="step_1",
        )
        content = {"steps": [{"label": "Receive"}]}

        render_timeline(component, prs, slide, content)

        texts = [shape.text_frame.text for shape in slide.shapes if shape.has_text_frame]
        self.assertTrue(any("Receive" in text for text in texts))

    def test_matrix_rendering(self):
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        component = ComponentSpecification(
            component_id="cell_1",
            type="cell",
            x=0.1,
            y=0.3,
            width=0.15,
            height=0.15,
            placeholder="cell_1",
        )
        content = {"cells": [{"value": "High Impact"}]}

        render_matrix(component, prs, slide, content)

        texts = [shape.text_frame.text for shape in slide.shapes if shape.has_text_frame]
        self.assertTrue(any("High Impact" in text for text in texts))

    def test_unknown_component_ignored_safely(self):
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        before = len(slide.shapes)
        component = ComponentSpecification(
            component_id="unknown_1",
            type="alien_shape",
            x=0.1,
            y=0.1,
            width=0.1,
            height=0.1,
            placeholder="x",
        )

        component_dispatcher.render(component, prs, slide, {})

        self.assertEqual(len(slide.shapes), before)

    def test_multiple_components_render_in_sequence(self):
        selection = VisualPatternSelection(
            pattern_id="CL-02",
            category="creative_listing",
            confidence=0.9,
            reasoning="Three strategy cards",
        )
        layout = generate_layout(selection)

        renderer = ProcessFlowRenderer()
        content = {
            "title": "Strategies",
            "cards": [
                {"title": "A", "description": "..."},
                {"title": "B", "description": "..."},
                {"title": "C", "description": "..."},
            ],
        }

        with tempfile.TemporaryDirectory() as tmp:
            path = os.path.join(tmp, "multi.pptx")
            prs = renderer.render(content, output_path=path, layout_spec=layout)

            shapes = list(prs.slides[0].shapes)
            self.assertGreaterEqual(len(shapes), 3)  # at least the three cards

    def test_backward_compatibility_public_api_unchanged(self):
        renderer = ProcessFlowRenderer()
        spec = {
            "title": "Legacy",
            "subtitle": "Test",
            "nodes": [{"id": "1", "label": "One"}],
            "connections": [],
            "pain_points": [],
        }

        with tempfile.TemporaryDirectory() as tmp:
            path = os.path.join(tmp, "compat.pptx")
            # Positional arguments only — no layout_spec.
            prs = renderer.render(spec, path)
            self.assertEqual(len(prs.slides), 1)


if __name__ == "__main__":
    unittest.main()
