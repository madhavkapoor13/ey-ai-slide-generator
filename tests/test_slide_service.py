from __future__ import annotations

import os
import unittest
from unittest.mock import MagicMock, patch

from pptx import Presentation

from backend.services import slide_service
from schemas.clarification import ClarificationQuestion, ClarificationResult
from schemas.deck_execution import DeckExecutionResult, SlideExecutionResult
from schemas.executive_card import ExecutiveCardContent
from schemas.layout import (
    BodySpecification,
    ComponentSpecification,
    FooterSpecification,
    HeaderSpecification,
    LayoutSpecification,
)
from schemas.pipeline_result import PipelineResult
from schemas.presentation import DeckSpec, SlidePlan
from schemas.slide_spec import SlideSpec
from schemas.visual import VisualPatternSelection


class SlideServiceTests(unittest.TestCase):
    def setUp(self):
        # Remove any leftover generated files from previous runs.
        for path in ["generated_slide_v2.pptx"]:
            if os.path.exists(path):
                os.remove(path)

    def tearDown(self):
        for path in ["generated_slide_v2.pptx"]:
            if os.path.exists(path):
                os.remove(path)

    def _deck_spec(self, slide_roles: list[str]) -> DeckSpec:
        return DeckSpec(
            presentation_type="Transformation Proposal",
            objective="Transform procurement.",
            audience="Senior leadership",
            narrative=" → ".join(slide_roles),
            estimated_slide_count=len(slide_roles),
            slides=[
                SlidePlan(
                    slide_number=index + 1,
                    slide_role=role,
                    purpose=f"Communicate {role}.",
                    required_inputs=[],
                    dependencies=[],
                    visualization_type="Executive Summary" if role == "Executive Summary" else "Process Flow",
                )
                for index, role in enumerate(slide_roles)
            ],
        )

    def _slide_spec(self, slide_role: str, slide_type: str = "operating_model") -> SlideSpec:
        return SlideSpec(
            slide_type=slide_type,
            raw_spec={
                "title": slide_role,
                "subtitle": f"Subtitle for {slide_role}",
                "description": f"Description for {slide_role}",
                "executive_summary": "Summary one. Summary two.",
                "summary": {"headline": slide_role, "description": "Summary one. Summary two.", "metrics": []},
                "stages": [
                    {"number": i + 1, "title": f"Stage {i + 1}", "label": f"Stage {i + 1}", "activities": ["A", "B"]}
                    for i in range(6)
                ],
                "pain_points": [
                    {"stage": f"Stage {i + 1}", "text": f"Pain point for stage {i + 1}."}
                    for i in range(6)
                ],
                "risks": [
                    {"stage": i + 1, "text": f"Risk for stage {i + 1}."}
                    for i in range(6)
                ],
                "metadata": {"slide_role": slide_role},
            },
            version="2.0",
            generated_by="test",
        )

    def _cl01_slide_spec(self, source: str = "business_benefits") -> SlideSpec:
        """Return a SlideSpec whose raw content exercises CL-01 card derivation."""
        if source == "cards":
            raw_spec = {
                "title": "Key Business Benefits",
                "subtitle": "Microsoft AI Procurement Transformation",
                "cards": [
                    {"title": "Cost Reduction", "description": "Reduce spend."},
                    {"title": "Compliance", "description": "Improve policy adherence."},
                    {"title": "Speed", "description": "Faster cycle times."},
                    {"title": "Visibility", "description": "Real-time dashboards."},
                ],
            }
        else:
            raw_spec = {
                "title": "Key Business Benefits",
                "subtitle": "Microsoft AI Procurement Transformation",
                source: [
                    {"title": "Cost Reduction", "description": "Reduce spend."},
                    {"title": "Compliance", "description": "Improve policy adherence."},
                    {"title": "Speed", "description": "Faster cycle times."},
                    {"title": "Visibility", "description": "Real-time dashboards."},
                ],
            }
        return SlideSpec(
            slide_type="operating_model",
            raw_spec=raw_spec,
            version="2.0",
            generated_by="test",
        )

    def _cl02_slide_spec(self, source: str = "key_insights") -> SlideSpec:
        """Return a SlideSpec whose raw content exercises CL-02 card derivation."""
        if source == "cards":
            raw_spec = {
                "title": "Top Three Strategic Priorities",
                "subtitle": "AI Procurement Transformation",
                "cards": [
                    {"title": "Automate", "description": "Automate sourcing."},
                    {"title": "Integrate", "description": "Integrate systems."},
                    {"title": "Scale", "description": "Scale adoption."},
                ],
            }
        else:
            raw_spec = {
                "title": "Top Three Strategic Priorities",
                "subtitle": "AI Procurement Transformation",
                source: [
                    {"title": "Automate", "description": "Automate sourcing."},
                    {"title": "Integrate", "description": "Integrate systems."},
                    {"title": "Scale", "description": "Scale adoption."},
                ],
            }
        return SlideSpec(
            slide_type="operating_model",
            raw_spec=raw_spec,
            version="2.0",
            generated_by="test",
        )

    def _cl03_slide_spec(self, source: str = "kpis") -> SlideSpec:
        """Return a SlideSpec whose raw content exercises CL-03 KPI derivation."""
        if source == "kpis":
            raw_spec = {
                "title": "Key KPIs",
                "subtitle": "Microsoft AI Procurement Transformation",
                "kpis": [
                    {"label": "Cost Reduction", "value": "15%"},
                    {"label": "Cycle Time", "value": "-30%"},
                    {"label": "Compliance", "value": "98%"},
                ],
            }
        elif source == "metrics":
            raw_spec = {
                "title": "Key KPIs",
                "subtitle": "Microsoft AI Procurement Transformation",
                "metrics": [
                    {"label": "Cost Reduction", "value": "15%"},
                    {"label": "Cycle Time", "value": "-30%"},
                    {"label": "Compliance", "value": "98%"},
                ],
            }
        else:
            raw_spec = {
                "title": "Key KPIs",
                "subtitle": "Microsoft AI Procurement Transformation",
                source: [
                    {"title": "Cost Reduction", "value": "15%"},
                    {"title": "Cycle Time", "value": "-30%"},
                    {"title": "Compliance", "value": "98%"},
                ],
            }
        return SlideSpec(
            slide_type="operating_model",
            raw_spec=raw_spec,
            version="2.0",
            generated_by="test",
        )

    def _cl04_slide_spec(self, source: str = "columns") -> SlideSpec:
        """Return a SlideSpec whose raw content exercises CL-04 comparison derivation."""
        if source == "columns":
            raw_spec = {
                "title": "Current vs Future",
                "subtitle": "Procurement Operating Model",
                "columns": [
                    {
                        "label": "Current",
                        "items": [
                            {"name": "Manual", "text": "Manual processes"},
                            {"name": "Siloed", "text": "Siloed data"},
                            {"name": "Slow", "text": "Slow cycle times"},
                            {"name": "Reactive", "text": "Reactive sourcing"},
                        ],
                    },
                    {
                        "label": "Future",
                        "items": [
                            {"name": "Automated", "text": "Automated workflows"},
                            {"name": "Integrated", "text": "Integrated data"},
                            {"name": "Fast", "text": "Fast cycle times"},
                            {"name": "Proactive", "text": "Proactive sourcing"},
                        ],
                    },
                ],
            }
        else:
            raw_spec = {
                "title": "Current vs Future",
                "subtitle": "Procurement Operating Model",
                source: [
                    {"name": "Manual", "text": "Manual processes"},
                    {"name": "Siloed", "text": "Siloed data"},
                    {"name": "Slow", "text": "Slow cycle times"},
                    {"name": "Reactive", "text": "Reactive sourcing"},
                ],
            }
        return SlideSpec(
            slide_type="operating_model",
            raw_spec=raw_spec,
            version="2.0",
            generated_by="test",
        )

    def _cl05_slide_spec(self, source: str = "columns") -> SlideSpec:
        """Return a SlideSpec whose raw content exercises CL-05 two-column derivation."""
        if source == "columns":
            raw_spec = {
                "title": "Challenges and Recommendations",
                "subtitle": "Procurement Transformation",
                "columns": [
                    {
                        "label": "Challenges",
                        "items": [
                            {"text": "Fragmented processes"},
                            {"text": "Limited visibility"},
                            {"text": "Manual compliance checks"},
                            {"text": "Slow vendor onboarding"},
                        ],
                    },
                    {
                        "label": "Recommendations",
                        "items": [
                            {"text": "Unified platform"},
                            {"text": "Real-time dashboards"},
                            {"text": "Automated policy enforcement"},
                            {"text": "Self-service portal"},
                        ],
                    },
                ],
            }
        else:
            raw_spec = {
                "title": "Challenges and Recommendations",
                "subtitle": "Procurement Transformation",
                "challenges": ["Fragmented processes", "Limited visibility"],
                "solutions": ["Unified platform", "Real-time dashboards"],
            }
        return SlideSpec(
            slide_type="operating_model",
            raw_spec=raw_spec,
            version="2.0",
            generated_by="test",
        )

    def _ig01_slide_spec(self) -> SlideSpec:
        return SlideSpec(
            slide_type="process_flow",
            raw_spec={
                "title": "Implementation Timeline",
                "subtitle": "Procurement Transformation",
                "timeline": [
                    {"date": "Q1", "title": "Assess", "description": "Current state"},
                    {"date": "Q2", "title": "Design", "description": "Target model"},
                    {"date": "Q3", "title": "Build", "description": "Platform rollout"},
                    {"date": "Q4", "title": "Scale", "description": "Adoption"},
                ],
            },
            version="2.0",
            generated_by="test",
        )

    def _ig02_slide_spec(self) -> SlideSpec:
        return SlideSpec(
            slide_type="process_flow",
            raw_spec={
                "title": "Transformation Roadmap",
                "subtitle": "Phased rollout",
                "phases": [
                    {"name": "Phase 1", "activities": ["Assess"]},
                    {"name": "Phase 2", "activities": ["Design"]},
                    {"name": "Phase 3", "activities": ["Build"]},
                    {"name": "Phase 4", "activities": ["Scale"]},
                ],
            },
            version="2.0",
            generated_by="test",
        )

    def _ig03_slide_spec(self) -> SlideSpec:
        return SlideSpec(
            slide_type="process_flow",
            raw_spec={
                "title": "Procurement Process Flow",
                "subtitle": "Procure-to-pay",
                "steps": [
                    {"label": "Request"},
                    {"label": "Approve"},
                    {"label": "Purchase"},
                    {"label": "Receive"},
                    {"label": "Pay"},
                ],
            },
            version="2.0",
            generated_by="test",
        )

    def _ig04_slide_spec(self) -> SlideSpec:
        return SlideSpec(
            slide_type="operating_model",
            raw_spec={
                "title": "Impact vs Effort Matrix",
                "subtitle": "Initiative prioritization",
                "rows": [
                    {
                        "label": "High Impact",
                        "cells": [
                            {"value": "Quick wins"},
                            {"value": "Major projects"},
                            {"value": "Strategic bets"},
                        ],
                    },
                    {
                        "label": "Low Impact",
                        "cells": [
                            {"value": "Fill-ins"},
                            {"value": "Avoid"},
                            {"value": "Deprioritize"},
                        ],
                    },
                    {
                        "label": "Medium",
                        "cells": [
                            {"value": "Consider"},
                            {"value": "Evaluate"},
                            {"value": "Plan"},
                        ],
                    },
                ],
            },
            version="2.0",
            generated_by="test",
        )

    def _ig05_slide_spec(self) -> SlideSpec:
        return SlideSpec(
            slide_type="process_flow",
            raw_spec={
                "title": "Customer Journey",
                "subtitle": "Procure-to-pay experience",
                "journey_stages": [
                    {"name": "Aware", "touchpoints": ["Campaign"]},
                    {"name": "Engage", "touchpoints": ["Workshop"]},
                    {"name": "Adopt", "touchpoints": ["Training"]},
                    {"name": "Optimize", "touchpoints": ["Review"]},
                ],
            },
            version="2.0",
            generated_by="test",
        )

    def _ig06_slide_spec(self) -> SlideSpec:
        return SlideSpec(
            slide_type="operating_model",
            raw_spec={
                "title": "Capability Map",
                "subtitle": "Procurement capabilities",
                "domains": [
                    {
                        "name": "Strategy",
                        "capabilities": [
                            {"name": "Demand planning"},
                            {"name": "Supplier strategy"},
                        ],
                    },
                    {
                        "name": "Sourcing",
                        "capabilities": [
                            {"name": "RFx"},
                            {"name": "Negotiation"},
                        ],
                    },
                    {
                        "name": "Operations",
                        "capabilities": [
                            {"name": "Purchase orders"},
                            {"name": "Invoice matching"},
                        ],
                    },
                    {
                        "name": "Analytics",
                        "capabilities": [
                            {"name": "Spend visibility"},
                            {"name": "Performance mgmt"},
                        ],
                    },
                ],
            },
            version="2.0",
            generated_by="test",
        )

    def _process_slide_spec(self, slide_role: str) -> SlideSpec:
        return SlideSpec(
            slide_type="process_flow",
            raw_spec={
                "title": slide_role,
                "subtitle": f"Subtitle for {slide_role}",
                "description": f"Description for {slide_role}",
                "nodes": [{"id": "a", "label": "A"}, {"id": "b", "label": "B"}],
                "connections": [{"source": "a", "target": "b"}],
                "pain_points": [{"stage": "A", "text": "Pain point."}],
                "metadata": {"slide_role": slide_role},
            },
            version="2.0",
            generated_by="test",
        )

    def _execution_result(
        self,
        slide_roles: list[str],
        slide_types: list[str] | None = None,
        failed_roles: list[str] | None = None,
    ) -> DeckExecutionResult:
        slide_types = slide_types or ["operating_model"] * len(slide_roles)
        failed_roles = failed_roles or []
        deck_spec = self._deck_spec(slide_roles + failed_roles)

        slides: list[SlideExecutionResult] = []
        successful_slides: list[SlideSpec] = []
        failed_slides: list[SlideExecutionResult] = []

        for role, slide_type in zip(slide_roles, slide_types):
            spec = (
                self._process_slide_spec(role)
                if slide_type == "process_flow"
                else self._slide_spec(role, slide_type)
            )
            result = SlideExecutionResult(
                slide_plan=deck_spec.slides[len(slides)],
                slide_spec=spec,
                validation_result=None,
                success=True,
            )
            slides.append(result)
            successful_slides.append(spec)

        for role in failed_roles:
            plan = deck_spec.slides[len(slides)]
            result = SlideExecutionResult(
                slide_plan=plan,
                slide_spec=None,
                validation_result=None,
                success=False,
                error=f"Failed to generate {role}.",
            )
            slides.append(result)
            failed_slides.append(result)

        total = len(slides)
        return DeckExecutionResult(
            deck_spec=deck_spec,
            slides=slides,
            successful_slides=successful_slides,
            failed_slides=failed_slides,
            all_succeeded=total > 0 and len(failed_slides) == 0,
            partial_success=len(successful_slides) > 0 and len(failed_slides) > 0,
        )

    def _completed_result(
        self,
        slide_roles: list[str],
        slide_types: list[str] | None = None,
        failed_roles: list[str] | None = None,
    ) -> PipelineResult:
        return PipelineResult(
            status="COMPLETED",
            needs_clarification=False,
            clarification_result=None,
            deck_execution_result=self._execution_result(slide_roles, slide_types, failed_roles),
            warnings=[],
        )

    def _completed_result_with_specs(self, specs: list[SlideSpec]) -> PipelineResult:
        """Build a completed PipelineResult from explicit SlideSpec objects."""
        deck_spec = DeckSpec(
            presentation_type="Transformation Proposal",
            objective="Transform procurement.",
            audience="Senior leadership",
            narrative=" → ".join(spec.raw_spec.get("title", f"Slide {i}") for i, spec in enumerate(specs, 1)),
            estimated_slide_count=len(specs),
            slides=[
                SlidePlan(
                    slide_number=index + 1,
                    slide_role=spec.raw_spec.get("title", "Slide"),
                    purpose=f"Communicate {spec.raw_spec.get('title', 'slide')}.",
                    required_inputs=[],
                    dependencies=[],
                    visualization_type="Insight Cards",
                )
                for index, spec in enumerate(specs)
            ],
        )
        slides: list[SlideExecutionResult] = []
        for index, spec in enumerate(specs):
            slides.append(
                SlideExecutionResult(
                    slide_plan=deck_spec.slides[index],
                    slide_spec=spec,
                    validation_result=None,
                    success=True,
                )
            )
        deck_result = DeckExecutionResult(
            deck_spec=deck_spec,
            slides=slides,
            successful_slides=specs,
            failed_slides=[],
            all_succeeded=True,
            partial_success=False,
        )
        return PipelineResult(
            status="COMPLETED",
            needs_clarification=False,
            clarification_result=None,
            deck_execution_result=deck_result,
            warnings=[],
        )

    def _waiting_result(self, content_questions: list[ClarificationQuestion]) -> PipelineResult:
        return PipelineResult(
            status="WAITING_FOR_USER",
            needs_clarification=True,
            clarification_result=ClarificationResult(
                needs_clarification=True,
                content_questions=content_questions,
                visualization_questions=[],
            ),
            deck_execution_result=None,
            warnings=["Missing information."],
        )

    def _assert_slide_count(self, expected: int):
        self.assertTrue(os.path.exists("generated_slide_v2.pptx"))
        prs = Presentation("generated_slide_v2.pptx")
        self.assertEqual(len(prs.slides), expected)

    def _visual_selection(self, pattern_id: str) -> VisualPatternSelection:
        return VisualPatternSelection(
            pattern_id=pattern_id,
            category="creative_listing",
            confidence=0.95,
            reasoning=f"Selected {pattern_id} for testing.",
        )

    def _cl06_layout(self) -> LayoutSpecification:
        return LayoutSpecification(
            layout_id="CL06",
            visual_pattern="CL-06",
            category="creative_listing",
            canvas_type="widescreen_16_9",
            header=HeaderSpecification(height=0.15),
            body=BodySpecification(x=0.05, y=0.18, width=0.9, height=0.62),
            footer=FooterSpecification(height=0.10),
            components=[
                ComponentSpecification(
                    component_id="card_1",
                    type="executive_card",
                    x=0.07,
                    y=0.22,
                    width=0.28,
                    height=0.54,
                    placeholder="card_1",
                    constraints=["title_bold", "icon_top", "max_4_lines", "stat_optional"],
                ),
                ComponentSpecification(
                    component_id="card_2",
                    type="executive_card",
                    x=0.36,
                    y=0.22,
                    width=0.28,
                    height=0.54,
                    placeholder="card_2",
                    constraints=["title_bold", "icon_top", "max_4_lines", "stat_optional"],
                ),
                ComponentSpecification(
                    component_id="card_3",
                    type="executive_card",
                    x=0.65,
                    y=0.22,
                    width=0.28,
                    height=0.54,
                    placeholder="card_3",
                    constraints=["title_bold", "icon_top", "max_4_lines", "stat_optional"],
                ),
            ],
            spacing="equal",
            alignment="center",
            supports_images=True,
            supports_percentages=True,
        )

    def test_generate_slide_v2_always_returns_string_path(self):
        result = self._completed_result(["Executive Summary"])

        with patch("backend.orchestrator.run_pipeline", return_value=result):
            path = slide_service.generate_slide_v2("Title", "Content")

        self.assertIsInstance(path, str)
        self.assertEqual(path, "generated_slide_v2.pptx")

    def test_one_slide_deck(self):
        result = self._completed_result(["Executive Summary"])

        with patch("backend.orchestrator.run_pipeline", return_value=result):
            path = slide_service.generate_slide_v2("Title", "Content")

        self.assertEqual(path, "generated_slide_v2.pptx")
        self._assert_slide_count(1)

    def test_six_slide_deck(self):
        roles = ["Executive Summary", "Current State", "Opportunities", "Future State", "Roadmap", "Next Steps"]
        result = self._completed_result(roles)

        with patch("backend.orchestrator.run_pipeline", return_value=result):
            path = slide_service.generate_slide_v2("Title", "Content")

        self.assertEqual(path, "generated_slide_v2.pptx")
        self._assert_slide_count(6)

    def test_ten_slide_deck(self):
        roles = [f"Slide {i}" for i in range(1, 11)]
        result = self._completed_result(roles)

        with patch("backend.orchestrator.run_pipeline", return_value=result):
            path = slide_service.generate_slide_v2("Title", "Content")

        self.assertEqual(path, "generated_slide_v2.pptx")
        self._assert_slide_count(10)

    def test_mixed_renderer_deck(self):
        roles = ["Current State", "Process Flow"]
        result = self._completed_result(roles, slide_types=["operating_model", "process_flow"])

        with patch("backend.orchestrator.run_pipeline", return_value=result):
            path = slide_service.generate_slide_v2("Title", "Content")

        self.assertEqual(path, "generated_slide_v2.pptx")
        self._assert_slide_count(2)

    def test_partial_failures_render_only_successful_slides(self):
        result = self._completed_result(
            ["Executive Summary", "Future State"],
            failed_roles=["Current State"],
        )

        with patch("backend.orchestrator.run_pipeline", return_value=result):
            path = slide_service.generate_slide_v2("Title", "Content")

        self.assertEqual(path, "generated_slide_v2.pptx")
        self._assert_slide_count(2)

    def test_all_slides_failed_returns_placeholder_deck(self):
        result = self._completed_result([], failed_roles=["Executive Summary", "Current State"])

        with patch("backend.orchestrator.run_pipeline", return_value=result):
            path = slide_service.generate_slide_v2("Title", "Content")

        self.assertEqual(path, "generated_slide_v2.pptx")
        self._assert_slide_count(1)
        prs = Presentation("generated_slide_v2.pptx")
        # Blank layout may not expose a title placeholder, so collect all text frame text.
        slide_text = ""
        for shape in prs.slides[0].shapes:
            if hasattr(shape, "text_frame"):
                slide_text += shape.text_frame.text + " "
        self.assertIn("Generation Failed", slide_text)

    def test_clarification_returns_placeholder_pptx(self):
        """When clarification is needed, generate_slide_v2 still returns a PPTX path."""
        result = self._waiting_result([
            ClarificationQuestion(
                id="company",
                category="content",
                question="Which company or client is this deck for?",
                required=True,
                reason="The request does not name a company or client.",
            )
        ])

        with patch("backend.orchestrator.run_pipeline", return_value=result):
            path = slide_service.generate_slide_v2("Title", "Content")

        self.assertIsInstance(path, str)
        self.assertEqual(path, "generated_slide_v2.pptx")
        self._assert_slide_count(1)

        prs = Presentation("generated_slide_v2.pptx")
        slide_text = ""
        for shape in prs.slides[0].shapes:
            if hasattr(shape, "text_frame"):
                slide_text += shape.text_frame.text + " "
        self.assertIn("More Information Needed", slide_text)
        self.assertIn("Content Questions", slide_text)
        self.assertIn("Which company or client is this deck for?", slide_text)

    def test_narrative_continuity_preserves_slide_order(self):
        """Slides are rendered in DeckSpec order and the narrative arc is preserved."""
        roles = ["Executive Summary", "Current State", "Opportunities", "Future State", "Roadmap", "Next Steps"]
        result = self._completed_result(roles)

        with patch("backend.orchestrator.run_pipeline", return_value=result):
            slide_service.generate_slide_v2("Title", "Content")

        prs = Presentation("generated_slide_v2.pptx")
        self.assertEqual(len(prs.slides), len(roles))

        # Collect slide text and verify the DeckSpec order is reflected.
        slide_texts: list[str] = []
        for slide in prs.slides:
            text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    text += shape.text_frame.text + " "
            slide_texts.append(text)

        for index, role in enumerate(roles):
            self.assertIn(role, slide_texts[index])

    def test_carried_visual_pattern_id_used_without_re_scoring(self):
        """A visual_pattern_id carried on the SlideSpec is honored directly —
        slide_service must NOT re-score via plan_visual_pattern (single source
        of truth owned by content generation)."""
        spec = SlideSpec(
            slide_type="operating_model",
            raw_spec={"title": "Roadmap", "subtitle": "x"},
            version="2.0",
            generated_by="test",
            visual_pattern_id="IG-02",
            visual_confidence=0.9,
        )
        with patch("backend.services.slide_service.plan_visual_pattern") as mock_plan:
            selection = slide_service._resolve_visual_selection(spec, SlidePlan(
                slide_number=1, slide_role="Roadmap", purpose="p",
                required_inputs=[], dependencies=[], visualization_type="Roadmap",
            ))
        mock_plan.assert_not_called()
        self.assertEqual(selection.pattern_id, "IG-02")
        self.assertEqual(selection.confidence, 0.9)

    def test_visual_pipeline_used_for_allowed_pattern(self):
        """CL-06 triggers the Visual Layout Engine and renders executive cards."""
        result = self._completed_result(["Executive Summary"])
        layout = self._cl06_layout()

        with patch("backend.orchestrator.run_pipeline", return_value=result), \
             patch("backend.services.slide_service.plan_visual_pattern", return_value=self._visual_selection("CL-06")) as mock_plan, \
             patch("backend.services.slide_service.generate_layout", return_value=layout) as mock_layout:
            path = slide_service.generate_slide_v2("Title", "Content")

        self.assertEqual(path, "generated_slide_v2.pptx")
        self._assert_slide_count(1)
        mock_plan.assert_called_once()
        mock_layout.assert_called_once()

        prs = Presentation("generated_slide_v2.pptx")
        slide_text = ""
        for shape in prs.slides[0].shapes:
            if hasattr(shape, "text_frame"):
                slide_text += shape.text_frame.text + " "
        # The derived cards use the first three stage labels and activities.
        self.assertIn("Stage 1", slide_text)
        self.assertIn("Stage 2", slide_text)
        self.assertIn("Stage 3", slide_text)

    def test_non_allowed_pattern_falls_back_to_legacy(self):
        """A non-allowed pattern skips the layout engine and uses legacy rendering."""
        result = self._completed_result(["Executive Summary"])

        with patch("backend.orchestrator.run_pipeline", return_value=result), \
             patch("backend.services.slide_service.plan_visual_pattern", return_value=self._visual_selection("IG-07")) as mock_plan, \
             patch("backend.services.slide_service.generate_layout") as mock_layout:
            path = slide_service.generate_slide_v2("Title", "Content")

        self.assertEqual(path, "generated_slide_v2.pptx")
        self._assert_slide_count(1)
        mock_plan.assert_called_once()
        mock_layout.assert_not_called()

    def test_visual_pipeline_failure_falls_back_to_legacy(self):
        """Any exception in the visual pipeline is caught and falls back safely."""
        result = self._completed_result(["Executive Summary"])

        with patch("backend.orchestrator.run_pipeline", return_value=result), \
             patch("backend.services.slide_service.plan_visual_pattern", side_effect=RuntimeError("planner failed")) as mock_plan, \
             patch("backend.services.slide_service.generate_layout") as mock_layout:
            path = slide_service.generate_slide_v2("Title", "Content")

        self.assertEqual(path, "generated_slide_v2.pptx")
        self._assert_slide_count(1)
        mock_plan.assert_called_once()
        mock_layout.assert_not_called()

    def test_build_executive_card_content_derives_cards_from_stages(self):
        """When no ``cards`` key exists, three cards are derived from ``stages``."""
        raw_spec = self._slide_spec("Executive Summary").raw_spec
        self.assertNotIn("cards", raw_spec)

        cards = slide_service._build_executive_card_content(raw_spec, "CL-06")

        self.assertEqual(len(cards), 3)
        self.assertEqual(cards[0].title, "Stage 1")
        self.assertEqual(cards[0].description, "A, B")
        # Original spec must not be mutated.
        self.assertNotIn("cards", raw_spec)

    def test_build_executive_card_content_returns_empty_when_no_sources(self):
        """No fallback padding: when no sources exist, the list is empty."""
        raw_spec = {"title": "Summary", "subtitle": "Key takeaways"}
        cards = slide_service._build_executive_card_content(raw_spec, "CL-06")

        self.assertEqual(len(cards), 0)

    def test_build_executive_card_content_preserves_existing_cards(self):
        """An existing ``cards`` key is returned without padding."""
        raw_spec = {"cards": [{"title": "Existing", "description": "Card"}]}
        cards = slide_service._build_executive_card_content(raw_spec, "CL-06")

        self.assertEqual(len(cards), 1)
        self.assertEqual(cards[0].title, "Existing")
        self.assertEqual(cards[0].description, "Card")

    def test_build_executive_card_content_cl01_returns_empty_when_no_sources(self):
        """CL-01 returns an empty list when no sources exist (no padding)."""
        raw_spec = {"title": "Benefits", "subtitle": "Overview"}
        cards = slide_service._build_executive_card_content(raw_spec, "CL-01")

        self.assertEqual(len(cards), 0)

    def test_build_executive_card_content_cl01_priority_order(self):
        """CL-01 derives cards using the documented priority order."""
        raw_spec = {
            "title": "Benefits",
            "subtitle": "Overview",
            "business_benefits": [
                {"title": "Benefit A", "description": "..."},
                {"title": "Benefit B", "description": "..."},
            ],
            "key_insights": [
                {"title": "Insight 1", "description": "..."},
                {"title": "Insight 2", "description": "..."},
            ],
            "recommendations": [
                {"title": "Rec 1", "description": "..."},
            ],
            "stages": [
                {"title": "Stage 1", "activities": ["A"]},
            ],
        }
        cards = slide_service._build_executive_card_content(raw_spec, "CL-01")

        self.assertEqual(len(cards), 4)
        self.assertEqual(cards[0].title, "Benefit A")
        self.assertEqual(cards[1].title, "Benefit B")
        self.assertEqual(cards[2].title, "Insight 1")
        self.assertEqual(cards[3].title, "Insight 2")

    def test_build_executive_card_content_cl01_derives_from_stages_without_padding(self):
        """CL-01 derives cards from stages but does not pad beyond available items."""
        raw_spec = {
            "title": "Benefits",
            "subtitle": "Overview",
            "stages": [
                {"title": "Stage 1", "activities": ["A", "B"]},
                {"title": "Stage 2", "activities": ["C"]},
            ],
        }
        cards = slide_service._build_executive_card_content(raw_spec, "CL-01")

        self.assertEqual(len(cards), 2)
        self.assertEqual(cards[0].title, "Stage 1")
        self.assertEqual(cards[0].description, "A, B")
        self.assertEqual(cards[1].title, "Stage 2")

    def test_build_executive_card_content_does_not_mutate_raw_spec(self):
        """The helper must not mutate the input SlideSpec dict."""
        raw_spec = {"title": "Benefits", "subtitle": "Overview"}
        slide_service._build_executive_card_content(raw_spec, "CL-01")

        self.assertNotIn("cards", raw_spec)

    def test_cl01_enters_visual_pipeline(self):
        """CL-01 routes through the Visual Planner, Layout Engine, and ExecutiveCardRenderer."""
        spec = self._cl01_slide_spec(source="cards")
        result = self._completed_result_with_specs([spec])

        with patch("backend.orchestrator.run_pipeline", return_value=result), \
             patch("backend.services.slide_service.plan_visual_pattern", return_value=self._visual_selection("CL-01")) as mock_plan, \
             patch("backend.services.slide_service.generate_layout") as mock_layout:
            path = slide_service.generate_slide_v2("Title", "Content")

        self.assertEqual(path, "generated_slide_v2.pptx")
        self._assert_slide_count(1)
        mock_plan.assert_called_once()
        mock_layout.assert_called_once()

    def test_cl01_renders_four_cards_with_real_layout(self):
        """CL-01 uses CL01_four_cards.json and renders four executive insight cards."""
        spec = self._cl01_slide_spec(source="business_benefits")
        result = self._completed_result_with_specs([spec])

        with patch("backend.orchestrator.run_pipeline", return_value=result), \
             patch("backend.services.slide_service.plan_visual_pattern", return_value=self._visual_selection("CL-01")) as mock_plan:
            path = slide_service.generate_slide_v2("Title", "Content")

        self.assertEqual(path, "generated_slide_v2.pptx")
        self._assert_slide_count(1)
        mock_plan.assert_called_once()

        prs = Presentation("generated_slide_v2.pptx")
        slide_text = ""
        for shape in prs.slides[0].shapes:
            if hasattr(shape, "text_frame"):
                slide_text += shape.text_frame.text + " "
        for benefit in ("Cost Reduction", "Compliance", "Speed", "Visibility"):
            self.assertIn(benefit, slide_text)

    def test_cl01_fallback_on_renderer_exception(self):
        """Any exception in the CL-01 visual pipeline falls back to the legacy renderer."""
        spec = self._cl01_slide_spec(source="cards")
        result = self._completed_result_with_specs([spec])

        with patch("backend.orchestrator.run_pipeline", return_value=result), \
             patch("backend.services.slide_service.plan_visual_pattern", return_value=self._visual_selection("CL-01")), \
             patch("backend.services.slide_service.generate_layout", side_effect=RuntimeError("layout engine failure")):
            path = slide_service.generate_slide_v2("Title", "Content")

        self.assertEqual(path, "generated_slide_v2.pptx")
        self._assert_slide_count(1)

    def test_mixed_cl01_and_cl06_deck(self):
        """A deck containing both CL-01 and CL-06 slides renders correctly."""
        cl01_spec = self._cl01_slide_spec(source="cards")
        cl06_spec = self._slide_spec("Executive Summary")
        result = self._completed_result_with_specs([cl01_spec, cl06_spec])

        with patch("backend.orchestrator.run_pipeline", return_value=result), \
             patch("backend.services.slide_service.plan_visual_pattern") as mock_plan:
            mock_plan.side_effect = [
                self._visual_selection("CL-01"),
                self._visual_selection("CL-06"),
            ]
            path = slide_service.generate_slide_v2("Title", "Content")

        self.assertEqual(path, "generated_slide_v2.pptx")
        self._assert_slide_count(2)
        self.assertEqual(mock_plan.call_count, 2)

    def test_cl02_enters_visual_pipeline_and_renders_three_cards(self):
        """CL-02 routes through the Layout Engine and renders three executive cards."""
        spec = self._cl02_slide_spec(source="key_insights")
        result = self._completed_result_with_specs([spec])

        with patch("backend.orchestrator.run_pipeline", return_value=result), \
             patch("backend.services.slide_service.plan_visual_pattern", return_value=self._visual_selection("CL-02")) as mock_plan:
            path = slide_service.generate_slide_v2("Title", "Content")

        self.assertEqual(path, "generated_slide_v2.pptx")
        self._assert_slide_count(1)
        mock_plan.assert_called_once()

        prs = Presentation("generated_slide_v2.pptx")
        slide_text = ""
        for shape in prs.slides[0].shapes:
            if hasattr(shape, "text_frame"):
                slide_text += shape.text_frame.text + " "
        for title in ("Automate", "Integrate", "Scale"):
            self.assertIn(title, slide_text)

    def test_cl02_builds_three_cards_from_stages(self):
        """CL-02 pads to three cards using stages when needed."""
        raw_spec = {
            "title": "Strategic Priorities",
            "subtitle": "Overview",
            "stages": [
                {"title": "Stage 1", "activities": ["A"]},
                {"title": "Stage 2", "activities": ["B"]},
            ],
        }
        cards = slide_service._build_executive_card_content(raw_spec, "CL-02")

        self.assertEqual(len(cards), 2)
        self.assertEqual(cards[0].title, "Stage 1")
        self.assertEqual(cards[1].title, "Stage 2")

    def test_cl03_enters_visual_pipeline_and_renders_three_kpis(self):
        """CL-03 routes through the Layout Engine and renders three KPI cards."""
        spec = self._cl03_slide_spec(source="kpis")
        result = self._completed_result_with_specs([spec])

        with patch("backend.orchestrator.run_pipeline", return_value=result), \
             patch("backend.services.slide_service.plan_visual_pattern", return_value=self._visual_selection("CL-03")) as mock_plan:
            path = slide_service.generate_slide_v2("Title", "Content")

        self.assertEqual(path, "generated_slide_v2.pptx")
        self._assert_slide_count(1)
        mock_plan.assert_called_once()

        prs = Presentation("generated_slide_v2.pptx")
        slide_text = ""
        for shape in prs.slides[0].shapes:
            if hasattr(shape, "text_frame"):
                slide_text += shape.text_frame.text + " "
        for label in ("Cost Reduction", "Cycle Time", "Compliance"):
            self.assertIn(label, slide_text)
        for value in ("15%", "-30%", "98%"):
            self.assertIn(value, slide_text)

    def test_cl03_kpi_fallback_from_metrics(self):
        """CL-03 falls back to metrics when kpis are absent."""
        spec = self._cl03_slide_spec(source="metrics")
        result = self._completed_result_with_specs([spec])

        with patch("backend.orchestrator.run_pipeline", return_value=result), \
             patch("backend.services.slide_service.plan_visual_pattern", return_value=self._visual_selection("CL-03")):
            path = slide_service.generate_slide_v2("Title", "Content")

        self._assert_slide_count(1)

    def test_cl03_kpis_return_empty_when_sources_missing(self):
        """CL-03 returns an empty KPI list when no sources exist (no padding)."""
        raw_spec = {"title": "Key KPIs"}
        kpis = slide_service._build_kpi_content(raw_spec)

        self.assertEqual(len(kpis), 0)

    def test_cl04_enters_visual_pipeline_and_renders_comparison(self):
        """CL-04 routes through the Layout Engine and renders a comparison layout."""
        spec = self._cl04_slide_spec(source="columns")
        result = self._completed_result_with_specs([spec])

        with patch("backend.orchestrator.run_pipeline", return_value=result), \
             patch("backend.services.slide_service.plan_visual_pattern", return_value=self._visual_selection("CL-04")) as mock_plan:
            path = slide_service.generate_slide_v2("Title", "Content")

        self.assertEqual(path, "generated_slide_v2.pptx")
        self._assert_slide_count(1)
        mock_plan.assert_called_once()

        prs = Presentation("generated_slide_v2.pptx")
        slide_text = ""
        for shape in prs.slides[0].shapes:
            if hasattr(shape, "text_frame"):
                slide_text += shape.text_frame.text + " "
        self.assertIn("Current", slide_text)
        self.assertIn("Future", slide_text)
        self.assertIn("Manual", slide_text)
        self.assertIn("Automated", slide_text)

    def test_cl04_derives_columns_from_key_pairs(self):
        """CL-04 derives left/right columns from common key pairs."""
        raw_spec = {
            "title": "Before vs After",
            "before": [
                {"name": "Manual", "text": "Manual processes"},
                {"name": "Slow", "text": "Slow cycle times"},
            ],
            "after": [
                {"name": "Automated", "text": "Automated workflows"},
                {"name": "Fast", "text": "Fast cycle times"},
            ],
        }
        content = slide_service._build_two_column_content(raw_spec)

        self.assertEqual(content["columns"], [{"label": "Before"}, {"label": "After"}])
        self.assertEqual(content["left_items"][0]["name"], "Manual")
        self.assertEqual(content["right_items"][0]["name"], "Automated")

    def test_cl05_enters_visual_pipeline_and_renders_two_columns(self):
        """CL-05 routes through the Layout Engine and renders a two-column text layout."""
        spec = self._cl05_slide_spec(source="columns")
        result = self._completed_result_with_specs([spec])

        with patch("backend.orchestrator.run_pipeline", return_value=result), \
             patch("backend.services.slide_service.plan_visual_pattern", return_value=self._visual_selection("CL-05")) as mock_plan:
            path = slide_service.generate_slide_v2("Title", "Content")

        self.assertEqual(path, "generated_slide_v2.pptx")
        self._assert_slide_count(1)
        mock_plan.assert_called_once()

        prs = Presentation("generated_slide_v2.pptx")
        slide_text = ""
        for shape in prs.slides[0].shapes:
            if hasattr(shape, "text_frame"):
                slide_text += shape.text_frame.text + " "
        self.assertIn("Challenges", slide_text)
        self.assertIn("Recommendations", slide_text)
        self.assertIn("Fragmented processes", slide_text)
        self.assertIn("Unified platform", slide_text)

    def test_cl05_derives_columns_from_challenges_solutions(self):
        """CL-05 derives left/right columns from challenges/solutions key pairs."""
        raw_spec = {
            "title": "Challenges and Recommendations",
            "challenges": ["Fragmented processes", "Limited visibility"],
            "solutions": ["Unified platform", "Real-time dashboards"],
        }
        content = slide_service._build_two_column_content(raw_spec)

        self.assertEqual(content["columns"], [{"label": "Challenges"}, {"label": "Solutions"}])
        self.assertEqual(content["left_items"], ["Fragmented processes", "Limited visibility"])
        self.assertEqual(content["right_items"], ["Unified platform", "Real-time dashboards"])

    def test_build_pattern_content_does_not_mutate_raw_spec(self):
        """The pattern content builder must not mutate the input SlideSpec dict."""
        raw_spec = {"title": "Benefits", "subtitle": "Overview"}
        slide_service._build_pattern_content(raw_spec, "CL-01")

        self.assertNotIn("cards", raw_spec)
        self.assertNotIn("kpis", raw_spec)
        self.assertNotIn("left_items", raw_spec)

    def test_mixed_all_creative_listing_patterns(self):
        """A deck with every Creative Listing pattern renders via the visual pipeline."""
        specs = [
            self._cl01_slide_spec(source="cards"),
            self._cl02_slide_spec(source="cards"),
            self._cl03_slide_spec(source="kpis"),
            self._cl04_slide_spec(source="columns"),
            self._cl05_slide_spec(source="columns"),
            self._slide_spec("Executive Summary"),
        ]
        result = self._completed_result_with_specs(specs)

        with patch("backend.orchestrator.run_pipeline", return_value=result), \
             patch("backend.services.slide_service.plan_visual_pattern") as mock_plan:
            mock_plan.side_effect = [
                self._visual_selection("CL-01"),
                self._visual_selection("CL-02"),
                self._visual_selection("CL-03"),
                self._visual_selection("CL-04"),
                self._visual_selection("CL-05"),
                self._visual_selection("CL-06"),
            ]
            path = slide_service.generate_slide_v2("Title", "Content")

        self.assertEqual(path, "generated_slide_v2.pptx")
        self._assert_slide_count(6)
        self.assertEqual(mock_plan.call_count, 6)

    def test_legacy_renderer_unchanged_for_non_whitelisted_patterns(self):
        """Non-whitelisted patterns (e.g. IG-07) still use the legacy renderer."""
        spec = self._slide_spec("Process Overview")
        result = self._completed_result_with_specs([spec])

        with patch("backend.orchestrator.run_pipeline", return_value=result), \
             patch("backend.services.slide_service.plan_visual_pattern", return_value=self._visual_selection("IG-07")) as mock_plan, \
             patch("backend.services.slide_service.generate_layout") as mock_layout:
            path = slide_service.generate_slide_v2("Title", "Content")

        self.assertEqual(path, "generated_slide_v2.pptx")
        self._assert_slide_count(1)
        mock_plan.assert_called_once()
        mock_layout.assert_not_called()

    def test_visual_pipeline_failure_falls_back_for_cl04(self):
        """Any exception in the CL-04 visual pipeline falls back to legacy renderer."""
        spec = self._cl04_slide_spec(source="columns")
        result = self._completed_result_with_specs([spec])

        with patch("backend.orchestrator.run_pipeline", return_value=result), \
             patch("backend.services.slide_service.plan_visual_pattern", return_value=self._visual_selection("CL-04")), \
             patch("backend.services.slide_service.generate_layout", side_effect=RuntimeError("layout failure")):
            path = slide_service.generate_slide_v2("Title", "Content")

        self.assertEqual(path, "generated_slide_v2.pptx")
        self._assert_slide_count(1)

    def test_ig01_enters_visual_pipeline_and_renders_timeline(self):
        """IG-01 routes through the Layout Engine and renders timeline nodes."""
        spec = self._ig01_slide_spec()
        result = self._completed_result_with_specs([spec])

        with patch("backend.orchestrator.run_pipeline", return_value=result), \
             patch("backend.services.slide_service.plan_visual_pattern", return_value=self._visual_selection("IG-01")) as mock_plan:
            path = slide_service.generate_slide_v2("Title", "Content")

        self.assertEqual(path, "generated_slide_v2.pptx")
        self._assert_slide_count(1)
        mock_plan.assert_called_once()

        prs = Presentation("generated_slide_v2.pptx")
        slide_text = ""
        for shape in prs.slides[0].shapes:
            if hasattr(shape, "text_frame"):
                slide_text += shape.text_frame.text + " "
        for event in ("Assess", "Design", "Build", "Scale"):
            self.assertIn(event, slide_text)

    def test_ig02_enters_visual_pipeline_and_renders_roadmap(self):
        """IG-02 routes through the Layout Engine and renders roadmap phases."""
        spec = self._ig02_slide_spec()
        result = self._completed_result_with_specs([spec])

        with patch("backend.orchestrator.run_pipeline", return_value=result), \
             patch("backend.services.slide_service.plan_visual_pattern", return_value=self._visual_selection("IG-02")) as mock_plan:
            path = slide_service.generate_slide_v2("Title", "Content")

        self.assertEqual(path, "generated_slide_v2.pptx")
        self._assert_slide_count(1)
        mock_plan.assert_called_once()

        prs = Presentation("generated_slide_v2.pptx")
        slide_text = ""
        for shape in prs.slides[0].shapes:
            if hasattr(shape, "text_frame"):
                slide_text += shape.text_frame.text + " "
        for phase in ("Phase 1", "Phase 2", "Phase 3", "Phase 4"):
            self.assertIn(phase, slide_text)

    def test_ig03_enters_visual_pipeline_and_renders_process_flow(self):
        """IG-03 routes through the Layout Engine and renders process steps."""
        spec = self._ig03_slide_spec()
        result = self._completed_result_with_specs([spec])

        with patch("backend.orchestrator.run_pipeline", return_value=result), \
             patch("backend.services.slide_service.plan_visual_pattern", return_value=self._visual_selection("IG-03")) as mock_plan:
            path = slide_service.generate_slide_v2("Title", "Content")

        self.assertEqual(path, "generated_slide_v2.pptx")
        self._assert_slide_count(1)
        mock_plan.assert_called_once()

        prs = Presentation("generated_slide_v2.pptx")
        slide_text = ""
        for shape in prs.slides[0].shapes:
            if hasattr(shape, "text_frame"):
                slide_text += shape.text_frame.text + " "
        for step in ("Request", "Approve", "Purchase", "Receive", "Pay"):
            self.assertIn(step, slide_text)

    def test_ig04_enters_visual_pipeline_and_renders_matrix(self):
        """IG-04 routes through the Layout Engine and renders matrix cells."""
        spec = self._ig04_slide_spec()
        result = self._completed_result_with_specs([spec])

        with patch("backend.orchestrator.run_pipeline", return_value=result), \
             patch("backend.services.slide_service.plan_visual_pattern", return_value=self._visual_selection("IG-04")) as mock_plan:
            path = slide_service.generate_slide_v2("Title", "Content")

        self.assertEqual(path, "generated_slide_v2.pptx")
        self._assert_slide_count(1)
        mock_plan.assert_called_once()

        prs = Presentation("generated_slide_v2.pptx")
        slide_text = ""
        for shape in prs.slides[0].shapes:
            if hasattr(shape, "text_frame"):
                slide_text += shape.text_frame.text + " "
        for cell in ("Quick wins", "Major projects", "Strategic bets", "Fill-ins"):
            self.assertIn(cell, slide_text)

    def test_ig05_enters_visual_pipeline_and_renders_journey(self):
        """IG-05 routes through the Layout Engine and renders journey stages."""
        spec = self._ig05_slide_spec()
        result = self._completed_result_with_specs([spec])

        with patch("backend.orchestrator.run_pipeline", return_value=result), \
             patch("backend.services.slide_service.plan_visual_pattern", return_value=self._visual_selection("IG-05")) as mock_plan:
            path = slide_service.generate_slide_v2("Title", "Content")

        self.assertEqual(path, "generated_slide_v2.pptx")
        self._assert_slide_count(1)
        mock_plan.assert_called_once()

        prs = Presentation("generated_slide_v2.pptx")
        slide_text = ""
        for shape in prs.slides[0].shapes:
            if hasattr(shape, "text_frame"):
                slide_text += shape.text_frame.text + " "
        for stage in ("Aware", "Engage", "Adopt", "Optimize"):
            self.assertIn(stage, slide_text)

    def test_ig06_enters_visual_pipeline_and_renders_capability_map(self):
        """IG-06 routes through the Layout Engine and renders capability domains."""
        spec = self._ig06_slide_spec()
        result = self._completed_result_with_specs([spec])

        with patch("backend.orchestrator.run_pipeline", return_value=result), \
             patch("backend.services.slide_service.plan_visual_pattern", return_value=self._visual_selection("IG-06")) as mock_plan:
            path = slide_service.generate_slide_v2("Title", "Content")

        self.assertEqual(path, "generated_slide_v2.pptx")
        self._assert_slide_count(1)
        mock_plan.assert_called_once()

        prs = Presentation("generated_slide_v2.pptx")
        slide_text = ""
        for shape in prs.slides[0].shapes:
            if hasattr(shape, "text_frame"):
                slide_text += shape.text_frame.text + " "
        for domain in ("Strategy", "Sourcing", "Operations", "Analytics"):
            self.assertIn(domain, slide_text)
        for capability in ("Demand planning", "RFx", "Purchase orders", "Spend visibility"):
            self.assertIn(capability, slide_text)

    def test_ig04_matrix_cells_flattened_from_rows(self):
        """IG-04 flattens nested row cells into the cell collection (no padding)."""
        raw_spec = {
            "title": "Matrix",
            "rows": [
                {"cells": [{"value": "A"}, {"value": "B"}]},
                {"cells": [{"value": "C"}, {"value": "D"}]},
            ],
        }
        content = slide_service._build_infographic_content(raw_spec, "IG-04")

        self.assertEqual(len(content["cells"]), 4)
        self.assertEqual(content["cells"][0]["value"], "A")
        self.assertEqual(content["cells"][1]["value"], "B")
        self.assertEqual(content["cells"][2]["value"], "C")

    def test_ig06_domains_fallback_from_capabilities_list(self):
        """IG-06 derives domains from a flat capabilities list if domains are absent (no padding)."""
        raw_spec = {
            "title": "Capability Map",
            "capabilities": [
                {"name": "Strategy", "capabilities": [{"name": "Planning"}]},
                {"name": "Sourcing", "capabilities": [{"name": "RFx"}]},
            ],
        }
        content = slide_service._build_infographic_content(raw_spec, "IG-06")

        self.assertEqual(len(content["domains"]), 2)
        self.assertEqual(content["domains"][0]["name"], "Strategy")
        self.assertEqual(content["domains"][0]["capabilities"][0]["name"], "Planning")

    def test_infographic_placeholder_no_padding(self):
        """Infographic builders do not pad missing items."""
        raw_spec = {"title": "Timeline"}
        content = slide_service._build_infographic_content(raw_spec, "IG-01")

        self.assertEqual(len(content["events"]), 0)

    def test_mixed_infographic_deck(self):
        """A deck with every infographic pattern renders via the visual pipeline."""
        specs = [
            self._ig01_slide_spec(),
            self._ig02_slide_spec(),
            self._ig03_slide_spec(),
            self._ig04_slide_spec(),
            self._ig05_slide_spec(),
            self._ig06_slide_spec(),
        ]
        result = self._completed_result_with_specs(specs)

        with patch("backend.orchestrator.run_pipeline", return_value=result), \
             patch("backend.services.slide_service.plan_visual_pattern") as mock_plan:
            mock_plan.side_effect = [
                self._visual_selection("IG-01"),
                self._visual_selection("IG-02"),
                self._visual_selection("IG-03"),
                self._visual_selection("IG-04"),
                self._visual_selection("IG-05"),
                self._visual_selection("IG-06"),
            ]
            path = slide_service.generate_slide_v2("Title", "Content")

        self.assertEqual(path, "generated_slide_v2.pptx")
        self._assert_slide_count(6)
        self.assertEqual(mock_plan.call_count, 6)

    def test_mixed_creative_and_infographic_deck(self):
        """A deck mixing Creative Listing and Infographic patterns renders correctly."""
        specs = [
            self._cl01_slide_spec(source="cards"),
            self._ig03_slide_spec(),
            self._slide_spec("Executive Summary"),
            self._ig06_slide_spec(),
        ]
        result = self._completed_result_with_specs(specs)

        with patch("backend.orchestrator.run_pipeline", return_value=result), \
             patch("backend.services.slide_service.plan_visual_pattern") as mock_plan:
            mock_plan.side_effect = [
                self._visual_selection("CL-01"),
                self._visual_selection("IG-03"),
                self._visual_selection("CL-06"),
                self._visual_selection("IG-06"),
            ]
            path = slide_service.generate_slide_v2("Title", "Content")

        self.assertEqual(path, "generated_slide_v2.pptx")
        self._assert_slide_count(4)
        self.assertEqual(mock_plan.call_count, 4)

    def test_visual_pipeline_failure_falls_back_for_ig03(self):
        """Any exception in the IG-03 visual pipeline falls back to legacy renderer."""
        spec = self._ig03_slide_spec()
        result = self._completed_result_with_specs([spec])

        with patch("backend.orchestrator.run_pipeline", return_value=result), \
             patch("backend.services.slide_service.plan_visual_pattern", return_value=self._visual_selection("IG-03")), \
             patch("backend.services.slide_service.generate_layout", side_effect=RuntimeError("layout failure")):
            path = slide_service.generate_slide_v2("Title", "Content")

        self.assertEqual(path, "generated_slide_v2.pptx")
        self._assert_slide_count(1)


if __name__ == "__main__":
    unittest.main()
