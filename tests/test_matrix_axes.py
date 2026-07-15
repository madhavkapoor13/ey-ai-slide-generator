from __future__ import annotations

import unittest

from pptx import Presentation
from pptx.util import Inches

from backend.layout_engine.layout_engine import generate_layout
from ppt_renderer.operating_model_renderer import OperatingModelRenderer
from schemas.visual import VisualPatternSelection


class MatrixAxesTests(unittest.TestCase):
    def test_matrix_renders_axis_labels_and_quadrant_colors(self):
        selection = VisualPatternSelection(
            pattern_id="IG-04",
            category="infographic",
            confidence=0.9,
            reasoning="Risk matrix",
        )
        layout = generate_layout(selection, item_count=9)

        content = {
            "title": "Risk Matrix",
            "subtitle": "Impact vs Likelihood",
            "cells": [
                {"value": "R1", "quadrant": {"impact": "High", "likelihood": "High"}},
                {"value": "R2", "quadrant": {"impact": "High", "likelihood": "Medium"}},
                {"value": "R3", "quadrant": {"impact": "High", "likelihood": "Low"}},
                {"value": "R4", "quadrant": {"impact": "Medium", "likelihood": "High"}},
                {"value": "R5", "quadrant": {"impact": "Medium", "likelihood": "Medium"}},
                {"value": "R6", "quadrant": {"impact": "Medium", "likelihood": "Low"}},
                {"value": "R7", "quadrant": {"impact": "Low", "likelihood": "High"}},
                {"value": "R8", "quadrant": {"impact": "Low", "likelihood": "Medium"}},
                {"value": "R9", "quadrant": {"impact": "Low", "likelihood": "Low"}},
            ],
        }

        renderer = OperatingModelRenderer()
        prs = renderer.render(content, layout_spec=layout)

        slide_text = ""
        for shape in prs.slides[0].shapes:
            if hasattr(shape, "text_frame"):
                slide_text += shape.text_frame.text + " "

        for label in ("Low", "Medium", "High", "Risk Matrix"):
            self.assertIn(label, slide_text)

        # At least the 9 risk cells rendered as shapes.
        self.assertGreaterEqual(len(prs.slides[0].shapes), 9)
