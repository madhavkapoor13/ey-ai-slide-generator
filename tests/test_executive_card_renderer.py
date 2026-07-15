import unittest

from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE_TYPE
from pptx.util import Inches

from backend.design_system.theme_loader import clear_cache, get_current_theme, set_active_theme
from ppt_renderer.components.component_dispatcher import render as dispatch_render
from ppt_renderer.components.executive_card_renderer import render as render_executive_card
from schemas.executive_card import ExecutiveCardContent
from schemas.layout import ComponentSpecification


class ExecutiveCardRendererTests(unittest.TestCase):

    def setUp(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

    def tearDown(self):
        clear_cache()
        set_active_theme("ey_blue")

    def _make_spec(self, component_id: str, placeholder: str, x: float, y: float, w: float, h: float):
        return ComponentSpecification(
            component_id=component_id,
            type="executive_card",
            placeholder=placeholder,
            x=x,
            y=y,
            width=w,
            height=h,
        )

    def _content(self, *cards: ExecutiveCardContent):
        return {"cards": [card.model_dump(mode="json") for card in cards]}

    def test_single_card(self):
        content = self._content(ExecutiveCardContent(title="Insight 1", description="A key insight."))
        spec = self._make_spec("ec1", "card_1", 0.1, 0.2, 0.3, 0.4)

        render_executive_card(spec, self.prs, self.slide, content)

        texts = [s.text_frame.text for s in self.slide.shapes if s.has_text_frame]
        self.assertIn("Insight 1", texts)
        self.assertIn("A key insight.", texts)

    def test_two_cards_equal_width(self):
        content = self._content(
            ExecutiveCardContent(title="Left", description="Left card."),
            ExecutiveCardContent(title="Right", description="Right card."),
        )
        left_spec = self._make_spec("ec_left", "card_1", 0.05, 0.2, 0.42, 0.5)
        right_spec = self._make_spec("ec_right", "card_2", 0.53, 0.2, 0.42, 0.5)

        render_executive_card(left_spec, self.prs, self.slide, content)
        render_executive_card(right_spec, self.prs, self.slide, content)

        texts = [s.text_frame.text for s in self.slide.shapes if s.has_text_frame]
        self.assertIn("Left", texts)
        self.assertIn("Right", texts)

        # Equal-width card backgrounds (auto shapes).
        card_shapes = [s for s in self.slide.shapes if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
        self.assertEqual(len(card_shapes), 2)
        self.assertAlmostEqual(card_shapes[0].width.inches, card_shapes[1].width.inches, places=5)

    def test_three_cards(self):
        content = self._content(
            ExecutiveCardContent(title="A", description="..."),
            ExecutiveCardContent(title="B", description="..."),
            ExecutiveCardContent(title="C", description="..."),
        )
        for i, x in enumerate([0.05, 0.35, 0.65]):
            spec = self._make_spec(f"ec_{i}", f"card_{i + 1}", x, 0.2, 0.25, 0.5)
            render_executive_card(spec, self.prs, self.slide, content)

        texts = [s.text_frame.text for s in self.slide.shapes if s.text_frame.text]
        for letter in ("A", "B", "C"):
            self.assertIn(letter, texts)

    def test_four_cards(self):
        content = self._content(
            ExecutiveCardContent(title="A", description="..."),
            ExecutiveCardContent(title="B", description="..."),
            ExecutiveCardContent(title="C", description="..."),
            ExecutiveCardContent(title="D", description="..."),
        )
        for i, x in enumerate([0.05, 0.27, 0.49, 0.71]):
            spec = self._make_spec(f"ec_{i}", f"card_{i + 1}", x, 0.2, 0.20, 0.5)
            render_executive_card(spec, self.prs, self.slide, content)

        texts = [s.text_frame.text for s in self.slide.shapes if s.text_frame.text]
        for letter in ("A", "B", "C", "D"):
            self.assertIn(letter, texts)

    def test_long_description(self):
        long_text = "This is a deliberately long description that should wrap inside the card bounds rather than overflow."
        content = self._content(ExecutiveCardContent(title="Long", description=long_text))
        spec = self._make_spec("ec_long", "card_1", 0.1, 0.2, 0.3, 0.4)

        render_executive_card(spec, self.prs, self.slide, content)

        texts = [s.text_frame.text for s in self.slide.shapes if s.has_text_frame]
        self.assertTrue(any(long_text in text for text in texts))

    def test_metric_present(self):
        content = self._content(ExecutiveCardContent(title="Metric", description="...", metric="37%"))
        spec = self._make_spec("ec_metric", "card_1", 0.1, 0.2, 0.3, 0.4)

        render_executive_card(spec, self.prs, self.slide, content)

        texts = [s.text_frame.text for s in self.slide.shapes if s.has_text_frame]
        self.assertIn("37%", texts)

    def test_metric_absent(self):
        content = self._content(ExecutiveCardContent(title="No Metric", description="No metric here."))
        spec = self._make_spec("ec_no_metric", "card_1", 0.1, 0.2, 0.3, 0.4)

        render_executive_card(spec, self.prs, self.slide, content)

        texts = [s.text_frame.text for s in self.slide.shapes if s.has_text_frame]
        self.assertNotIn("37%", texts)

    def test_highlight_present(self):
        content = self._content(ExecutiveCardContent(title="Highlighted", description="...", highlight="Priority"))
        spec = self._make_spec("ec_highlight", "card_1", 0.1, 0.2, 0.3, 0.4)

        render_executive_card(spec, self.prs, self.slide, content)

        texts = [s.text_frame.text for s in self.slide.shapes if s.has_text_frame]
        self.assertIn("Priority", texts)

    def test_theme_styling_applied(self):
        set_active_theme("ey_parthenon")
        theme = get_current_theme()
        content = self._content(ExecutiveCardContent(title="Styled", description="...", highlight="Important"))
        spec = self._make_spec("ec_style", "card_1", 0.1, 0.2, 0.3, 0.4)

        render_executive_card(spec, self.prs, self.slide, content)

        # The highlight badge should use the active theme's primary color.
        badge = None
        for shape in self.slide.shapes:
            if shape.has_text_frame and shape.text_frame.text == "Important":
                badge = shape
                break
        self.assertIsNotNone(badge)
        self.assertEqual(badge.fill.fore_color.rgb, theme.color("primary"))

    def test_placeholder_mapping(self):
        content = self._content(
            ExecutiveCardContent(title="First", description="..."),
            ExecutiveCardContent(title="Second", description="..."),
        )
        spec = self._make_spec("ec_second", "card_2", 0.1, 0.2, 0.3, 0.4)

        render_executive_card(spec, self.prs, self.slide, content)

        texts = [s.text_frame.text for s in self.slide.shapes if s.has_text_frame]
        self.assertIn("Second", texts)
        self.assertNotIn("First", texts)

    def test_dispatcher_routes_executive_card(self):
        content = self._content(ExecutiveCardContent(title="Dispatched", description="..."))
        spec = self._make_spec("ec_dispatched", "card_1", 0.1, 0.2, 0.3, 0.4)

        dispatch_render(spec, self.prs, self.slide, content)

        texts = [s.text_frame.text for s in self.slide.shapes if s.has_text_frame]
        self.assertIn("Dispatched", texts)


if __name__ == "__main__":
    unittest.main()
