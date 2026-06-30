from __future__ import annotations

from datetime import date
from typing import Any

from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches

from ppt_renderer.operating_model_theme import OperatingModelTheme as Theme


class TextStyle:
    @staticmethod
    def apply(paragraph, size, color, bold: bool = False) -> None:
        paragraph.font.name = Theme.FONT
        paragraph.font.size = size
        paragraph.font.color.rgb = color
        paragraph.font.bold = bold


class Geometry:
    @staticmethod
    def textbox(slide, layout: dict[str, Any]):
        return slide.shapes.add_textbox(
            layout["x"],
            layout["y"],
            layout["width"],
            layout["height"],
        )

    @staticmethod
    def shape(slide, shape_type, layout: dict[str, Any]):
        return slide.shapes.add_shape(
            shape_type,
            layout["x"],
            layout["y"],
            layout["width"],
            layout["height"],
        )


class HeaderComponent:
    @staticmethod
    def draw(slide, spec: dict[str, Any], layout: dict[str, Any]) -> None:
        title = Geometry.textbox(slide, layout["title"])
        title_tf = title.text_frame
        title_tf.clear()
        title_tf.word_wrap = True
        title_p = title_tf.paragraphs[0]
        title_p.text = spec.get("title", "")
        title_p.font.name = Theme.TITLE_FONT
        title_p.font.size = Theme.TITLE_SIZE
        title_p.font.color.rgb = Theme.INK

        subtitle = Geometry.textbox(slide, layout["subtitle"])
        subtitle_p = subtitle.text_frame.paragraphs[0]
        subtitle_p.text = spec.get("subtitle", "")
        TextStyle.apply(subtitle_p, Theme.SUBTITLE_SIZE, Theme.MUTED_TEXT, bold=True)

        description = Geometry.textbox(slide, layout["description"])
        description_tf = description.text_frame
        description_tf.word_wrap = True
        description_p = description_tf.paragraphs[0]
        description_p.text = spec.get("description", "")
        TextStyle.apply(description_p, Theme.DESCRIPTION_SIZE, Theme.MUTED_TEXT)

        divider = Geometry.shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, layout["divider"])
        divider.fill.solid()
        divider.fill.fore_color.rgb = Theme.ACCENT
        divider.line.fill.background()


class SummaryRibbonComponent:
    @staticmethod
    def draw(slide, summary: dict[str, Any], layout: dict[str, Any]) -> None:
        ribbon = Geometry.shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, layout["ribbon"])
        ribbon.fill.solid()
        ribbon.fill.fore_color.rgb = Theme.ACCENT
        ribbon.line.color.rgb = Theme.ACCENT_DARK
        ribbon.line.width = Inches(0.006)

        text_box = Geometry.textbox(slide, layout["summary_text"])
        tf = text_box.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.margin_left = 0
        tf.margin_right = 0

        headline = tf.paragraphs[0]
        headline.text = summary.get("headline", "")
        TextStyle.apply(headline, Theme.RIBBON_HEADLINE_SIZE, Theme.LIGHT_TEXT, bold=True)

        body = tf.add_paragraph()
        body.text = summary.get("description", "")
        TextStyle.apply(body, Theme.RIBBON_BODY_SIZE, Theme.LIGHT_TEXT)

        for metric_layout in layout["metrics"]:
            MetricCardComponent.draw(slide, metric_layout["metric"], metric_layout)


class MetricCardComponent:
    @staticmethod
    def draw(slide, metric: dict[str, str], layout: dict[str, Any]) -> None:
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            layout["x"],
            layout["y"],
            layout["width"],
            layout["height"],
        )
        card.fill.solid()
        card.fill.fore_color.rgb = Theme.ACCENT_LIGHT
        card.line.fill.background()

        tf = card.text_frame
        tf.clear()
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.06)
        tf.margin_right = Inches(0.06)
        tf.word_wrap = True

        value = tf.paragraphs[0]
        value.text = metric.get("value", "")
        value.alignment = PP_ALIGN.CENTER
        TextStyle.apply(value, Theme.METRIC_VALUE_SIZE, Theme.ACCENT_DARK, bold=True)

        label = tf.add_paragraph()
        label.text = metric.get("label", "")
        label.alignment = PP_ALIGN.CENTER
        TextStyle.apply(label, Theme.METRIC_LABEL_SIZE, Theme.MUTED_TEXT)


class StageComponent:
    @staticmethod
    def draw(slide, layout: dict[str, Any]) -> None:
        stage = layout["stage"]

        container = Geometry.shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, layout["container"])
        container.fill.solid()
        container.fill.fore_color.rgb = Theme.PANEL
        container.line.color.rgb = Theme.BORDER
        container.line.width = Inches(0.005)

        header = Geometry.shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, layout["header"])
        header.fill.solid()
        header.fill.fore_color.rgb = Theme.ACCENT
        header.line.color.rgb = Theme.ACCENT

        number = Geometry.shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, layout["number"])
        number.fill.solid()
        number.fill.fore_color.rgb = Theme.BACKGROUND
        number.line.fill.background()
        number_tf = number.text_frame
        number_tf.clear()
        number_tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        number_p = number_tf.paragraphs[0]
        number_p.text = str(stage.get("number", ""))
        number_p.alignment = PP_ALIGN.CENTER
        TextStyle.apply(number_p, Theme.STAGE_NUMBER_SIZE, Theme.ACCENT_DARK, bold=True)

        title = Geometry.textbox(slide, layout["title"])
        title_tf = title.text_frame
        title_tf.clear()
        title_tf.word_wrap = True
        title_tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        title_p = title_tf.paragraphs[0]
        title_p.text = str(stage.get("title", "")).upper()
        TextStyle.apply(title_p, Theme.STAGE_TITLE_SIZE, Theme.LIGHT_TEXT, bold=True)

        for activity_layout in layout["activities"]:
            ActivityComponent.draw(slide, activity_layout)


class ActivityComponent:
    @staticmethod
    def draw(slide, layout: dict[str, Any]) -> None:
        activity = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            layout["x"],
            layout["y"],
            layout["width"],
            layout["height"],
        )
        activity.fill.solid()
        activity.fill.fore_color.rgb = Theme.CARD
        activity.line.color.rgb = Theme.BORDER
        activity.line.width = Inches(0.004)

        tf = activity.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.06)
        tf.margin_right = Inches(0.06)
        tf.margin_top = Inches(0.03)
        tf.margin_bottom = Inches(0.03)

        title = tf.paragraphs[0]
        title.text = layout["text"]
        TextStyle.apply(title, Theme.ACTIVITY_TITLE_SIZE, Theme.ACCENT_DARK, bold=True)


class ConnectorComponent:
    @staticmethod
    def draw(slide, layout: dict[str, Any]):
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            layout["start_x"],
            layout["start_y"],
            layout["end_x"],
            layout["end_y"],
        )
        connector.line.color.rgb = Theme.ACCENT_DARK
        connector.line.width = Theme.CONNECTOR_WIDTH
        ConnectorComponent._add_arrowhead(connector)
        return connector

    @staticmethod
    def _add_arrowhead(connector) -> None:
        line = connector.element.spPr.ln
        head = OxmlElement("a:headEnd")
        head.set("type", "triangle")
        head.set("w", "sm")
        head.set("len", "sm")
        line.append(head)


class RiskStripComponent:
    @staticmethod
    def draw(slide, layout: dict[str, Any]) -> None:
        strip = Geometry.shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, layout["strip"])
        strip.fill.solid()
        strip.fill.fore_color.rgb = Theme.RISK_FILL
        strip.line.color.rgb = Theme.RISK_BORDER
        strip.line.width = Inches(0.005)

        label = Geometry.textbox(slide, layout["label"])
        label_tf = label.text_frame
        label_tf.clear()
        label_tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        label_p = label_tf.paragraphs[0]
        label_p.text = "RISKS / PAIN POINTS"
        TextStyle.apply(label_p, Theme.RISK_SIZE, Theme.RISK_TEXT, bold=True)

        for cell in layout["cells"]:
            RiskStripComponent._draw_cell(slide, cell)

    @staticmethod
    def _draw_cell(slide, layout: dict[str, Any]) -> None:
        if not layout["texts"]:
            return

        text = " | ".join(layout["texts"])
        box = slide.shapes.add_textbox(
            layout["x"],
            layout["y"],
            layout["width"],
            layout["height"],
        )
        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.05)
        tf.margin_right = Inches(0.05)

        paragraph = tf.paragraphs[0]
        paragraph.text = text
        paragraph.alignment = PP_ALIGN.CENTER
        TextStyle.apply(paragraph, Theme.RISK_SIZE, Theme.ACCENT_DARK, bold=True)


class FooterComponent:
    @staticmethod
    def draw(slide, layout: dict[str, Any], slide_number: int = 1) -> None:
        rule = Geometry.shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, layout["rule"])
        rule.fill.solid()
        rule.fill.fore_color.rgb = Theme.BORDER
        rule.line.fill.background()

        FooterComponent._draw_text(slide, layout["left"], "Generated by EY AI Pitch Prototype", PP_ALIGN.LEFT)
        FooterComponent._draw_text(slide, layout["center"], date.today().strftime("%d %b %Y"), PP_ALIGN.CENTER)
        FooterComponent._draw_text(slide, layout["right"], f"Slide {slide_number}", PP_ALIGN.RIGHT)

    @staticmethod
    def _draw_text(slide, layout: dict[str, Any], text: str, alignment) -> None:
        box = Geometry.textbox(slide, layout)
        paragraph = box.text_frame.paragraphs[0]
        paragraph.text = text
        paragraph.alignment = alignment
        TextStyle.apply(paragraph, Theme.FOOTER_SIZE, Theme.MUTED_TEXT)
