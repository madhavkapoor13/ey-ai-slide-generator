from __future__ import annotations

from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches, Pt

from backend.design_system.theme_loader import get_current_theme


class OperatingModelRenderer:
    """Generates editable consulting-style operating model PowerPoint slides."""

    def __init__(self):
        self.prs = Presentation()
        theme = get_current_theme()
        self.prs.slide_width = theme.SLIDE_WIDTH if hasattr(theme, "SLIDE_WIDTH") else Inches(13.333)
        self.prs.slide_height = theme.SLIDE_HEIGHT if hasattr(theme, "SLIDE_HEIGHT") else Inches(7.5)

    def render(
        self,
        slide_spec: dict[str, Any],
        output_path: str = "operating_model.pptx",
        presentation: Presentation | None = None,
        layout_spec=None,
    ) -> Presentation:
        """
        Render a single operating-model slide.

        If ``presentation`` is supplied, the slide is appended to that
        presentation and no save occurs. This lets callers build multi-slide
        decks from multiple SlideSpecs while reusing one ``Presentation``
        object. When ``presentation`` is omitted, behaviour is identical to
        the original single-slide contract.

        If ``layout_spec`` is supplied, the new Visual Layout Engine path is
        used and the legacy hardcoded drawing is skipped.
        """
        prs = presentation if presentation is not None else self.prs
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        if layout_spec is not None:
            self._render_layout(slide_spec, layout_spec, prs, slide)
        else:
            self._render_legacy(slide_spec, slide, prs)

        if presentation is None:
            prs.save(output_path)

        return prs

    def _render_legacy(self, slide_spec: dict[str, Any], slide, prs) -> None:
        """Clean fallback: title, subtitle, and up to six bullet lines."""
        self._draw_background(slide)
        _render_text_fallback(slide, slide_spec)

    def _render_layout(
        self,
        slide_spec: dict[str, Any],
        layout_spec,
        presentation: Presentation,
        slide,
    ) -> None:
        """Visual Layout Engine path: delegate each component to the dispatcher."""
        from ppt_renderer.components import component_dispatcher
        from ppt_renderer.components.header_renderer import render_header
        from ppt_renderer.components.footer_renderer import render_footer

        self._draw_background(slide)

        pattern_id = getattr(layout_spec, "visual_pattern", None)
        render_header(layout_spec.header, presentation, slide, slide_spec, pattern_id=pattern_id)

        layout_context = {"pattern_id": pattern_id}
        for component in layout_spec.components:
            component_dispatcher.render(
                component, presentation, slide, slide_spec, layout_context=layout_context
            )

        render_footer(
            layout_spec.footer,
            presentation,
            slide,
            slide_spec,
            len(presentation.slides),
            pattern_id=pattern_id,
        )

    @staticmethod
    def _draw_background(slide) -> None:
        theme = get_current_theme()
        width = theme.SLIDE_WIDTH if hasattr(theme, "SLIDE_WIDTH") else Inches(13.333)
        height = theme.SLIDE_HEIGHT if hasattr(theme, "SLIDE_HEIGHT") else Inches(7.5)
        background = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            0,
            0,
            width,
            height,
        )
        background.fill.solid()
        background.fill.fore_color.rgb = theme.color("background")
        background.line.fill.background()


def _render_text_fallback(slide, slide_spec: dict[str, Any]) -> None:
    """Draw a clean title + subtitle + bullet fallback when no layout is supplied."""
    theme = get_current_theme()
    title = str(slide_spec.get("title", ""))
    subtitle = str(slide_spec.get("subtitle", "") or slide_spec.get("description", ""))
    bullets = _extract_fallback_bullets(slide_spec)

    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(12.3), Inches(0.9)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = theme.color("ink")

    if subtitle:
        sub_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.6)
        )
        tf = sub_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(14)
        p.font.color.rgb = theme.color("charcoal")

    if bullets:
        body_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(2.4), Inches(11.9), Inches(4.0)
        )
        tf = body_box.text_frame
        tf.word_wrap = True
        for index, bullet in enumerate(bullets):
            para = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
            para.text = f"• {bullet}"
            para.font.size = Pt(14)
            para.font.color.rgb = theme.color("ink")
            if index > 0:
                para.space_before = Pt(8)


def _extract_fallback_bullets(slide_spec: dict[str, Any]) -> list[str]:
    """Return up to six bullet lines from executive_summary or stages."""
    import re

    text = slide_spec.get("executive_summary", "") or slide_spec.get("description", "")
    bullets: list[str] = []
    if isinstance(text, str) and text.strip():
        parts = re.split(r"(?<=[.!?])\s+", text.strip())
        for part in parts:
            part = part.strip()
            if part:
                bullets.append(part)
    stages = slide_spec.get("stages", [])
    if isinstance(stages, list):
        for stage in stages:
            if isinstance(stage, dict):
                label = stage.get("label") or stage.get("title", "")
                if label:
                    bullets.append(str(label))
            elif stage:
                bullets.append(str(stage))
    return bullets[:6]
