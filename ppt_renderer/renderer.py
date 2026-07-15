from pptx import Presentation
from pptx.util import Inches

from backend.design_system.theme_loader import get_current_theme
from ppt_renderer.operating_model_renderer import _render_text_fallback


class ProcessFlowRenderer:

    def __init__(self):

        self.prs = Presentation()

        # Widescreen 16:9
        theme = get_current_theme()
        self.prs.slide_width = theme.SLIDE_WIDTH if hasattr(theme, "SLIDE_WIDTH") else Inches(13.333)
        self.prs.slide_height = theme.SLIDE_HEIGHT if hasattr(theme, "SLIDE_HEIGHT") else Inches(7.5)

    def render(
        self,
        slide_spec,
        output_path="output.pptx",
        presentation=None,
        layout_spec=None,
    ):
        """
        Render a single process-flow slide.

        If ``presentation`` is supplied, the slide is appended to that
        presentation and no save occurs, enabling multi-slide deck
        rendering. When ``presentation`` is omitted, behaviour matches the
        original single-slide contract.

        If ``layout_spec`` is supplied, the new Visual Layout Engine path is
        used and the legacy hardcoded drawing is skipped. This preserves
        full backward compatibility when ``layout_spec`` is omitted.
        """
        prs = presentation if presentation is not None else self.prs

        slide = prs.slides.add_slide(
            prs.slide_layouts[6]
        )

        if layout_spec is not None:
            self._render_layout(slide_spec, layout_spec, prs, slide)
        else:
            self._render_legacy(slide_spec, slide, prs)

        # -------------------------
        # Save
        # -------------------------

        if presentation is None:
            prs.save(output_path)

        return prs

    def _render_legacy(self, slide_spec, slide, prs):
        """Clean fallback: title, subtitle, and up to six bullet lines."""
        _render_text_fallback(slide, slide_spec)

    def _render_layout(self, slide_spec, layout_spec, presentation, slide):
        """Visual Layout Engine path: delegate each component to the dispatcher."""
        from ppt_renderer.components import component_dispatcher
        from ppt_renderer.components.header_renderer import render_header
        from ppt_renderer.components.footer_renderer import render_footer

        pattern_id = getattr(layout_spec, "visual_pattern", None)
        render_header(layout_spec.header, presentation, slide, slide_spec, pattern_id=pattern_id)

        layout_context = {"pattern_id": pattern_id}
        for component in layout_spec.components:
            component_dispatcher.render(
                component, presentation, slide, slide_spec, layout_context=layout_context
            )

        render_footer(
            layout_spec.footer,
            presentation,
            slide,
            slide_spec,
            len(presentation.slides),
            pattern_id=pattern_id,
        )
