"""
ppt_renderer/components/matrix_renderer.py
==========================================
Matrix component renderer.

Draws matrix cells and capability-map columns using normalized bounds.
"""

from __future__ import annotations

from typing import Any

from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches

from pptx.dml.color import RGBColor

from ppt_renderer.components.coordinates import convert_bounds
from ppt_renderer.components.placeholder_resolver import resolve_placeholder
from backend.design_system import design_language
from backend.design_system.theme_loader import get_current_theme


def render(
    component_specification,
    presentation,
    slide,
    content: dict[str, Any],
    *,
    layout_context: dict[str, Any] | None = None,
) -> None:
    """Render a matrix cell or capability column."""
    theme = get_current_theme()
    left, top, width, height = convert_bounds(component_specification, presentation)
    pattern_id = (layout_context or {}).get("pattern_id")

    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        left,
        top,
        width,
        height,
    )
    shape.fill.solid()

    item = resolve_placeholder(component_specification.placeholder, content)
    bg_color = _quadrant_color(item) or theme.color("panel_grey")
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = theme.color("border")

    padding_x = design_language.get(
        "matrix", "padding_x", pattern_id=pattern_id, category="spacing",
        default=theme.space("box_padding_x").inches,
    )
    padding_y = design_language.get(
        "matrix", "padding_y", pattern_id=pattern_id, category="spacing",
        default=theme.space("box_padding_y").inches,
    )
    alignment = _to_pp_align(design_language.get_alignment("matrix", pattern_id=pattern_id))

    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    tf.margin_left = Inches(padding_x)
    tf.margin_right = Inches(padding_x)
    tf.margin_top = Inches(padding_y)
    tf.margin_bottom = Inches(padding_y)

    if component_specification.type == "column" and isinstance(item, dict):
        # Capability-map column: domain name + capabilities.
        p1 = tf.paragraphs[0]
        p1.text = str(item.get("name", ""))
        p1.alignment = alignment
        p1.font.bold = True
        p1.font.name = theme.font("body")
        p1.font.size = theme.size("body")
        p1.font.color.rgb = theme.color("ink")

        capabilities = item.get("capabilities", [])
        for capability in capabilities[:4]:
            cap_name = capability.get("name", "") if isinstance(capability, dict) else str(capability)
            p = tf.add_paragraph()
            p.text = f"• {cap_name}"
            p.alignment = alignment
            p.font.name = theme.font("body")
            p.font.size = theme.size("small")
            p.font.color.rgb = theme.color("charcoal")
    else:
        value = item.get("value", "") if isinstance(item, dict) else str(item)
        p = tf.paragraphs[0]
        p.text = str(value)
        p.alignment = alignment
        p.font.bold = True
        p.font.name = theme.font("body")
        p.font.size = theme.size("body")
        p.font.color.rgb = theme.color("ink")


def _to_pp_align(value: str):
    """Map a design-language alignment name to a python-pptx constant."""
    mapping = {
        "center": PP_ALIGN.CENTER,
        "left": PP_ALIGN.LEFT,
        "right": PP_ALIGN.RIGHT,
        "justify": PP_ALIGN.JUSTIFY,
    }
    return mapping.get(str(value).lower(), PP_ALIGN.CENTER)


def _quadrant_color(item: Any) -> Any:
    """Return a risk-quadrant tint for matrix cells, or None for default fill."""
    if not isinstance(item, dict):
        return None
    quadrant = item.get("quadrant")
    if not isinstance(quadrant, dict):
        return None
    impact = str(quadrant.get("impact", "")).lower()
    likelihood = str(quadrant.get("likelihood", "")).lower()

    if impact == "high" and likelihood == "high":
        return RGBColor(255, 120, 120)
    if impact == "low" and likelihood == "low":
        return RGBColor(120, 220, 120)
    if impact in {"high", "medium"} and likelihood in {"high", "medium"}:
        return RGBColor(255, 190, 100)
    if impact == "high" or likelihood == "high":
        return RGBColor(255, 210, 130)
    if impact == "medium" or likelihood == "medium":
        return RGBColor(255, 230, 150)
    return None
