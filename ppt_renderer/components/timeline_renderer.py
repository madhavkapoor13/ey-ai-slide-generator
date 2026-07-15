"""
ppt_renderer/components/timeline_renderer.py
============================================
Timeline / process component renderer.

Draws nodes, axis lines, and phase bars for timelines, roadmaps, journeys,
and process flows.
"""

from __future__ import annotations

from typing import Any

from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches

from ppt_renderer.components.coordinates import convert_bounds
from ppt_renderer.components.placeholder_resolver import resolve_placeholder
from backend.design_system import design_language
from backend.design_system.theme_loader import get_current_theme


def render(
    component_specification,
    presentation,
    slide,
    content: dict[str, Any],
    *,
    layout_context: dict[str, Any] | None = None,
) -> None:
    """Render a timeline/process component (node, axis, or bar)."""
    component_type = component_specification.type

    if component_type == "axis":
        _render_axis(component_specification, presentation, slide, layout_context)
    elif component_type == "bar":
        _render_bar(component_specification, presentation, slide, content, layout_context)
    else:
        _render_node(component_specification, presentation, slide, content, layout_context)


def _render_axis(component_specification, presentation, slide, layout_context) -> None:
    theme = get_current_theme()
    left, top, width, height = convert_bounds(component_specification, presentation)
    pattern_id = (layout_context or {}).get("pattern_id")
    constraints = component_specification.constraints or []
    # Center the line vertically within the component bounds.
    line_y = top + height // 2
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        left,
        line_y,
        left + width,
        line_y,
    )
    connector.line.color.rgb = theme.color("connector")
    connector_width = design_language.get(
        "timeline", "connector_width", pattern_id=pattern_id, category="spacing",
        default=theme.space("connector_width").inches,
    )
    connector.line.width = Inches(connector_width)

    if "arrow_ends" in constraints or "arrow_right" in constraints:
        _set_arrowhead(connector)


def _set_arrowhead(connector) -> None:
    """Add a triangle arrowhead to the end of a connector via OOXML."""
    from lxml import etree

    P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
    A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"

    spPr = connector._element.find(f"{{{P_NS}}}spPr")
    if spPr is None:
        return
    ln = spPr.find(f"{{{A_NS}}}ln")
    if ln is None:
        return
    tail = ln.find(f"{{{A_NS}}}tailEnd")
    if tail is None:
        tail = etree.SubElement(ln, f"{{{A_NS}}}tailEnd")
    tail.set("type", "triangle")
    tail.set("w", "med")
    tail.set("len", "med")


def _render_bar(component_specification, presentation, slide, content: dict[str, Any], layout_context) -> None:
    theme = get_current_theme()
    left, top, width, height = convert_bounds(component_specification, presentation)
    pattern_id = (layout_context or {}).get("pattern_id")

    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        left,
        top,
        width,
        height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = theme.color("panel_grey")
    shape.line.color.rgb = theme.color("border")

    item = resolve_placeholder(component_specification.placeholder, content)
    if isinstance(item, dict):
        label = item.get("name", "")
        description = item.get("description", "")
    else:
        label = str(item)
        description = ""

    alignment = _to_pp_align(design_language.get_alignment("roadmap", pattern_id=pattern_id))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = str(label)
    p.alignment = alignment
    p.font.bold = True
    p.font.name = theme.font("body")
    p.font.size = theme.size("small")
    p.font.color.rgb = theme.color("ink")

    if description:
        p2 = tf.add_paragraph()
        p2.text = str(description)
        p2.alignment = alignment
        p2.font.name = theme.font("body")
        p2.font.size = theme.size("tiny")
        p2.font.color.rgb = theme.color("charcoal")
        p2.space_before = theme.size("tiny")


def _render_node(component_specification, presentation, slide, content: dict[str, Any], layout_context) -> None:
    theme = get_current_theme()
    left, top, width, height = convert_bounds(component_specification, presentation)
    pattern_id = (layout_context or {}).get("pattern_id")

    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        left,
        top,
        width,
        height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = theme.color("light_grey")
    shape.line.color.rgb = theme.color("border")
    shape.line.width = theme.space("box_border_width")

    item = resolve_placeholder(component_specification.placeholder, content)
    if isinstance(item, dict):
        label = item.get("label", item.get("title", item.get("name", "")))
        description = item.get("description", "")
    else:
        label = str(item)
        description = ""

    padding_x = design_language.get(
        "matrix", "padding_x", pattern_id=pattern_id, category="spacing",
        default=theme.space("box_padding_x").inches,
    )
    alignment = _to_pp_align(design_language.get_alignment("timeline", pattern_id=pattern_id))

    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    tf.margin_left = Inches(padding_x)
    tf.margin_right = Inches(padding_x)
    p = tf.paragraphs[0]
    p.text = str(label)
    p.alignment = alignment
    p.font.bold = True
    p.font.name = theme.font("body")
    p.font.size = theme.size("body")
    p.font.color.rgb = theme.color("ink")

    if description:
        p2 = tf.add_paragraph()
        p2.text = str(description)
        p2.alignment = alignment
        p2.font.name = theme.font("body")
        p2.font.size = theme.size("tiny")
        p2.font.color.rgb = theme.color("charcoal")
        p2.space_before = theme.size("tiny")


def _to_pp_align(value: str):
    """Map a design-language alignment name to a python-pptx constant."""
    mapping = {
        "center": PP_ALIGN.CENTER,
        "left": PP_ALIGN.LEFT,
        "right": PP_ALIGN.RIGHT,
        "justify": PP_ALIGN.JUSTIFY,
    }
    return mapping.get(str(value).lower(), PP_ALIGN.CENTER)
