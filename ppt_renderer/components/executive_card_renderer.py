"""
ppt_renderer/components/executive_card_renderer.py
===================================================
Design Sprint D1 — Executive Insight Card renderer.

A production-quality reusable card component for Executive Summary, Key
Benefits, Strategic Pillars, Business Outcomes, Value Drivers,
Recommendations, Risks, and Opportunities.

The renderer consumes an ``ExecutiveCardContent`` object (or plain dict) and
draws a themed card inside the normalized component bounds supplied by the
Layout Engine. No business logic is inferred.
"""

from __future__ import annotations

from typing import Any

from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches

from backend.design_system import design_language
from backend.design_system.theme_loader import get_current_theme
from ppt_renderer.components.coordinates import convert_bounds
from ppt_renderer.components.placeholder_resolver import resolve_placeholder
from schemas.executive_card import ExecutiveCardContent


def render(
    component_specification,
    presentation,
    slide,
    content: dict[str, Any],
    *,
    layout_context: dict[str, Any] | None = None,
) -> None:
    """Render an Executive Insight Card component."""
    theme = get_current_theme()
    left, top, width, height = convert_bounds(component_specification, presentation)
    pattern_id = (layout_context or {}).get("pattern_id")

    # Resolve content. Accept either an ExecutiveCardContent model or a dict.
    raw = resolve_placeholder(component_specification.placeholder, content)
    if isinstance(raw, ExecutiveCardContent):
        card = raw
    elif isinstance(raw, dict):
        card = ExecutiveCardContent.model_validate(raw)
    else:
        card = ExecutiveCardContent(title=str(raw))

    spacing = design_language.get_card_spacing(pattern_id=pattern_id)
    padding = spacing.get("padding", theme.space("card_padding").inches)
    small_gap = spacing.get("gap", theme.space("small").inches)

    # Card background.
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        left,
        top,
        width,
        height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = theme.color("panel_grey")
    shape.line.color.rgb = theme.color("border")
    shape.line.width = Inches(spacing.get("border_width", theme.space("card_border_width").inches))

    inner_left = left + Inches(padding)
    inner_width = width - Inches(padding * 2)
    cursor_y = top + Inches(padding)

    # Optional highlight badge.
    if card.highlight:
        badge_spec = design_language.get("badge", pattern_id=pattern_id, category="spacing", default={})
        badge_height_inches = badge_spec.get("height", 0.25)
        badge_width_ratio = badge_spec.get("width_ratio", 0.45)
        badge_height = Inches(badge_height_inches)
        badge = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            inner_left,
            cursor_y,
            inner_width * badge_width_ratio,
            badge_height,
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = theme.color("primary")
        badge.line.fill.background()

        tf = badge.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = str(card.highlight)
        p.alignment = PP_ALIGN.CENTER
        p.font.name = theme.font("body")
        p.font.size = theme.size("card_highlight")
        p.font.bold = True
        p.font.color.rgb = theme.color("ink")

        cursor_y += badge_height + Inches(small_gap)

    # Title.
    title_height = Inches(design_language.get("card_title", "height", pattern_id=pattern_id, category="spacing", default=0.35))
    title_box = slide.shapes.add_textbox(
        inner_left,
        cursor_y,
        inner_width,
        title_height,
    )
    tf = title_box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = card.title
    p.font.name = theme.font("title")
    p.font.size = theme.size("card_title")
    p.font.bold = True
    p.font.color.rgb = theme.color("ink")
    cursor_y += title_height + Inches(small_gap)

    # Description.
    desc_height_inches = design_language.get("card_description", "height", pattern_id=pattern_id, category="spacing", default=0.65)
    if card.metric:
        desc_height_inches = design_language.get("card_description", "metric_height", pattern_id=pattern_id, category="spacing", default=0.45)
    desc_height = Inches(desc_height_inches)
    desc_box = slide.shapes.add_textbox(
        inner_left,
        cursor_y,
        inner_width,
        desc_height,
    )
    tf = desc_box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = card.description
    p.font.name = theme.font("body")
    p.font.size = theme.size("card_description")
    p.font.color.rgb = theme.color("charcoal")
    cursor_y += desc_height + Inches(small_gap)

    # Optional metric.
    if card.metric:
        metric_height = Inches(design_language.get("card_metric", "height", pattern_id=pattern_id, category="spacing", default=0.35))
        metric_box = slide.shapes.add_textbox(
            inner_left,
            cursor_y,
            inner_width,
            metric_height,
        )
        tf = metric_box.text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = str(card.metric)
        p.font.name = theme.font("title")
        p.font.size = theme.size("card_metric")
        p.font.bold = True
        p.font.color.rgb = theme.color("ink")
        cursor_y += metric_height + Inches(small_gap)

    # Optional footer tag.
    if card.tag:
        tag_height = Inches(design_language.get("card_tag", "height", pattern_id=pattern_id, category="spacing", default=0.20))
        tag_box = slide.shapes.add_textbox(
            inner_left,
            cursor_y,
            inner_width,
            tag_height,
        )
        tf = tag_box.text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = str(card.tag)
        p.font.name = theme.font("body")
        p.font.size = theme.size("card_tag")
        p.font.color.rgb = theme.color("grey")
