"""
ppt_renderer/components/header_renderer.py
==========================================
Header component renderer.

Draws the slide title, subtitle, and description inside the normalized
header region.
"""

from __future__ import annotations

from typing import Any

from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt

from ppt_renderer.components.coordinates import (
    convert_height,
    convert_width,
    convert_x,
    convert_y,
)
from ppt_renderer.components.placeholder_resolver import resolve_placeholder
from backend.design_system import design_language
from backend.design_system.theme_loader import get_current_theme


def render(
    component_specification,
    presentation,
    slide,
    content: dict[str, Any],
    *,
    layout_context: dict[str, Any] | None = None,
) -> None:
    """Render a header component from a ComponentSpecification bounding box."""
    theme = get_current_theme()
    pattern_id = (layout_context or {}).get("pattern_id")
    title = resolve_placeholder("title", content)
    subtitle = resolve_placeholder("subtitle", content)
    description = resolve_placeholder("description", content)

    left = convert_x(component_specification.x, presentation)
    top = convert_y(component_specification.y, presentation)
    width = convert_width(component_specification.width, presentation)
    height = convert_height(component_specification.height, presentation)

    title_spec = design_language.get_typography_hierarchy("executive_title", pattern_id=pattern_id)
    subtitle_spec = design_language.get_typography_hierarchy("section_title", pattern_id=pattern_id)
    body_spec = design_language.get_typography_hierarchy("body", pattern_id=pattern_id)

    _draw_textbox(
        slide, left, top, width, height * 0.45, title,
        font=theme.font(title_spec.get("font", "title")),
        size=Pt(title_spec.get("size", theme.size("title").pt)),
        bold=title_spec.get("bold", True),
    )
    if subtitle:
        _draw_textbox(
            slide,
            left,
            top + height * 0.45,
            width,
            height * 0.30,
            subtitle,
            font=theme.font(subtitle_spec.get("font", "title")),
            size=Pt(subtitle_spec.get("size", theme.size("subtitle").pt)),
            bold=subtitle_spec.get("bold", True),
        )
    if description:
        _draw_textbox(
            slide,
            left,
            top + height * 0.75,
            width,
            height * 0.25,
            description,
            font=theme.font(body_spec.get("font", "body")),
            size=Pt(body_spec.get("size", theme.size("description").pt)),
            bold=body_spec.get("bold", False),
        )


def render_header(
    header_spec,
    presentation,
    slide,
    content: dict[str, Any],
    *,
    pattern_id: str | None = None,
) -> None:
    """Render the full header using a HeaderSpecification."""
    theme = get_current_theme()
    title = resolve_placeholder("title", content)
    subtitle = resolve_placeholder("subtitle", content)
    description = resolve_placeholder("description", content)

    title_spec = design_language.get_typography_hierarchy("executive_title", pattern_id=pattern_id)
    subtitle_spec = design_language.get_typography_hierarchy("section_title", pattern_id=pattern_id)
    body_spec = design_language.get_typography_hierarchy("body", pattern_id=pattern_id)

    title_area = header_spec.title_area or {}
    subtitle_area = header_spec.subtitle_area or {}
    description_area = _resolve_description_area(header_spec, subtitle_area)

    _draw_textbox(
        slide,
        convert_x(title_area.get("x", 0.05), presentation),
        convert_y(title_area.get("y", 0.04), presentation),
        convert_width(title_area.get("width", 0.9), presentation),
        convert_height(title_area.get("height", 0.06), presentation),
        title,
        font=theme.font(title_spec.get("font", "title")),
        size=Pt(title_spec.get("size", theme.size("title").pt)),
        bold=title_spec.get("bold", True),
    )

    if subtitle:
        _draw_textbox(
            slide,
            convert_x(subtitle_area.get("x", 0.05), presentation),
            convert_y(subtitle_area.get("y", 0.10), presentation),
            convert_width(subtitle_area.get("width", 0.9), presentation),
            convert_height(subtitle_area.get("height", 0.04), presentation),
            subtitle,
            font=theme.font(subtitle_spec.get("font", "title")),
            size=Pt(subtitle_spec.get("size", theme.size("subtitle").pt)),
            bold=subtitle_spec.get("bold", True),
        )

    if description:
        _draw_textbox(
            slide,
            convert_x(description_area.get("x", 0.05), presentation),
            convert_y(description_area.get("y", 0.14), presentation),
            convert_width(description_area.get("width", 0.9), presentation),
            convert_height(description_area.get("height", 0.04), presentation),
            description,
            font=theme.font(body_spec.get("font", "body")),
            size=Pt(body_spec.get("size", theme.size("description").pt)),
            bold=body_spec.get("bold", False),
        )

    # Accent divider sits at the bottom of the header, just above the body.
    spacing = design_language.get_header_spacing(pattern_id=pattern_id)
    divider_left = convert_x(0.05, presentation)
    divider_top = convert_y(_divider_y(header_spec, description_area), presentation)
    divider_width = convert_width(0.18, presentation)
    divider_height = convert_height(spacing.get("divider_height", 0.008), presentation)

    divider = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        divider_left,
        divider_top,
        divider_width,
        divider_height,
    )
    divider.fill.solid()
    divider.fill.fore_color.rgb = theme.color("ey_yellow")
    divider.line.fill.background()


def _draw_textbox(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    *,
    font: str,
    size,
    bold: bool = False,
):
    theme = get_current_theme()
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text or ""
    p.font.name = font
    p.font.size = size
    p.font.bold = bold
    p.font.color.rgb = theme.color("ink")
    return box


def _resolve_description_area(header_spec, subtitle_area: dict) -> dict:
    """
    Return the description bounding box, deriving a non-overlapping default
    when the layout does not specify ``description_area``.

    Subtitle and description used to share ``subtitle_area`` which caused
    them to render at the same y on every slide. Now we shift the
    description to sit just below the subtitle (subtitle.y + subtitle.height),
    keeping the same x/width as the subtitle.
    """
    explicit = getattr(header_spec, "description_area", None)
    if explicit:
        return dict(explicit)
    sub = dict(subtitle_area)
    sub_y = float(sub.get("y", 0.10))
    sub_h = float(sub.get("height", 0.04))
    # No extra gap: description sits flush below the subtitle (its own
    # height reserves whitespace) and stays above the accent divider.
    derived = dict(sub)
    derived["y"] = sub_y + sub_h + 0.01
    derived["height"] = 0.02
    return derived


def _divider_y(header_spec, description_area: dict) -> float:
    """Place the accent divider just below the description, clamped to the
    header height so it never intrudes into the body region."""
    desc_y = float(description_area.get("y", 0.14))
    desc_h = float(description_area.get("height", 0.02))
    return max(desc_y + desc_h, float(getattr(header_spec, "height", 0.15)) - 0.01)
