"""
ppt_renderer/components/card_renderer.py
==========================================
Card component renderer.

Draws card-style components such as insight cards, strategy cards, and
KPI cards using normalized bounds.
"""

from __future__ import annotations

from typing import Any

from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN

from ppt_renderer.components.coordinates import convert_bounds
from ppt_renderer.components.placeholder_resolver import resolve_placeholder
from backend.design_system import design_language
from backend.design_system.theme_loader import get_current_theme


def render(
    component_specification,
    presentation,
    slide,
    content: dict[str, Any],
    *,
    layout_context: dict[str, Any] | None = None,
) -> None:
    """Render a card or KPI card component."""
    theme = get_current_theme()
    left, top, width, height = convert_bounds(component_specification, presentation)
    pattern_id = (layout_context or {}).get("pattern_id")

    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        left,
        top,
        width,
        height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = theme.color("light_grey")
    shape.line.color.rgb = theme.color("border")
    shape.line.width = theme.space("box_border_width")

    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    tf.margin_left = theme.space("box_padding_x")
    tf.margin_right = theme.space("box_padding_x")
    tf.margin_top = theme.space("box_padding_y")
    tf.margin_bottom = theme.space("box_padding_y")

    item = resolve_placeholder(component_specification.placeholder, content)

    if component_specification.type == "kpi_card" and isinstance(item, dict):
        value = item.get("value", "")
        label = item.get("label", "")
        trend = item.get("trend", "")
        description = item.get("description", "")
        alignment = _to_pp_align(design_language.get_alignment("kpi_card", pattern_id=pattern_id))
        p1 = tf.paragraphs[0]
        p1.text = str(value)
        p1.alignment = alignment
        p1.font.bold = True
        p1.font.name = theme.font("body")
        p1.font.size = theme.size("kpi_value")
        p1.font.color.rgb = theme.color("ink")

        if label:
            p2 = tf.add_paragraph()
            p2.text = str(label)
            p2.alignment = alignment
            p2.font.name = theme.font("body")
            p2.font.size = theme.size("body")
            p2.font.color.rgb = theme.color("charcoal")
        if trend:
            p3 = tf.add_paragraph()
            p3.text = str(trend)
            p3.alignment = alignment
            p3.font.name = theme.font("body")
            p3.font.size = theme.size("small")
            p3.font.color.rgb = theme.color("accent")
        if description:
            p4 = tf.add_paragraph()
            p4.text = str(description)
            p4.alignment = alignment
            p4.font.name = theme.font("body")
            p4.font.size = theme.size("tiny")
            p4.font.color.rgb = theme.color("grey")
    elif isinstance(item, dict):
        title = item.get("title", "")
        description = item.get("description", "")
        p1 = tf.paragraphs[0]
        p1.text = str(title)
        p1.alignment = PP_ALIGN.LEFT
        p1.font.bold = True
        p1.font.name = theme.font("body")
        p1.font.size = theme.size("body")
        p1.font.color.rgb = theme.color("ink")

        if description:
            p2 = tf.add_paragraph()
            p2.text = str(description)
            p2.alignment = PP_ALIGN.LEFT
            p2.font.name = theme.font("body")
            p2.font.size = theme.size("small")
            p2.font.color.rgb = theme.color("grey")
    else:
        p = tf.paragraphs[0]
        p.text = str(item)
        p.alignment = PP_ALIGN.LEFT
        p.font.name = theme.font("body")
        p.font.size = theme.size("body")
        p.font.color.rgb = theme.color("ink")


def _to_pp_align(value: str):
    """Map a design-language alignment name to a python-pptx constant."""
    mapping = {
        "center": PP_ALIGN.CENTER,
        "left": PP_ALIGN.LEFT,
        "right": PP_ALIGN.RIGHT,
        "justify": PP_ALIGN.JUSTIFY,
    }
    return mapping.get(str(value).lower(), PP_ALIGN.LEFT)
