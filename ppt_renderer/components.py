from __future__ import annotations

from datetime import date
from typing import Any

from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from ppt_renderer.theme import EYTheme


class SlideComponents:
    """Reusable native PowerPoint drawing primitives."""

    @staticmethod
    def draw_title(slide, title: str, subtitle: str = "", description: str | None = None) -> None:
        title_box = slide.shapes.add_textbox(
            EYTheme.HEADER_TITLE_X,
            EYTheme.HEADER_TITLE_Y,
            EYTheme.HEADER_TITLE_W,
            EYTheme.HEADER_TITLE_H,
        )
        title_tf = title_box.text_frame
        title_tf.clear()
        title_tf.word_wrap = True

        title_p = title_tf.paragraphs[0]
        title_p.text = title
        title_p.font.bold = True
        title_p.font.name = EYTheme.TITLE_FONT_FAMILY
        title_p.font.size = EYTheme.TITLE_SIZE
        title_p.font.color.rgb = EYTheme.INK

        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                EYTheme.LEFT_MARGIN,
                EYTheme.HEADER_SUBTITLE_Y,
                EYTheme.HEADER_SUBTITLE_W,
                EYTheme.HEADER_SUBTITLE_H,
            )
            subtitle_tf = subtitle_box.text_frame
            subtitle_tf.clear()
            subtitle_tf.word_wrap = True
            subtitle_p = subtitle_tf.paragraphs[0]
            subtitle_p.text = subtitle
            subtitle_p.font.bold = True
            subtitle_p.font.name = EYTheme.FONT_FAMILY
            subtitle_p.font.size = EYTheme.SUBTITLE_SIZE
            subtitle_p.font.color.rgb = EYTheme.CHARCOAL

        if description:
            description_box = slide.shapes.add_textbox(
                EYTheme.LEFT_MARGIN,
                EYTheme.HEADER_DESCRIPTION_Y,
                EYTheme.HEADER_DESCRIPTION_W,
                EYTheme.HEADER_DESCRIPTION_H,
            )
            description_tf = description_box.text_frame
            description_tf.clear()
            description_tf.word_wrap = True
            description_p = description_tf.paragraphs[0]
            description_p.text = description
            description_p.font.name = EYTheme.FONT_FAMILY
            description_p.font.size = EYTheme.DESCRIPTION_SIZE
            description_p.font.color.rgb = EYTheme.GREY

        divider = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            EYTheme.LEFT_MARGIN,
            EYTheme.HEADER_DIVIDER_Y,
            EYTheme.HEADER_DIVIDER_W,
            EYTheme.HEADER_DIVIDER_H,
        )
        divider.fill.solid()
        divider.fill.fore_color.rgb = EYTheme.EY_YELLOW
        divider.line.fill.background()

    @staticmethod
    def draw_process_box(
        slide,
        node: dict[str, Any],
        x,
        y,
        width,
        height,
    ):
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            x,
            y,
            width,
            height,
        )

        shape.fill.solid()
        shape.fill.fore_color.rgb = EYTheme.LIGHT_GREY

        shape.line.color.rgb = EYTheme.BORDER
        shape.line.width = EYTheme.BOX_BORDER_WIDTH

        tf = shape.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        tf.margin_left = EYTheme.BOX_PADDING_X
        tf.margin_right = EYTheme.BOX_PADDING_X
        tf.margin_top = EYTheme.BOX_PADDING_Y
        tf.margin_bottom = EYTheme.BOX_PADDING_Y

        p = tf.paragraphs[0]
        p.text = node.get("label", "")
        p.alignment = PP_ALIGN.CENTER
        p.font.bold = True
        p.font.name = EYTheme.FONT_FAMILY
        p.font.size = EYTheme.BODY_SIZE
        p.font.color.rgb = EYTheme.INK

        return shape

    @staticmethod
    def draw_connector(slide, start_x, start_y, end_x, end_y):
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            start_x,
            start_y,
            end_x,
            end_y,
        )

        connector.line.color.rgb = EYTheme.CONNECTOR
        connector.line.width = EYTheme.CONNECTOR_WIDTH
        SlideComponents._add_arrowhead(connector)

        return connector

    @staticmethod
    def draw_pain_point(
        slide,
        x,
        y,
        width,
        height,
        text: str,
        anchor_x=None,
        anchor_y=None,
    ):
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            x,
            y,
            width,
            height,
        )

        shape.fill.solid()
        shape.fill.fore_color.rgb = EYTheme.LIGHT_RED

        shape.line.color.rgb = EYTheme.RED
        shape.line.width = EYTheme.PAIN_CALLOUT_BORDER_WIDTH

        tf = shape.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.margin_left = EYTheme.PAIN_CALLOUT_PADDING
        tf.margin_right = EYTheme.PAIN_CALLOUT_PADDING
        tf.margin_top = EYTheme.PAIN_CALLOUT_MARGIN_Y
        tf.margin_bottom = EYTheme.PAIN_CALLOUT_MARGIN_Y

        title = tf.paragraphs[0]
        title.text = "PAIN POINT"
        title.font.bold = True
        title.font.name = EYTheme.FONT_FAMILY
        title.font.size = EYTheme.CALLOUT_TITLE_SIZE
        title.font.color.rgb = EYTheme.RED

        body = tf.add_paragraph()
        body.text = text
        body.font.name = EYTheme.FONT_FAMILY
        body.font.size = EYTheme.CALLOUT_BODY_SIZE
        body.font.color.rgb = EYTheme.INK

        if anchor_x is not None and anchor_y is not None:
            SlideComponents.draw_callout_anchor(slide, anchor_x, anchor_y, x + width // 2, y)

        return shape

    @staticmethod
    def draw_callout_anchor(slide, start_x, start_y, end_x, end_y):
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            start_x,
            start_y + EYTheme.PAIN_ANCHOR_OFFSET,
            end_x,
            end_y - EYTheme.PAIN_ANCHOR_OFFSET,
        )
        connector.line.color.rgb = EYTheme.RED
        connector.line.width = EYTheme.PAIN_ANCHOR_WIDTH
        return connector

    @staticmethod
    def draw_footer(slide, slide_number: int = 1) -> None:
        rule = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            EYTheme.LEFT_MARGIN,
            EYTheme.FOOTER_RULE_Y,
            EYTheme.SLIDE_WIDTH - EYTheme.LEFT_MARGIN - EYTheme.RIGHT_MARGIN,
            EYTheme.FOOTER_RULE_H,
        )
        rule.fill.solid()
        rule.fill.fore_color.rgb = EYTheme.BORDER
        rule.line.fill.background()

        left = slide.shapes.add_textbox(
            EYTheme.LEFT_MARGIN,
            EYTheme.FOOTER_Y,
            EYTheme.FOOTER_LEFT_W,
            EYTheme.FOOTER_TEXT_H,
        )
        left_p = left.text_frame.paragraphs[0]
        left_p.text = "Generated by EY AI Pitch Prototype"
        SlideComponents._format_footer(left_p)

        center = slide.shapes.add_textbox(
            EYTheme.FOOTER_CENTER_X,
            EYTheme.FOOTER_Y,
            EYTheme.FOOTER_CENTER_W,
            EYTheme.FOOTER_TEXT_H,
        )
        center_p = center.text_frame.paragraphs[0]
        center_p.text = date.today().strftime("%d %b %Y")
        center_p.alignment = PP_ALIGN.CENTER
        SlideComponents._format_footer(center_p)

        right = slide.shapes.add_textbox(
            EYTheme.SLIDE_WIDTH - EYTheme.RIGHT_MARGIN - EYTheme.FOOTER_RIGHT_W,
            EYTheme.FOOTER_Y,
            EYTheme.FOOTER_RIGHT_W,
            EYTheme.FOOTER_TEXT_H,
        )
        right_p = right.text_frame.paragraphs[0]
        right_p.text = f"Slide {slide_number}"
        right_p.alignment = PP_ALIGN.RIGHT
        SlideComponents._format_footer(right_p)

    @staticmethod
    def _format_footer(paragraph) -> None:
        paragraph.font.name = EYTheme.FONT_FAMILY
        paragraph.font.size = EYTheme.FOOTER_SIZE
        paragraph.font.color.rgb = EYTheme.MID_GREY

    @staticmethod
    def _add_arrowhead(connector) -> None:
        # python-pptx keeps connectors editable but does not expose arrowheads
        # in this version, so the native PowerPoint line XML is set directly.
        line = connector.element.spPr.ln
        head = OxmlElement("a:headEnd")
        head.set("type", "triangle")
        head.set("w", "sm")
        head.set("len", "sm")
        line.append(head)
