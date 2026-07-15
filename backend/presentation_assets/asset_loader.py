"""
backend/presentation_assets/asset_loader.py
============================================
Thin python-pptx reader used by the Asset Inspector and the Populator.

Two responsibilities:

1. ``enumerate_shapes`` — given a ``.pptx`` path, list every shape on its
   first slide (recursing into groups), classifying each as a native
   placeholder or a free-form custom shape. Used by the Asset Inspector
   to author a manifest, and by tests / diagnostics.

2. ``open_for_population`` — given an ``asset_id``, resolve its
   ``asset.pptx`` via the registry, open it with python-pptx, and return
   the source presentation + its first slide. The Populator copies this
   slide into the target deck and fills placeholders.

This module never mutates files. It only reads.
"""

from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from backend.presentation_assets import asset_registry


@dataclass
class ShapeInfo:
    """
    Inspectable description of one shape on an asset slide.

    Attributes
    ----------
    name:
        ``shape.name`` (stable identifier for free-form shape binding).
    shape_type:
        MSO_SHAPE_TYPE name (e.g. ``"TEXT_BOX"``, ``"PLACEHOLDER"``,
        ``"GROUP"``, ``"AUTO_SHAPE"``).
    is_placeholder:
        ``True`` when the shape is a native PowerPoint placeholder
        (bindable by ``placeholder_idx``).
    placeholder_idx:
        The native placeholder index (``placeholder_format.idx``) when
        ``is_placeholder`` is ``True``; otherwise ``None``.
    has_text_frame:
        Whether the shape has an editable text frame.
    current_text:
        Current text content of the shape's text frame (empty string when
        no text frame).
    group_path:
        Ordered list of containing group shape names (empty for top-level
        shapes). Used to locate shapes nested inside groups by name.
    """

    name: str
    shape_type: str
    is_placeholder: bool
    placeholder_idx: int | None
    has_text_frame: bool
    current_text: str
    group_path: list[str] = field(default_factory=list)


def _shape_type_name(shape) -> str:
    try:
        return shape.shape_type.name
    except AttributeError:
        return str(getattr(shape, "shape_type", "UNKNOWN"))


def _collect_shapes(shapes, group_path: list[str], out: list[ShapeInfo]) -> None:
    for shape in shapes:
        name = shape.name or ""
        stype = _shape_type_name(shape)

        if stype == "GROUP":
            _collect_shapes(shape.shapes, group_path + [name], out)
            continue

        is_placeholder = bool(getattr(shape, "is_placeholder", False))
        placeholder_idx: int | None = None
        if is_placeholder:
            try:
                placeholder_idx = int(shape.placeholder_format.idx)
            except (AttributeError, TypeError):
                placeholder_idx = None

        has_text_frame = shape.has_text_frame
        current_text = ""
        if has_text_frame:
            current_text = shape.text_frame.text or ""

        out.append(
            ShapeInfo(
                name=name,
                shape_type=stype,
                is_placeholder=is_placeholder,
                placeholder_idx=placeholder_idx,
                has_text_frame=has_text_frame,
                current_text=current_text,
                group_path=list(group_path),
            )
        )


def enumerate_shapes(pptx_path: str | Path) -> list[ShapeInfo]:
    """
    Inspect the first slide of ``pptx_path`` and return all shapes on it.

    Group shapes are flattened recursively: every leaf shape is returned
    with its ``group_path`` recording the nesting. The slide's shape order
    is preserved.
    """
    prs = Presentation(str(pptx_path))
    if not prs.slides:
        return []
    slide = prs.slides[0]
    out: list[ShapeInfo] = []
    _collect_shapes(slide.shapes, [], out)
    return out


def open_for_population(asset_id: str, assets_dir: Path | None = None):
    """
    Open the ``asset.pptx`` for ``asset_id`` and return
    ``(presentation, slide)`` for the populator to copy and fill.

    When ``assets_dir`` is supplied, a one-off scan of that directory is
    used (test isolation). Otherwise the cached registry index is used.

    Raises:
        FileNotFoundError: if the asset id is not registered or its
            ``asset.pptx`` file is missing on disk.
    """
    asset_dir = asset_registry.get_asset_path(asset_id, assets_dir=assets_dir)
    if asset_dir is None:
        raise FileNotFoundError(f"Asset not registered: {asset_id!r}")

    pptx_path = asset_dir / "asset.pptx"
    if not pptx_path.exists():
        raise FileNotFoundError(f"asset.pptx not found for asset {asset_id!r}: {pptx_path}")

    prs = Presentation(str(pptx_path))
    if not prs.slides:
        raise ValueError(f"Asset {asset_id!r} has no slides in {pptx_path}")

    return prs, prs.slides[0]