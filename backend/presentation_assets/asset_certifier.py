"""
backend/presentation_assets/asset_certifier.py
==============================================
Certification checks for Presentation Assets.

Certification does not mutate assets. It validates that an asset manifest and
its PPTX can safely participate in production generation.
"""

from __future__ import annotations

import hashlib
import tempfile
from datetime import datetime, timezone
from pathlib import Path
from typing import Any
from zipfile import ZipFile

from pptx import Presentation

from backend.presentation_assets import asset_loader, asset_populator, asset_registry
from backend.presentation_assets.asset_loader import enumerate_shapes
from backend.presentation_assets.manifest_conformance import check_conformance
from backend.presentation_assets.text_fit import check_text_fit
from schemas.presentation_asset import AssetCertification, AssetManifest
from schemas.slide_spec import SlideSpec


SUPPORTED_WARN_TAGS = ("chart", "oleObject", "externalLink")


def certify_asset(asset_id: str, *, assets_dir: Path | None = None) -> AssetCertification:
    """Certify one registered asset and return metadata suitable for the manifest."""
    manifest = asset_registry.get(asset_id, assets_dir=assets_dir)
    if manifest is None:
        return AssetCertification(
            certified=False,
            certified_at=_now(),
            errors=[f"asset {asset_id!r} is not registered"],
        )

    asset_dir = asset_registry.get_asset_path(asset_id, assets_dir=assets_dir)
    if asset_dir is None:
        return AssetCertification(
            certified=False,
            certified_at=_now(),
            errors=[f"asset path not found for {asset_id!r}"],
        )

    pptx_path = asset_dir / "asset.pptx"
    warnings: list[str] = []
    errors: list[str] = []

    if not pptx_path.exists():
        errors.append(f"asset.pptx missing: {pptx_path}")
    else:
        warnings.extend(_unsupported_feature_warnings(pptx_path))

    if not errors:
        errors.extend(_binding_errors(pptx_path, manifest))
        errors.extend(_repeating_errors(pptx_path, manifest))
        warnings.extend(_text_constraint_warnings(manifest))
        errors.extend(_placeholder_leakage_errors(pptx_path, manifest))
        errors.extend(_smoke_population_errors(asset_id, manifest, assets_dir=assets_dir))

    preview_hash = _hash_file(pptx_path) if pptx_path.exists() else None
    return AssetCertification(
        certified=not errors,
        certified_at=_now(),
        preview_hash=preview_hash,
        warnings=warnings,
        errors=errors,
    )


def certify_all(*, assets_dir: Path | None = None) -> dict[str, AssetCertification]:
    """Certify every registered asset."""
    return {
        asset_id: certify_asset(asset_id, assets_dir=assets_dir)
        for asset_id in sorted(asset_registry.load_assets(assets_dir))
    }


def _binding_errors(pptx_path: Path, manifest: AssetManifest) -> list[str]:
    shapes = enumerate_shapes(pptx_path)
    names = {shape.name for shape in shapes}
    placeholder_indices = {
        shape.placeholder_idx for shape in shapes if shape.placeholder_idx is not None
    }
    errors: list[str] = []
    for placeholder in manifest.placeholders:
        binding = placeholder.binding
        if binding.shape_name:
            concrete_names = _binding_names(placeholder, manifest)
            missing = [name for name in concrete_names if name not in names]
            if missing and placeholder.required:
                errors.append(f"placeholder {placeholder.id!r} missing shapes: {missing}")
        elif binding.native_placeholder_idx is not None:
            if binding.native_placeholder_idx not in placeholder_indices and placeholder.required:
                errors.append(
                    f"placeholder {placeholder.id!r} missing native idx {binding.native_placeholder_idx}"
                )
        elif placeholder.required:
            errors.append(f"placeholder {placeholder.id!r} has no binding")
    return errors


def _repeating_errors(pptx_path: Path, manifest: AssetManifest) -> list[str]:
    if manifest.repeating is None:
        return []
    shapes = enumerate_shapes(pptx_path)
    names = {shape.name for shape in shapes}
    errors: list[str] = []
    lo, hi = manifest.density_range
    if not (lo <= manifest.repeating.count <= hi):
        errors.append(
            f"repeating count {manifest.repeating.count} outside density_range {manifest.density_range}"
        )
    for placeholder in manifest.placeholders:
        if placeholder.cardinality != "N" or not placeholder.binding.shape_name:
            continue
        missing = [name for name in _binding_names(placeholder, manifest) if name not in names]
        if missing and placeholder.required:
            errors.append(f"repeating placeholder {placeholder.id!r} missing shapes: {missing}")
    return errors


def _binding_names(placeholder, manifest: AssetManifest) -> list[str]:
    shape_name = placeholder.binding.shape_name
    if not shape_name:
        return []
    if placeholder.cardinality == "N" and manifest.repeating is not None:
        token = manifest.repeating.index_token
        return [
            shape_name.replace(token, str(index))
            for index in range(1, manifest.repeating.count + 1)
        ]
    return [shape_name]


def _text_constraint_warnings(manifest: AssetManifest) -> list[str]:
    warnings: list[str] = []
    for placeholder in manifest.placeholders:
        if placeholder.kind.value in {"text", "title", "subtitle", "body"}:
            constraints = placeholder.constraints or {}
            if "max_chars" not in constraints and "max_lines" not in constraints:
                warnings.append(f"placeholder {placeholder.id!r} missing text-fit constraints")
    return warnings


def _placeholder_leakage_errors(pptx_path: Path, manifest: AssetManifest) -> list[str]:
    bound_names: set[str] = set()
    for placeholder in manifest.placeholders:
        bound_names.update(_binding_names(placeholder, manifest))
    leaks = {"text", "title", "subtitle", "placeholder", "lorem ipsum"}
    errors: list[str] = []
    try:
        prs = Presentation(str(pptx_path))
        slide = prs.slides[0]
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            normalized = " ".join(str(text).strip().lower().split())
            if normalized in leaks and getattr(shape, "name", "") not in bound_names:
                errors.append(f"unbound placeholder leakage shape {getattr(shape, 'name', '?')!r}: {text!r}")
    except Exception as exc:  # noqa: BLE001
        errors.append(f"could not inspect placeholder leakage: {exc}")
    return errors


def _smoke_population_errors(
    asset_id: str,
    manifest: AssetManifest,
    *,
    assets_dir: Path | None = None,
) -> list[str]:
    content = _synthetic_content(manifest)
    errors = check_conformance(content, manifest)
    fit = check_text_fit(content, manifest)
    if not fit.passed:
        errors.extend(f"text-fit failed in smoke content: {f.placeholder_id}" for f in fit.failures)
    if errors:
        return errors

    try:
        target = Presentation()
        slide_spec = SlideSpec(
            slide_type="operating_model",
            raw_spec=content,
            asset_id=asset_id,
        )
        asset_populator.populate_asset_slide(target, slide_spec, manifest, assets_dir=assets_dir)
        with tempfile.NamedTemporaryFile(suffix=".pptx") as tmp:
            target.save(tmp.name)
            Presentation(tmp.name)
    except Exception as exc:  # noqa: BLE001 - certification should report, not raise.
        return [f"smoke population failed: {exc}"]
    return []


def _synthetic_content(manifest: AssetManifest) -> dict[str, Any]:
    content: dict[str, Any] = {}
    count = manifest.density
    for placeholder in manifest.placeholders:
        max_chars = placeholder.constraints.get("max_chars")
        text = _sample_text(max_chars if isinstance(max_chars, int) else 24)
        if placeholder.cardinality == "N":
            content[placeholder.id] = [text for _ in range(count)]
        else:
            content[placeholder.id] = text
    return content


def _sample_text(max_chars: int) -> str:
    base = "Certified placeholder"
    if max_chars <= 0:
        return base
    return base[:max_chars]


def _unsupported_feature_warnings(pptx_path: Path) -> list[str]:
    warnings: list[str] = []
    try:
        with ZipFile(pptx_path) as package:
            for name in package.namelist():
                if name.endswith(".rels"):
                    text = package.read(name).decode("utf-8", errors="ignore")
                    for tag in SUPPORTED_WARN_TAGS:
                        if tag in text:
                            warnings.append(f"unsupported or fragile relationship feature detected: {tag}")
    except Exception as exc:  # noqa: BLE001
        warnings.append(f"could not inspect relationships: {exc}")
    return sorted(set(warnings))


def _hash_file(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as f:
        for chunk in iter(lambda: f.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()[:16]


def _now() -> str:
    return datetime.now(timezone.utc).isoformat()
