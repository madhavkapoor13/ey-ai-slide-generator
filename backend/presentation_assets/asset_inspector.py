"""
backend/presentation_assets/asset_inspector.py
==============================================
Asset Inspector — turns an opaque ``.pptx`` into a draft ``AssetManifest``.

The Inspector is **reusable**: it powers the offline ``ingest_asset.py``
CLI as well as a future admin UI and diagnostics ("why did my asset
populate wrong?"). It performs pure shape analysis on the first slide
and emits:

  - the list of ``ShapeInfo`` produced by ``asset_loader.enumerate_shapes``,
  - detected repeating groups (Phase1Label / Phase2Label / ...),
  - a proposed ``AssetManifest`` with auto-filled ``binding``, ``density``,
    ``repeating``, and ``kind`` proposals. Authored metadata (``family``,
    ``purpose``, ``audience_tags``, ``style_tags``, ``recommended_for``,
    ``avoid_for``, ``fits_content_kinds``, ``density_range``) is left blank
    or minimal — a human fills these during review.

The Inspector never writes files. Persisting the manifest is the CLI's
job; humans review and fill the authored metadata.

Design notes
------------
- Group shapes are flattened by ``asset_loader``; the Inspector sees leaf
  shapes only.
- Repeating detection: scans shape names for a trailing digit suffix
  matching ``^(.+?)(\\d+)$`` and groups by the stripped prefix. A group
  is "repeating" if it has ≥ 2 occurrences with sequential indices
  starting at 1.
- The Inspector does NOT try to guess family / purpose / tags — these are
  consulting-judgement decisions that a human must make.
"""

from __future__ import annotations

import re
from dataclasses import dataclass
from pathlib import Path
from typing import Optional

from pptx.enum.shapes import MSO_SHAPE_TYPE

from backend.presentation_assets.asset_loader import ShapeInfo, enumerate_shapes
from schemas.presentation_asset import (
    AssetManifest,
    AssetPlaceholder,
    PlaceholderBinding,
    PlaceholderKind,
    RepeatingGroup,
)


# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

# Match a single run of digits anywhere in the shape name. Group 1 is the
# prefix, group 2 is the digits, group 3 is the suffix. Used for repeating
# detection: e.g. "Phase1Label" -> ("Phase", "1", "Label").
_EMBEDDED_DIGITS_RE = re.compile(r"^(?P<prefix>.*?)(?P<digits>\d+)(?P<suffix>.*)$")
_DATE_HINT_RE = re.compile(
    r"\b(Q[1-4]\s*\d{4}|H[12]\s*\d{4}|\d{4}|Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)\b",
    re.IGNORECASE,
)
_CURRENCY_RE = re.compile(r"[$€£¥]")
_NUMERIC_RE = re.compile(r"^-?\d+(?:\.\d+)?$")
_TITLE_NATIVE_IDX = 0

# Family folders list — used to auto-derive a draft ``family`` when the
# .pptx is placed under presentation_assets/<known_family>/<asset_id>/.
_KNOWN_FAMILIES = frozenset(
    {
        "executive_summary",
        "operating_model",
        "roadmap",
        "timeline",
        "capability_map",
        "matrix",
        "process",
        "comparison",
        "strategy",
        "journey",
        "kpi",
        "business_benefits",
        "case_for_change",
        "risk",
        "next_steps",
        "opportunity_matrix",
    }
)


# ---------------------------------------------------------------------------
# Dataclass outputs (runtime-only; not boundary objects)
# ---------------------------------------------------------------------------


@dataclass
class RepeatingDetection:
    """One detected repeating pattern in shape names."""

    stripped_id: str          # snake_case of stripped prefix, e.g. "phase_label"
    shape_name_template: str  # "Phase{N}Label"
    indices: list[int]       # [1, 2, 3]


@dataclass
class InspectionReport:
    """Full output of ``inspect()``."""

    pptx_path: Path
    shapes: list[ShapeInfo]
    repeating: list[RepeatingDetection]
    supports_images: bool
    proposed_manifest: AssetManifest


# ---------------------------------------------------------------------------
# Heuristics
# ---------------------------------------------------------------------------


def _snake(s: str) -> str:
    """Convert 'PhaseLabel'/'Phase Label'/'phase-label'/'QDate' to 'phase_label'/'q_date'."""
    s = re.sub(r"[\s\-]+", "_", s.strip())
    s = re.sub(r"(.)([A-Z][a-z]+)", r"\1_\2", s)
    s = re.sub(r"([a-z0-9])([A-Z])", r"\1_\2", s)
    return s.lower()


def _propose_kind(shape: ShapeInfo) -> PlaceholderKind:
    """Heuristically infer the content kind from a shape's text/name."""
    text = (shape.current_text or "").strip()
    name = (shape.name or "").lower()

    if shape.is_placeholder and shape.placeholder_idx == _TITLE_NATIVE_IDX:
        return PlaceholderKind.TITLE

    if "chevron" in name:
        return PlaceholderKind.CHEVRON
    if "icon" in name:
        return PlaceholderKind.ICON
    if "node" in name or "axis" in name:
        return PlaceholderKind.TIMELINE_NODE

    if text:
        if "%" in text:
            return PlaceholderKind.PERCENTAGE
        if _CURRENCY_RE.search(text):
            return PlaceholderKind.CURRENCY
        if _DATE_HINT_RE.search(text):
            return PlaceholderKind.DATE
        if _NUMERIC_RE.match(text):
            return PlaceholderKind.METRIC

    return PlaceholderKind.BODY


def _propose_role(shape: ShapeInfo, stripped_id: Optional[str] = None) -> str:
    """Infer a role string for a placeholder."""
    if shape.is_placeholder and shape.placeholder_idx == _TITLE_NATIVE_IDX:
        return "title"
    if stripped_id:
        return stripped_id
    return _snake(shape.name)


def _detect_repeating(shapes: list[ShapeInfo]) -> dict[str, RepeatingDetection]:
    """
    Detect repeating patterns by a single embedded digit run in shape names.

    Handles names like "Phase1Label", "Phase2Label", ... (digits in the
    middle) as well as "Card1", "Card2", ... (digits at the end) and
    "Q1Metric" / "Q2Metric". Shapes are grouped by their full ``prefix`` +
    ``suffix`` (digits removed); a group is "repeating" if it has ≥ 2
    occurrences on a sequential 1..N run (so Phase1 + Phase3 with no
    Phase2 is NOT detected — the asset is malformed and the human reviews).

    Returns ``{stripped_id: RepeatingDetection}`` keyed by a snake_case
    id derived from the (prefix + suffix).
    """
    by_stripped: dict[str, list[tuple[int, ShapeInfo]]] = {}
    for shape in shapes:
        if shape.is_placeholder or not shape.name:
            continue
        m = _EMBEDDED_DIGITS_RE.match(shape.name)
        if not m:
            continue
        prefix, digits, suffix = m.group("prefix"), m.group("digits"), m.group("suffix")
        # Avoid matching pure-noise shapes (e.g. a stray digit in "FY24"
        # is not a repeating slot on its own). Require the prefix or
        # suffix to provide a non-trivial name stem.
        if not prefix and not suffix:
            continue
        n = int(digits)
        base = prefix + suffix
        stripped_id = _snake(base)
        by_stripped.setdefault(stripped_id, []).append((n, shape, m.group("digits")))

    out: dict[str, RepeatingDetection] = {}
    for stripped_id, entries in by_stripped.items():
        entries.sort(key=lambda pair: pair[0])
        indices = [pair[0] for pair in entries]
        if len(indices) < 2:
            continue
        if indices != list(range(1, indices[-1] + 1)):
            continue  # not a 1..N sequence
        first_shape = entries[0][1]
        # Build template by replacing the matched digit group with {N}.
        first_match_digits = entries[0][2]
        template_name = first_shape.name.replace(first_match_digits, "{N}", 1)
        out[stripped_id] = RepeatingDetection(
            stripped_id=stripped_id,
            shape_name_template=template_name,
            indices=indices,
        )
    return out


def _detect_picture(shapes: list[ShapeInfo]) -> bool:
    return any(s.shape_type == MSO_SHAPE_TYPE.PICTURE.name for s in shapes)


# ---------------------------------------------------------------------------
# Public entry point
# ---------------------------------------------------------------------------


def _derive_asset_id(pptx_path: Path) -> str:
    """
    Default asset_id: parent dir name when the file is named ``asset.pptx``,
    otherwise the stem of the file (uppercased, ``_``→``-``).
    """
    stem = pptx_path.stem
    if stem.lower() == "asset":
        return pptx_path.parent.name.upper()
    return stem.upper().replace("_", "-")


def _derive_family(pptx_path: Path) -> str:
    """
    Infer the family from the directory structure if it matches the
    Presentation Asset Library layout ``presentation_assets/<family>/<asset_id>/asset.pptx``.
    Returns the lowercased family name, or ``""`` if no match.
    """
    # Direct parent is the asset_id folder; the grandparent should be the family.
    family_dir_name = pptx_path.parent.parent.name
    if family_dir_name.lower() in _KNOWN_FAMILIES:
        return family_dir_name.lower()
    return ""


def inspect(
    pptx_path: str | Path,
    *,
    asset_id: Optional[str] = None,
    family: Optional[str] = None,
) -> InspectionReport:
    """
    Inspect a ``.pptx`` and return an :class:`InspectionReport` with a draft manifest.

    The draft auto-fills ``binding``, ``density``, ``repeating``, and
    ``kind`` proposals; ``family`` is pre-filled only when the file is
    placed under a known family directory. Authored metadata (``purpose``,
    ``audience_tags``, ``style_tags``, ``recommended_for``, ``avoid_for``,
    ``density_range``) is the human's responsibility during review.

    Parameters
    ----------
    pptx_path:
        Path to the asset ``.pptx`` file.
    asset_id:
        Optional asset id to stamp on the draft manifest. If omitted,
        derived from the filename stem (or parent directory when the file
        is the conventional ``asset.pptx``).
    family:
        Optional family name to pre-fill. Defaults to the parent
        directory name when it matches a known family folder; otherwise blank.
    """
    pptx_path = Path(pptx_path)
    shapes = enumerate_shapes(pptx_path)
    repeating_map = _detect_repeating(shapes)
    supports_images = _detect_picture(shapes)
    repeating_ids = set(repeating_map.keys())

    # Build the placeholder list: one placeholder per shape, except for
    # repeating members which collapse into a single template placeholder.
    placeholders: list[AssetPlaceholder] = []
    seen_repeating_ids: set[str] = set()

    for shape in shapes:
        # Only shapes with a text frame are fillable (text). Pictures,
        # pure decorators, and invisible layout shapes are skipped.
        if not shape.has_text_frame:
            continue

        # Repeating member? Emit one template placeholder per group.
        m = _EMBEDDED_DIGITS_RE.match(shape.name or "")
        if m and not shape.is_placeholder:
            base = m.group("prefix") + m.group("suffix")
            stripped_id = _snake(base)
            if stripped_id in repeating_ids:
                if stripped_id in seen_repeating_ids:
                    continue
                seen_repeating_ids.add(stripped_id)
                placeholders.append(
                    AssetPlaceholder(
                        id=stripped_id,
                        role=_propose_role(shape, stripped_id),
                        kind=_propose_kind(shape),
                        cardinality="N",
                        required=False,
                        binding=PlaceholderBinding(
                            shape_name=repeating_map[stripped_id].shape_name_template
                        ),
                    )
                )
                continue

        # Regular single-bind placeholder (native or named).
        if shape.is_placeholder:
            binding = PlaceholderBinding(native_placeholder_idx=shape.placeholder_idx)
        else:
            binding = PlaceholderBinding(shape_name=shape.name)

        if shape.is_placeholder and shape.placeholder_idx == _TITLE_NATIVE_IDX:
            pid = "title"
        else:
            pid = _snake(shape.name)

        placeholders.append(
            AssetPlaceholder(
                id=pid,
                role=_propose_role(shape),
                kind=_propose_kind(shape),
                cardinality="1",
                binding=binding,
            )
        )

    # Density: max repeating count (if any repeating detected);
    # else number of single-slot placeholders.
    if repeating_map:
        density = max(det.indices[-1] for det in repeating_map.values())
        repeating_group_meta = RepeatingGroup(
            group_template=next(iter(repeating_map.values())).shape_name_template,
            placeholders_per_group=sorted(repeating_map.keys()),
            index_token="{N}",
            count=density,
        )
    else:
        density = max(len(placeholders), 1)
        repeating_group_meta = None

    # Asset id and family (used in Sprint D); empty authored metadata.
    chosen_id = asset_id if asset_id else _derive_asset_id(pptx_path)
    chosen_family = family if family is not None else _derive_family(pptx_path)

    manifest = AssetManifest(
        asset_id=chosen_id,
        schema_version="1.0.0",
        family=chosen_family,
        family_aliases=[],
        purpose="",
        audience_tags=[],
        style_tags=[],
        recommended_for=[],
        avoid_for=[],
        density=density,
        density_range=[density, density],
        fits_content_kinds=[],
        supports_images=supports_images,
        placeholders=placeholders,
        repeating=repeating_group_meta,
    )

    return InspectionReport(
        pptx_path=pptx_path,
        shapes=shapes,
        repeating=list(repeating_map.values()),
        supports_images=supports_images,
        proposed_manifest=manifest,
    )
