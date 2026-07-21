"""
backend/presentation_assets/asset_populator.py
==============================================
Sprint E — Asset Populator.

Fills a Presentation Asset's ``.pptx`` slide using a manifest-shaped
``SlideSpec``. The Populator is the bridge between the Content Generator
(placeholder-keyed ``raw_spec``) and the actual PowerPoint shapes:

1. Resolve each ``AssetPlaceholder`` binding to a concrete shape on the slide.
2. Coerce the matching content value to text according to ``PlaceholderKind``.
3. Write the text into the shape's text frame.

For repeating placeholders the binding shape name may contain the repeating
 group's ``index_token`` (e.g. ``"Phase{N}Label"``); the token is replaced by
 the 1-based item index at population time. Items beyond the supplied content
 count have their target shapes cleared so stale template text is never
 rendered.

The module exposes two primary APIs:

- ``populate_slide(slide, content, manifest)`` — mutate an already-opened
  slide in place.
- ``populate_asset_slide(target_prs, slide_spec, manifest=None)`` — open the
  asset's ``asset.pptx``, copy its first slide into ``target_prs``, populate
  it, and return the new slide.
"""

from __future__ import annotations

import copy
import logging
from typing import Any

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt

from backend.presentation_assets import asset_loader, asset_registry
from schemas.presentation_asset import AssetManifest, AssetPlaceholder, PlaceholderKind
from schemas.slide_spec import SlideSpec

logger = logging.getLogger(__name__)

_DEFAULT_TEXT_COLOR = RGBColor(35, 35, 35)
_PLACEHOLDER_LEAKAGE_VALUES = {
    "text",
    "title",
    "subtitle",
    "placeholder",
    "lorem ipsum",
    "item 1",
    "step 1",
    "phase 1",
}

_PLACEHOLDER_LEAKAGE_PATTERNS = (
    "uc10 employee sentiment analysis",
    "uc11 automated meeting-summary generation",
    "uc12 synthetic-data generation",
    "uc13 autonomous procurement negotiation",
)


class PopulatorError(Exception):
    """Raised when the populator cannot resolve or populate an asset slide."""


def populate_asset_slide(
    target_prs: Presentation,
    slide_spec: SlideSpec,
    manifest: AssetManifest | None = None,
    *,
    assets_dir=None,
) -> Any:
    """
    Populate a Presentation Asset slide into ``target_prs``.

    Steps:

    1. Look up the asset manifest (from the supplied ``manifest`` or from the
       registry using ``slide_spec.asset_id``).
    2. Open the asset's ``asset.pptx`` via the loader.
    3. Copy the asset's first slide into ``target_prs``.
    4. Populate the copied slide using ``slide_spec.raw_spec`` and the manifest.

    Returns the newly added slide. Raises ``PopulatorError`` when the asset or
    its manifest cannot be found.
    """
    asset_id = slide_spec.asset_id
    if not asset_id:
        raise PopulatorError("SlideSpec has no asset_id; cannot populate a Presentation Asset slide.")

    resolved_manifest = manifest or asset_registry.get(asset_id, assets_dir=assets_dir)
    if resolved_manifest is None:
        raise PopulatorError(f"Asset manifest not found for {asset_id!r}")

    source_prs, source_slide = asset_loader.open_for_population(asset_id, assets_dir=assets_dir)
    new_slide, ph_name_map = _copy_slide(source_slide, target_prs)

    # Placeholders demoted to regular shapes can no longer be resolved by
    # native placeholder index. Patch the manifest to use shape names instead.
    patched_manifest = _patch_placeholder_bindings(resolved_manifest, ph_name_map)

    content = slide_spec.raw_spec if isinstance(slide_spec.raw_spec, dict) else {}
    populate_slide(new_slide, content, patched_manifest)
    return new_slide


def populate_slide(
    slide,
    content: dict[str, Any],
    manifest: AssetManifest,
) -> None:
    """
    Populate ``slide`` in place using placeholder-keyed ``content``.

    ``content`` must be a dict whose keys match the placeholder ids declared
    in ``manifest``. The manifest drives binding resolution, cardinality, and
    text coercion.
    """
    if not isinstance(content, dict):
        logger.warning("slide content is not a dict; nothing to populate")
        return

    for placeholder in manifest.placeholders:
        value = content.get(placeholder.id)
        _populate_placeholder(slide, placeholder, value, manifest)
    _clear_unbound_placeholder_leakage(slide)
    _log_layout_bounds_warnings(slide)


def _populate_placeholder(
    slide,
    placeholder: AssetPlaceholder,
    value: Any,
    manifest: AssetManifest,
) -> None:
    """Populate one placeholder, handling both single and repeating bindings."""
    if placeholder.cardinality == "1":
        target = _resolve_binding_target(slide, placeholder, index=None, manifest=manifest)
        if target is None:
            if placeholder.required:
                logger.warning("placeholder %r: bound target not found", placeholder.id)
            else:
                logger.debug("placeholder %r: optional bound target not found", placeholder.id)
            return
        _set_target_text(target, value, placeholder)
        return

    if placeholder.cardinality == "N":
        if value is None:
            value = []
        if not isinstance(value, list):
            logger.warning(
                "placeholder %r expects a list (cardinality N), got %s",
                placeholder.id,
                type(value).__name__,
            )
            value = []

        for index, item in enumerate(value, start=1):
            target = _resolve_binding_target(slide, placeholder, index=index, manifest=manifest)
            if target is None:
                if placeholder.required:
                    logger.warning(
                        "placeholder %r: bound target not found for index %d",
                        placeholder.id,
                        index,
                    )
                else:
                    logger.debug(
                        "placeholder %r: optional bound target not found for index %d",
                        placeholder.id,
                        index,
                    )
                continue
            _set_target_text(target, item, placeholder)

        # Clear any template shapes beyond the supplied content count so stale
        # text is not rendered.
        _clear_unused_repeating_shapes(slide, placeholder, len(value), manifest)
        return

    logger.warning("placeholder %r has unsupported cardinality %r", placeholder.id, placeholder.cardinality)


def _set_target_text(target, value: Any, placeholder: AssetPlaceholder) -> None:
    if hasattr(target, "text") and not hasattr(target, "shape_type"):
        target.text = _coerce_text(value, placeholder)
        return
    if hasattr(target, "text_frame"):
        _set_shape_text(target, value, placeholder)
        return
    if hasattr(target, "text"):
        target.text = _coerce_text(value, placeholder)
        return
    logger.debug("placeholder %r target cannot accept text", placeholder.id)


def _resolve_binding_target(
    slide,
    placeholder: AssetPlaceholder,
    index: int | None,
    manifest: AssetManifest,
):
    """Return a bound shape or table cell for ``placeholder``."""
    cell = _resolve_table_cell(slide, placeholder, index=index, manifest=manifest)
    if cell is not None:
        return cell
    return _resolve_shape(slide, placeholder, index=index, manifest=manifest)


def _resolve_table_cell(
    slide,
    placeholder: AssetPlaceholder,
    index: int | None,
    manifest: AssetManifest,
):
    binding = placeholder.binding
    table_name = binding.table_shape_name
    if not table_name:
        return None
    table_shape = _find_shape_by_name(slide.shapes, table_name)
    if table_shape is None or not getattr(table_shape, "has_table", False):
        return None
    row_index = binding.row_index
    col_index = binding.col_index
    if row_index is None or col_index is None:
        return None
    try:
        return table_shape.table.cell(row_index, col_index)
    except Exception:
        logger.warning("placeholder %r: table cell (%s, %s) not found", placeholder.id, row_index, col_index)
        return None


def _resolve_shape(
    slide,
    placeholder: AssetPlaceholder,
    index: int | None,
    manifest: AssetManifest,
):
    """Return the concrete shape for ``placeholder`` (and optional repeating index)."""
    binding = placeholder.binding

    if binding.native_placeholder_idx is not None:
        try:
            return slide.placeholders[binding.native_placeholder_idx]
        except (KeyError, IndexError):
            logger.debug(
                "native placeholder idx %d not found for %r",
                binding.native_placeholder_idx,
                placeholder.id,
            )
            return None

    shape_name = binding.shape_name or placeholder.id
    if index is not None and manifest.repeating is not None:
        token = manifest.repeating.index_token
        if token in shape_name:
            shape_name = shape_name.replace(token, str(index))

    return _find_shape_by_name(slide.shapes, shape_name)


def _find_shape_by_name(shapes, name: str):
    """Recursively search ``shapes`` for a shape whose ``name`` matches."""
    for shape in shapes:
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
            found = _find_shape_by_name(shape.shapes, name)
            if found is not None:
                return found
        if shape.name == name:
            return shape
    return None


def _set_shape_text(shape, value: Any, placeholder: AssetPlaceholder) -> None:
    """Write the coerced text for ``value`` into ``shape``."""
    if not getattr(shape, "has_text_frame", False):
        logger.debug("shape %r has no text frame; skipping", getattr(shape, "name", "?"))
        return

    text = _coerce_text(value, placeholder)
    text_frame = shape.text_frame
    style = _capture_text_style(text_frame)
    text_frame.clear()
    paragraph = text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text

    # Apply deterministic text styling so copied assets do not inherit
    # invisible white-on-white or arbitrary template placeholder styling.
    try:
        if placeholder.constraints.get("preserve_font_color", True):
            _apply_preserved_font_color(run, style)
        else:
            run.font.color.rgb = _DEFAULT_TEXT_COLOR
        if placeholder.constraints.get("preserve_font_size", True):
            _apply_preserved_font_size(run, style)
        else:
            size = placeholder.constraints.get("font_size_pt")
            run.font.size = Pt(int(size)) if isinstance(size, int) else Pt(12)
        if placeholder.constraints.get("preserve_bold", True):
            if style.get("bold") is not None:
                run.font.bold = style["bold"]
        else:
            run.font.bold = bool(placeholder.constraints.get("bold", False))
        if style.get("name"):
            run.font.name = style["name"]
        if style.get("italic") is not None:
            run.font.italic = style["italic"]
    except Exception:
        logger.debug("could not apply deterministic text style for %r", placeholder.id)


def _capture_text_style(text_frame) -> dict[str, Any]:
    """Capture the first concrete run style before clearing a text frame."""
    style: dict[str, Any] = {}
    try:
        paragraph = text_frame.paragraphs[0]
        for run in paragraph.runs:
            if not run.text:
                continue
            font = run.font
            style["size"] = font.size
            style["bold"] = font.bold
            style["italic"] = font.italic
            style["name"] = font.name
            try:
                style["color"] = font.color.rgb
            except Exception:
                style["color"] = None
            return style
        font = paragraph.font
        style["size"] = font.size
        style["bold"] = font.bold
        style["italic"] = font.italic
        style["name"] = font.name
        try:
            style["color"] = font.color.rgb
        except Exception:
            style["color"] = None
    except Exception:
        pass
    return style


def _apply_preserved_font_size(run, style: dict[str, Any]) -> None:
    size = style.get("size")
    if size is not None:
        run.font.size = size


def _apply_preserved_font_color(run, style: dict[str, Any]) -> None:
    color = style.get("color")
    if color is not None:
        run.font.color.rgb = color


def _clear_unbound_placeholder_leakage(slide) -> None:
    """Clear visible template defaults that were not bound in the manifest."""
    for shape in _iter_shapes(slide.shapes):
        if not getattr(shape, "has_text_frame", False):
            continue
        text = " ".join((shape.text or "").strip().split())
        normalized = text.lower()
        if _is_placeholder_leakage(normalized):
            shape.text_frame.clear()
            logger.warning("cleared leaked placeholder text from shape %r", getattr(shape, "name", "?"))


def _is_placeholder_leakage(normalized_text: str) -> bool:
    if normalized_text in _PLACEHOLDER_LEAKAGE_VALUES:
        return True
    stripped = normalized_text
    stripped = stripped.removeprefix("1. ").removeprefix("1) ").strip()
    return any(pattern in stripped for pattern in _PLACEHOLDER_LEAKAGE_PATTERNS)


def _log_layout_bounds_warnings(slide) -> None:
    """Best-effort layout bounds check for copied shapes."""
    # python-pptx does not expose the owning Presentation from a slide. Use the
    # standard widescreen bounds as a conservative smoke check.
    slide_width = 12192000
    slide_height = 6858000
    for shape in _iter_shapes(slide.shapes):
        try:
            if shape.left < 0 or shape.top < 0 or shape.left + shape.width > slide_width or shape.top + shape.height > slide_height:
                logger.warning("shape outside slide bounds: %s", getattr(shape, "name", "?"))
        except Exception:
            continue


def _iter_shapes(shapes):
    for shape in shapes:
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes(shape.shapes)
        else:
            yield shape


def _coerce_text(value: Any, placeholder: AssetPlaceholder) -> str:
    """Convert a content value to a string suitable for the shape text frame."""
    if value is None:
        return ""

    if isinstance(value, str):
        return value

    if isinstance(value, (int, float)):
        return str(value)

    if isinstance(value, dict):
        # Prefer a field named after the placeholder role or id.
        for key in (placeholder.role, placeholder.id):
            if key and key in value:
                return str(value[key])
        # Fallback: first string-ish value in the dict.
        for v in value.values():
            if isinstance(v, str):
                return v
            if isinstance(v, (int, float)):
                return str(v)
        return str(value)

    if isinstance(value, list):
        return ", ".join(str(item) for item in value)

    return str(value)


def _clear_unused_repeating_shapes(
    slide,
    placeholder: AssetPlaceholder,
    used_count: int,
    manifest: AssetManifest,
) -> None:
    """Remove text from repeating shapes that have no corresponding content item."""
    if manifest.repeating is None:
        return

    max_count = manifest.density_range[1]
    for index in range(used_count + 1, max_count + 1):
        shape = _resolve_shape(slide, placeholder, index=index, manifest=manifest)
        if shape is None:
            continue
        if getattr(shape, "has_text_frame", False):
            text_frame = shape.text_frame
            text_frame.clear()
            paragraph = text_frame.paragraphs[0]
            run = paragraph.add_run()
            run.text = ""


_R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
_P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"

_GRAPHIC_FRAME_TAG = f"{{{_P_NS}}}graphicFrame"
_CUST_DATA_LST_TAG = f"{{{_P_NS}}}custDataLst"
_BLIP_TAG = f"{{{_A_NS}}}blip"
_GRAPHIC_DATA_TAG = f"{{{_A_NS}}}graphicData"
_C_NV_PR_TAG = f"{{{_P_NS}}}cNvPr"
_NV_SP_PR_TAG = f"{{{_P_NS}}}nvSpPr"
_C_NV_SP_PR_TAG = f"{{{_P_NS}}}cNvSpPr"
_NV_PR_TAG = f"{{{_P_NS}}}nvPr"
_PH_TAG = f"{{{_P_NS}}}ph"
_SP_PR_TAG = f"{{{_P_NS}}}spPr"
_GRP_SP_PR_TAG = f"{{{_P_NS}}}grpSpPr"
_XFRM_TAG = f"{{{_A_NS}}}xfrm"
_OFF_TAG = f"{{{_A_NS}}}off"
_EXT_TAG = f"{{{_A_NS}}}ext"
_PRST_GEOM_TAG = f"{{{_A_NS}}}prstGeom"
_AV_LST_TAG = f"{{{_A_NS}}}avLst"

_REL_ATTRS = [
    f"{{{_R_NS}}}embed",
    f"{{{_R_NS}}}link",
    f"{{{_R_NS}}}id",
]


def _strip_all_rel_refs(el):
    """Recursively remove every r:embed, r:link, and r:id attribute."""
    for attr in _REL_ATTRS:
        el.attrib.pop(attr, None)
    for child in el:
        _strip_all_rel_refs(child)


def _restore_embedded_image_refs(source_el, copied_el, source_part, target_part) -> None:
    """Restore image relationships from ``source_el`` onto stripped ``copied_el``.

    SVG icons are stored as ``asvg:svgBlip r:embed`` extension elements rather
    than normal ``a:blip r:embed`` attributes. Walking the source and copied XML
    trees together preserves both forms without re-importing or re-parsing the
    image payload.
    """
    for source_node, copied_node in zip(source_el.iter(), copied_el.iter()):
        for attr in _REL_ATTRS:
            rel_id = source_node.get(attr)
            if not rel_id:
                continue
            try:
                rel = source_part.rels[rel_id]
                related_part = source_part.related_part(rel_id)
            except Exception:
                continue
            content_type = getattr(related_part, "content_type", "")
            if not content_type.startswith("image/"):
                continue
            try:
                copied_rel_id = target_part.relate_to(related_part, rel.reltype)
                copied_node.set(attr, copied_rel_id)
            except Exception:
                logger.warning("could not restore embedded image relationship %s", rel_id)


def _is_hidden(shape) -> bool:
    """Return True if the shape's non-visual properties mark it hidden."""
    el = getattr(shape, "element", None)
    if el is None:
        return False
    raw = etree.fromstring(etree.tostring(el))
    for cnv_pr in raw.iter(_C_NV_PR_TAG):
        if cnv_pr.get("hidden") == "1":
            return True
    return False


def _is_table_graphic_frame(shape) -> bool:
    """Return True for native PowerPoint tables encoded as graphicFrames."""
    el = getattr(shape, "element", None)
    if el is None:
        return False
    raw = etree.fromstring(etree.tostring(el))
    for graphic_data in raw.iter(_GRAPHIC_DATA_TAG):
        if graphic_data.get("uri") == "http://schemas.openxmlformats.org/drawingml/2006/table":
            return True
    return False


def _shape_has_placeholder(shape) -> bool:
    """Return True if the shape element declares a placeholder."""
    el = getattr(shape, "element", None)
    if el is None:
        return False
    return getattr(el, "has_ph_elm", False)


def _native_placeholder_idx(shape) -> int | None:
    """Return the native placeholder index (idx) for a source placeholder."""
    if not _shape_has_placeholder(shape):
        return None
    el = getattr(shape, "element", None)
    if el is None:
        return None
    raw = etree.fromstring(etree.tostring(el))
    for nv_pr in raw.iter(_NV_PR_TAG):
        ph = nv_pr.find(_PH_TAG)
        if ph is not None:
            idx = ph.get("idx")
            if idx is not None:
                return int(idx)
    return None


def _remove_placeholder_decl(shape_el) -> None:
    """Remove <p:ph> from a shape so it becomes a regular shape."""
    for nv_pr in shape_el.iter(_NV_PR_TAG):
        for ph in list(nv_pr.findall(_PH_TAG)):
            nv_pr.remove(ph)


def _patch_placeholder_bindings(manifest: AssetManifest, ph_name_map: dict[int, str]) -> AssetManifest:
    """Rewrite idx-based bindings to name-based bindings for demoted placeholders."""
    if not ph_name_map:
        return manifest

    new_placeholders = []
    for placeholder in manifest.placeholders:
        binding = placeholder.binding
        if binding.native_placeholder_idx is not None and binding.shape_name is None:
            name = ph_name_map.get(binding.native_placeholder_idx)
            if name is not None:
                new_binding = binding.model_copy(update={"shape_name": name, "native_placeholder_idx": None})
                new_placeholders.append(placeholder.model_copy(update={"binding": new_binding}))
                continue
        new_placeholders.append(placeholder)

    return manifest.model_copy(update={"placeholders": new_placeholders})


def _demote_placeholder_to_textbox(shape, shape_el) -> bool:
    """Convert a layout placeholder into a regular text box with explicit bounds.

    Source template shapes often rely on the slide layout for their geometry
    (their ``<p:spPr>`` is empty and they carry a ``<p:ph>`` declaration).
    When copied to a different presentation's blank layout those placeholders
    lose their geometry, so we materialise the effective bounds and add the
    geometry markup required for python-pptx (and PowerPoint) to treat them as
    normal text boxes.

    Returns ``True`` when geometry could be applied; ``False`` when the shape
    has no resolvable bounds and should be skipped.
    """
    try:
        left = int(shape.left)
        top = int(shape.top)
        width = int(shape.width)
        height = int(shape.height)
    except Exception:
        # Some shapes (e.g. placeholders with no resolved geometry) cannot be
        # measured; skip them so they do not become invalid shapes.
        return False

    # Mark the non-visual properties as a text box.
    cnv_sp_pr = shape_el.find(_C_NV_SP_PR_TAG)
    if cnv_sp_pr is None:
        nv_sp_pr = shape_el.find(_NV_SP_PR_TAG)
        if nv_sp_pr is None:
            return False
        cnv_sp_pr = etree.SubElement(nv_sp_pr, _C_NV_SP_PR_TAG)
    cnv_sp_pr.set("txBox", "1")

    # Ensure shape properties exist.
    pr_el = shape_el.find(_SP_PR_TAG)
    if pr_el is None:
        pr_el = etree.SubElement(shape_el, _SP_PR_TAG)

    # Replace any existing transform with explicit bounds.
    for existing in list(pr_el.findall(_XFRM_TAG)):
        pr_el.remove(existing)
    xfrm = etree.SubElement(pr_el, _XFRM_TAG)
    etree.SubElement(xfrm, _OFF_TAG, x=str(left), y=str(top))
    etree.SubElement(xfrm, _EXT_TAG, cx=str(width), cy=str(height))

    # Add a rectangle geometry so the shape is recognised as a valid auto shape.
    if pr_el.find(_PRST_GEOM_TAG) is None:
        geom = etree.SubElement(pr_el, _PRST_GEOM_TAG, prst="rect")
        etree.SubElement(geom, _AV_LST_TAG)

    return True


def _copy_slide(source_slide, target_prs: Presentation):
    """Return a new slide in ``target_prs`` cloned from ``source_slide``.

    The source EYP template slides contain think-cell OLE objects, embedded
    charts, and tag metadata whose package parts cannot be cleanly copied into
    a fresh presentation. Native PowerPoint tables are also graphicFrames but
    are inline DrawingML and are safe to copy. Rather than attempting to copy
    unsafe external parts (which produces broken relationship references and
    corrupts the file in PowerPoint), this function:

    1. **Skips** non-table ``<p:graphicFrame>`` shapes (OLE objects, charts)
       and any hidden think-cell tag shapes.
    2. **Demotes** layout placeholders to regular shapes and writes their
       effective geometry explicitly so positions survive the copy.
    3. **Strips** all ``r:link`` and unrelated ``r:id`` attributes from the
       remaining vector shapes and removes their ``<p:custDataLst>`` tag
       containers so no dangling references remain. Embedded image refs are
       copied into the target package and reattached so fixed asset icons stay
       exactly as authored.

    The visual formatting (geometry, fills, fonts, text) lives in the shape
    XML itself, so the slides retain their full appearance.

    Returns a tuple ``(new_slide, placeholder_name_map)`` where
    ``placeholder_name_map`` maps native placeholder indices to the copied
    shape names.  This lets callers patch manifest bindings that referenced
    placeholders by index.
    """
    new_slide = target_prs.slides.add_slide(_matching_layout(source_slide, target_prs))

    # Drop any placeholders that the blank layout inserted so they do not
    # overlap or conflict with the copied template shapes.
    for placeholder in list(new_slide.shapes):
        sp = getattr(placeholder, "element", None)
        if sp is not None and getattr(sp, "has_ph_elm", False):
            new_slide.shapes._spTree.remove(sp)

    placeholder_name_map: dict[int, str] = {}
    try:
        for idx, ph_shape in enumerate(source_slide.placeholders):
            placeholder_name_map[idx] = ph_shape.name
    except Exception:
        pass

    for shape in source_slide.shapes:
        tag = shape.element.tag

        # Skip OLE objects and charts — their parts can't be copied cleanly.
        # Keep native tables; they are inline DrawingML and do not need copied
        # package relationships.
        if tag == _GRAPHIC_FRAME_TAG and not _is_table_graphic_frame(shape):
            logger.debug("skipping graphicFrame: %s", getattr(shape, "name", "?"))
            continue

        # Skip hidden think-cell tag containers and other invisible shapes.
        if _is_hidden(shape):
            logger.debug("skipping hidden shape: %s", getattr(shape, "name", "?"))
            continue

        # Work on a raw lxml copy; ``copy.deepcopy()`` of a CT_* element
        # returns another CT_* wrapper whose child-finding API is awkward for
        # our namespace-aware stripping.
        new_el = etree.fromstring(etree.tostring(copy.deepcopy(shape.element)))

        # Remove <p:custDataLst> (tag metadata containers with r:id refs).
        for cust in list(new_el.iter(_CUST_DATA_LST_TAG)):
            cust.getparent().remove(cust)

        # Strip every remaining r:embed / r:link / r:id attribute so no
        # unresolved relationship references survive in the XML.
        _strip_all_rel_refs(new_el)
        _restore_embedded_image_refs(shape.element, new_el, source_slide.part, new_slide.part)

        # Placeholders copied from the source layout lose their geometry in the
        # target layout; convert them to regular text boxes with explicit bounds.
        if _shape_has_placeholder(shape):
            _remove_placeholder_decl(new_el)
            if not _demote_placeholder_to_textbox(shape, new_el):
                logger.debug(
                    "skipping placeholder with no geometry: %s",
                    getattr(shape, "name", "?"),
                )
                continue

        new_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")

    return new_slide, placeholder_name_map


def _matching_layout(source_slide, target_prs: Presentation):
    """Return the target layout matching the source slide's authored layout."""
    source_layout = getattr(source_slide, "slide_layout", None)
    source_name = getattr(source_layout, "name", None)
    if source_name:
        for layout in target_prs.slide_layouts:
            if getattr(layout, "name", None) == source_name:
                return layout

    try:
        source_layouts = list(source_slide.part.package.presentation_part.presentation.slide_layouts)
        source_index = source_layouts.index(source_layout)
        return target_prs.slide_layouts[source_index]
    except Exception:
        pass

    try:
        return target_prs.slide_layouts[6]
    except IndexError:
        return target_prs.slide_layouts[0]
