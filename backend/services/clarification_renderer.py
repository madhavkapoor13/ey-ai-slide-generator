"""
backend/services/clarification_renderer.py
==========================================
Sprint H.1 — Clarification rendering helper.

When the pipeline determines that clarification is required before a deck
can be generated, this module renders a minimal placeholder deck containing
the structured clarification questions. This keeps ``slide_service.py`` free
of placeholder-rendering logic while preserving the existing external API
contract: ``/generate/v2`` always returns a ``.pptx`` file.
"""

from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches, Pt

from schemas.clarification import ClarificationResult

# Widescreen 16:9 dimensions shared across renderers.
_SLIDE_WIDTH_IN = 13.333
_SLIDE_HEIGHT_IN = 7.5


def render_clarification_deck(output_path: str, clarification_result: ClarificationResult) -> str:
    """
    Render a single-slide placeholder deck containing clarification questions.

    Parameters
    ----------
    output_path:
        File path where the placeholder .pptx will be saved.
    clarification_result:
        Structured clarification questions from the Clarification Engine.

    Returns
    -------
    str
        The ``output_path`` that was written.
    """
    prs = Presentation()
    prs.slide_width = Inches(_SLIDE_WIDTH_IN)
    prs.slide_height = Inches(_SLIDE_HEIGHT_IN)

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.75), Inches(0.5), Inches(11.8), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_para = title_frame.paragraphs[0]
    title_para.text = "More Information Needed"
    title_para.font.size = Pt(32)
    title_para.font.bold = True

    # Body intro
    body_box = slide.shapes.add_textbox(
        Inches(0.75), Inches(1.4), Inches(11.8), Inches(0.6)
    )
    body_frame = body_box.text_frame
    body_frame.word_wrap = True
    body_para = body_frame.paragraphs[0]
    body_para.text = (
        "Please answer the following questions so we can generate the right deck."
    )
    body_para.font.size = Pt(14)

    # Content questions
    y_pos = 2.2
    if clarification_result.content_questions:
        y_pos = _add_question_section(
            slide,
            "Content Questions",
            clarification_result.content_questions,
            y_pos,
        )

    # Visualization questions
    if clarification_result.visualization_questions:
        y_pos = _add_question_section(
            slide,
            "Visualization Questions",
            clarification_result.visualization_questions,
            y_pos + 0.15,
        )

    prs.save(output_path)
    return output_path


def _add_question_section(slide, heading: str, questions, y_pos: float) -> float:
    """Add a categorized question block and return the next y position."""
    heading_box = slide.shapes.add_textbox(
        Inches(0.75), Inches(y_pos), Inches(11.8), Inches(0.4)
    )
    heading_frame = heading_box.text_frame
    heading_para = heading_frame.paragraphs[0]
    heading_para.text = heading
    heading_para.font.size = Pt(16)
    heading_para.font.bold = True

    y_pos += 0.4
    for question in questions:
        question_box = slide.shapes.add_textbox(
            Inches(0.95), Inches(y_pos), Inches(11.6), Inches(0.55)
        )
        question_frame = question_box.text_frame
        question_frame.word_wrap = True
        question_para = question_frame.paragraphs[0]
        question_para.text = f"• {question.question}"
        question_para.font.size = Pt(12)
        y_pos += 0.45

    return y_pos
