import logging
import os
import tempfile
from pathlib import Path
from typing import Any, Optional, Tuple
from uuid import uuid4

from pptx import Presentation
from pptx.util import Inches, Pt

from backend.layout_engine.layout_engine import generate_layout
from backend.llm.planner import create_operating_model_spec, create_slide_spec
from backend.modules.visual_planner import load_pattern_registry, plan_visual_pattern
from backend.services.clarification_renderer import render_clarification_deck
from backend.presentation_assets import asset_populator
from backend.presentation_assets import asset_registry
from backend.presentation_assets.fallbacks import resolve_fallback_asset_id
from backend.services.failure_renderer import (
    _SLIDE_HEIGHT_IN,
    _SLIDE_WIDTH_IN,
    render_placeholder_deck,
)
from ppt_renderer.operating_model_renderer import OperatingModelRenderer
from ppt_renderer.renderer import ProcessFlowRenderer
from schemas.executive_card import ExecutiveCardContent
from schemas.presentation import SlidePlan
from schemas.slide_spec import SlideSpec
from schemas.visual import VisualPatternSelection

logger = logging.getLogger(__name__)

_GENERATION_MODE_ENV = "EY_GENERATION_MODE"
_DEVELOPMENT_MODE = "development"
_PRODUCTION_MODE = "production"
_DEMO_MODE = "demo"

# Visual Pattern Library patterns allowed to use the production visual pipeline.
# Everything else falls back to a clean text slide.
_ALLOWED_VISUAL_PATTERNS = {
    "CL-01", "CL-02", "CL-03", "CL-04", "CL-05", "CL-06",
    "IG-01", "IG-02", "IG-03", "IG-04", "IG-05", "IG-06",
    "SECTION-DIVIDER",
}

# Infographic pattern IDs (for dispatch).
_INFOGRAPHIC_PATTERNS = {"IG-01", "IG-02", "IG-03", "IG-04", "IG-05", "IG-06"}

# Number of executive cards required by each enabled card-based visual pattern.
_CARD_COUNTS: dict[str, int] = {"CL-01": 4, "CL-02": 3, "CL-06": 3}

# Source priority for deriving executive cards when they are not already present
# in the SlideSpec. The helper walks each source in order, accumulates cards,
# and stops once the target count is reached; any shortfall is padded.
_CARD_SOURCES: dict[str, list[str]] = {
    "CL-01": ["cards", "business_benefits", "key_insights", "recommendations", "stages"],
    "CL-02": ["cards", "key_insights", "recommendations", "stages"],
    "CL-06": ["cards", "stages"],
}

# Number of KPI cards required by CL-03.
_KPI_CARD_COUNT = 3

# Candidate sources for deriving KPI content, in priority order.
_KPI_SOURCES = ["kpis", "metrics", "cards", "stages"]

# Candidate column key pairs for CL-04 / CL-05 two-column layouts.
# Each tuple is (left_key, right_key). The first pair ("columns", None) is
# handled specially because the data is already structured as columns.
_COLUMN_PAIRS: list[Tuple[str, Optional[str]]] = [
    ("columns", None),
    ("before", "after"),
    ("current", "future"),
    ("option_a", "option_b"),
    ("challenges", "solutions"),
    ("problems", "recommendations"),
    ("benefits", "actions"),
    ("left", "right"),
]


# ── Phase 1 (unchanged) ───────────────────────────────────────────────────────

def generate_slide(title, content):

    if _is_operating_model_request(title, content):
        logger.info("operating model generation selected")
        slide_spec = create_operating_model_spec(
            title,
            content
        )
        logger.info("LLM completed: operating_model stages=%s", len(slide_spec.get("stages", [])))

        renderer = OperatingModelRenderer()
    else:
        logger.info("process flow generation selected")
        slide_spec = create_slide_spec(
            title,
            content
        )
        logger.info("LLM completed: process_flow nodes=%s", len(slide_spec.get("nodes", [])))

        renderer = ProcessFlowRenderer()

    output_path = "generated_slide.pptx"

    renderer.render(
        slide_spec,
        output_path
    )
    logger.info("PPT generated: path=%s", output_path)

    return output_path


def _is_operating_model_request(title, content):
    text = f"{title} {content}".lower()

    operating_model_signals = [
        "operating model",
        "current state",
        "business stages",
        "detailed business activities",
        "executive summary",
        "value leakage",
        "kpis",
        "pain points, business risks",
        "risks, and inefficiencies",
    ]

    return any(signal in text for signal in operating_model_signals)


# ── Phase 2 ───────────────────────────────────────────────────────────────────

def _resolve_visual_selection(
    spec: SlideSpec, slide_plan: SlidePlan
) -> VisualPatternSelection:
    """
    Return the VisualPatternSelection to use for this slide.

    Honors a pattern_id carried on the SlideSpec (set once during content
    generation by the Deck Executor) so layout + content shaping share one
    source of truth. When no pattern is carried (e.g. legacy callers or specs
    built directly in tests), the Visual Planner re-scores as before.

    This ``try/except`` does NOT wrap this function — callers catch exceptions.
    """
    carried_id = getattr(spec, "visual_pattern_id", None)
    if carried_id:
        registry = load_pattern_registry()
        meta = registry.get(carried_id, {})
        carried_confidence = getattr(spec, "visual_confidence", None)
        return VisualPatternSelection(
            pattern_id=carried_id,
            category=meta.get("category", "creative_listing"),
            confidence=carried_confidence if carried_confidence is not None else 0.5,
            reasoning="Carried from content generation; single source of truth.",
            recommended_variant=None,
        )
    return plan_visual_pattern(slide_plan, spec)


def generate_slide_v2(title: str, content: str) -> str:
    """
    Phase 2 deck generation — orchestrator-driven pipeline.

    Routes the request through the full AI orchestration pipeline and returns
    a path to a ``.pptx`` file. The returned file is either:

    - A completed multi-slide consulting deck (``PipelineResult.status COMPLETED``).
    - A clarification placeholder deck listing questions that must be answered
      before generation can proceed (``PipelineResult.status WAITING_FOR_USER``).
    - A failure placeholder deck when the pipeline completed but produced no
      successful slides.

    This function always returns a file path so that the external ``/generate/v2``
    contract remains unchanged.

    Parameters
    ----------
    title:
        Raw slide title from the user request.
    content:
        Raw slide content / description from the user request.

    Returns
    -------
    str
        Path to the generated ``.pptx`` file.
    """
    from backend.orchestrator import run_pipeline  # local import avoids circular deps at module load
    from schemas.pipeline_result import PipelineResult

    logger.info("v2 generation started: title=%s", title)

    result = run_pipeline(title, content)
    mode = _generation_mode()
    output_path = _new_output_path("generated_slide_v2", mode=mode)

    if result.status == "WAITING_FOR_USER":
        logger.info(
            "v2 pipeline waiting for clarification: content=%d visualization=%d",
            len(result.clarification_result.content_questions) if result.clarification_result else 0,
            len(result.clarification_result.visualization_questions) if result.clarification_result else 0,
        )
        return render_clarification_deck(
            output_path,
            result.clarification_result,
        )

    deck_result = result.deck_execution_result
    if deck_result is None or not deck_result.successful_slides:
        logger.error(
            "v2 pipeline produced no successful slides: failed=%d",
            len(deck_result.failed_slides) if deck_result else 0,
        )
        errors = []
        if deck_result:
            errors = [slide.error for slide in deck_result.failed_slides if slide.error]
        return render_placeholder_deck(output_path, errors)

    if mode == _DEMO_MODE and not deck_result.evaluation_report.demo_ready:
        errors = _demo_failure_messages(deck_result)
        _persist_generation_metadata(output_path, deck_result)
        logger.error("v2 demo mode blocked non-demo-ready deck: %s", errors)
        return render_placeholder_deck(output_path, errors)

    prs = Presentation()
    prs.slide_width = Inches(_SLIDE_WIDTH_IN)
    prs.slide_height = Inches(_SLIDE_HEIGHT_IN)

    for slide_result in deck_result.slides:
        if not slide_result.success:
            continue

        spec = slide_result.slide_spec
        slide_plan = slide_result.slide_plan

        # Sprint F — Asset Populator path: when the slide carries a selected
        # Presentation Asset, copy and populate that asset slide directly.
        if spec.asset_id:
            slide_count_before = len(prs.slides)
            try:
                asset_populator.populate_asset_slide(prs, spec)
                _mark_population(slide_result, "success")
                logger.info(
                    "v2 asset slide populated: slide=%d asset=%s",
                    slide_plan.slide_number,
                    spec.asset_id,
                )
                continue
            except Exception as exc:  # noqa: BLE001
                _mark_population(slide_result, "asset_failed", warning=str(exc))
                logger.warning(
                    "v2 asset population failed for slide %d; mode=%s error=%s",
                    slide_plan.slide_number,
                    mode,
                    exc,
                )
                _remove_partial_slides(prs, slide_count_before)
                if mode == _PRODUCTION_MODE:
                    if _try_family_fallback(prs, spec, slide_result):
                        continue
                    logger.error(
                        "v2 production fallback failed for slide %d; skipping legacy renderer",
                        slide_plan.slide_number,
                    )
                    continue

        if mode == _PRODUCTION_MODE:
            if _try_family_fallback(prs, spec, slide_result):
                continue
            logger.error(
                "v2 production slide has no usable asset/fallback; skipping legacy renderer: slide=%d",
                slide_plan.slide_number,
            )
            continue

        # Development mode keeps the legacy renderer as an explicit debugging
        # fallback and records that downgrade in the evaluation report.
        _mark_population(slide_result, "legacy_fallback", fallback_used=True)
        renderer = _select_renderer(spec.slide_type)

        raw_spec = spec.raw_spec
        layout_spec = None
        try:
            visual_selection = _resolve_visual_selection(spec, slide_plan)
            logger.info(
                "v2 visual pattern selected: slide=%d pattern=%s confidence=%s source=%s",
                slide_plan.slide_number,
                visual_selection.pattern_id,
                visual_selection.confidence,
                "carried" if spec.visual_pattern_id else "re-scored",
            )
            if visual_selection.pattern_id in _ALLOWED_VISUAL_PATTERNS:
                # Build the pattern-native content first so we know the real item count,
                # then synthesize a layout that fits exactly (no empty rectangles,
                # no truncation).
                raw_spec = _build_pattern_content(
                    raw_spec, visual_selection.pattern_id
                )
                item_count = _count_content_items(raw_spec, visual_selection.pattern_id)
                layout_spec = generate_layout(visual_selection, item_count=item_count)
                logger.info(
                    "v2 layout selected: slide=%d layout=%s items=%d",
                    slide_plan.slide_number,
                    layout_spec.layout_id,
                    item_count,
                )
        except Exception as exc:  # noqa: BLE001
            logger.warning(
                "v2 visual pipeline failed for slide %d; falling back to legacy renderer: %s",
                slide_plan.slide_number,
                exc,
            )
            layout_spec = None

        slide_count_before = len(prs.slides)
        try:
            renderer.render(raw_spec, output_path, presentation=prs, layout_spec=layout_spec)
        except Exception as render_exc:  # noqa: BLE001
            logger.warning(
                "v2 render failed for slide %d; rendering clean text fallback slide: %s",
                slide_plan.slide_number,
                render_exc,
            )
            # The renderer commits a slide before rendering. Remove any partial
            # slide it left behind so a failure produces exactly one fallback slide.
            _remove_partial_slides(prs, slide_count_before)
            _render_text_fallback_slide(prs, raw_spec)
            _mark_population(slide_result, "text_fallback", fallback_used=True, warning=str(render_exc))
        logger.info(
            "v2 slide appended: slide_type=%s",
            spec.slide_type,
        )

    prs.save(output_path)
    _persist_generation_metadata(output_path, deck_result)

    rendered_count = len(prs.slides)
    expected = len(deck_result.successful_slides)
    logger.info(
        "v2 PPT generated: path=%s slides=%d expected=%d",
        output_path,
        rendered_count,
        expected,
    )
    if rendered_count != expected:
        logger.warning(
            "v2 slide count mismatch: rendered=%d expected=%d — extra slides "
            "may have been added by an external tool (e.g., PowerPoint).",
            rendered_count,
            expected,
        )
    return output_path


def _generation_mode() -> str:
    mode = os.getenv(_GENERATION_MODE_ENV, _DEVELOPMENT_MODE).strip().lower()
    if mode in {_PRODUCTION_MODE, _DEMO_MODE}:
        return mode
    return _DEVELOPMENT_MODE


def _new_output_path(prefix: str, *, mode: str) -> str:
    if mode == _DEVELOPMENT_MODE:
        return f"{prefix}.pptx"
    base = Path(tempfile.gettempdir()) / "ey-ai-slide-generator"
    request_dir = base / uuid4().hex
    request_dir.mkdir(parents=True, exist_ok=True)
    return str(request_dir / f"{prefix}.pptx")


def _persist_generation_metadata(output_path: str, deck_result) -> None:
    metadata_path = Path(output_path).with_suffix(".evaluation.json")
    try:
        metadata_path.write_text(
            deck_result.evaluation_report.model_dump_json(indent=2),
            encoding="utf-8",
        )
    except Exception as exc:  # noqa: BLE001
        logger.warning("could not persist generation metadata: %s", exc)


def _demo_failure_messages(deck_result) -> list[str]:
    report = deck_result.evaluation_report
    messages: list[str] = []
    if deck_result.failed_slides:
        messages.extend(slide.error or f"Slide {slide.slide_plan.slide_number} failed" for slide in deck_result.failed_slides)
    if report.duplicate_roles:
        messages.append(f"Duplicate roles: {', '.join(report.duplicate_roles)}")
    if report.repeated_slide_titles:
        messages.append(f"Repeated titles: {', '.join(report.repeated_slide_titles)}")
    if report.consulting_language_warnings:
        messages.append(f"Consulting language warnings: {len(report.consulting_language_warnings)}")
    if report.placeholder_leakage:
        messages.append(f"Placeholder leakage: {len(report.placeholder_leakage)}")
    if report.overflow_slides:
        messages.append(f"Text-fit failures on slides: {report.overflow_slides}")
    return messages or ["Deck did not meet demo readiness gates."]


def _mark_population(slide_result, population: str, *, fallback_used: bool = False, warning: Optional[str] = None) -> None:
    report = getattr(slide_result, "evaluation_report", None)
    if report is None:
        return
    report.population = population
    if fallback_used:
        report.fallback_used = True
    if warning:
        report.warnings.append(warning)


def _try_family_fallback(prs: Presentation, spec: SlideSpec, slide_result) -> bool:
    family = _family_for_spec(spec)
    fallback_id = resolve_fallback_asset_id(family, require_certified=True)
    report = getattr(slide_result, "evaluation_report", None)
    if fallback_id is None:
        if report is not None:
            report.warnings.append(f"no certified fallback asset for family {family!r}")
        return False
    fallback_spec = spec.model_copy(update={"asset_id": fallback_id})
    try:
        asset_populator.populate_asset_slide(prs, fallback_spec)
    except Exception as exc:  # noqa: BLE001
        if report is not None:
            report.warnings.append(f"certified fallback {fallback_id} failed: {exc}")
        return False
    if report is not None:
        report.population = "family_fallback"
        report.fallback_used = True
        report.fallback_asset = fallback_id
    return True


def _family_for_spec(spec: SlideSpec) -> Optional[str]:
    if spec.asset_id:
        manifest = asset_registry.get(spec.asset_id)
        if manifest is not None:
            return manifest.family
    pattern_id = getattr(spec, "visual_pattern_id", None)
    if not pattern_id:
        return None
    try:
        from backend.presentation_assets.asset_selector import family_for_pattern

        return family_for_pattern(pattern_id)
    except Exception:
        return None


def _count_content_items(raw_spec: dict[str, Any], pattern_id: str) -> int:
    """Return the number of content items available for ``pattern_id``."""
    if pattern_id in _CARD_COUNTS:
        return len(raw_spec.get("cards", []))
    if pattern_id == "CL-03":
        return len(raw_spec.get("kpis", []))
    if pattern_id in ("CL-04", "CL-05"):
        # Two-column layouts are fixed at two columns; item count is the
        # maximum number of rows in either column.
        left, right, _, _ = _derive_two_columns(raw_spec)
        return max(len(left), len(right))
    if pattern_id == "IG-01":
        return len(raw_spec.get("events", []))
    if pattern_id == "IG-02":
        return len(raw_spec.get("phases", []))
    if pattern_id == "IG-03":
        return len(raw_spec.get("steps", []))
    if pattern_id == "IG-04":
        return len(raw_spec.get("cells", []))
    if pattern_id == "IG-05":
        return len(raw_spec.get("journey_stages", []))
    if pattern_id == "IG-06":
        return len(raw_spec.get("domains", []))
    if pattern_id == "SECTION-DIVIDER":
        return 1
    return 0


def _build_pattern_content(
    raw_spec: dict[str, Any],
    pattern_id: str,
) -> dict[str, Any]:
    """
    Build a renderer-ready content dict for a production-enabled visual pattern.

    The original ``raw_spec`` is never mutated. Padding/truncation is avoided;
    the layout engine synthesizes component positions for the actual item count.
    """
    content = dict(raw_spec)

    if pattern_id in _CARD_COUNTS:
        cards = _build_executive_card_content(raw_spec, pattern_id)
        content["cards"] = [card.model_dump(mode="json") for card in cards]
    elif pattern_id == "CL-03":
        content["kpis"] = _build_kpi_content(raw_spec)
    elif pattern_id in ("CL-04", "CL-05"):
        content.update(_build_two_column_content(raw_spec))
    elif pattern_id in _INFOGRAPHIC_PATTERNS:
        content.update(_build_infographic_content(raw_spec, pattern_id))

    return content


def _build_executive_card_content(
    raw_spec: dict[str, Any],
    pattern_id: str,
) -> list[ExecutiveCardContent]:
    """
    Derive executive cards for the requested card-based visual pattern.

    Returns a list of ``ExecutiveCardContent`` objects suitable for the
    Executive Insight Card renderer. The target count and derivation priority
    depend on ``pattern_id``:

    - CL-01 (Four Insight Cards): 4 cards.
      Priority: existing cards → business benefits → key insights →
      recommendations → stages → title/subtitle placeholders.
    - CL-02 (Three Cards): 3 cards.
      Priority: existing cards → key insights → recommendations → stages →
      title/subtitle placeholders.
    - CL-06 (Executive Summary Cards): 3 cards.
      Priority: existing cards → stages → title/subtitle placeholders.

    The original ``raw_spec`` is never mutated.
    """
    target_count = _CARD_COUNTS.get(pattern_id, 3)
    sources = _CARD_SOURCES.get(pattern_id, _CARD_SOURCES["CL-06"])
    cards: list[ExecutiveCardContent] = []

    for source in sources:
        items = _extract_card_source_items(raw_spec, source)
        for item in items:
            cards.append(_normalize_to_executive_card(item, source))
            if len(cards) >= target_count:
                break
        if len(cards) >= target_count:
            break

    # Do not pad short lists; the layout engine synthesizes the component
    # count to match the actual content.
    return cards[:target_count]


def _build_kpi_content(raw_spec: dict[str, Any]) -> list[dict[str, str]]:
    """
    Derive KPI cards for CL-03.

    Priority: existing ``kpis`` → ``metrics`` → ``cards`` → ``stages`` →
    title placeholders.
    """
    kpis: list[dict[str, str]] = []

    for source in _KPI_SOURCES:
        items = _extract_card_source_items(raw_spec, source)
        for item in items:
            kpis.append(_normalize_to_kpi(item))
            if len(kpis) >= _KPI_CARD_COUNT:
                break
        if len(kpis) >= _KPI_CARD_COUNT:
            break

    # Do not pad short lists; the layout engine adapts to the actual count.
    return kpis[:_KPI_CARD_COUNT]


def _build_two_column_content(raw_spec: dict[str, Any]) -> dict[str, Any]:
    """
    Derive left/right column content for CL-04 (Comparison) and CL-05
    (Two Column).

    Tries structured ``columns`` first, then common key pairs such as
    ``before``/``after``, ``current``/``future``, ``option_a``/``option_b``,
    ``challenges``/``solutions``, ``problems``/``recommendations``,
    ``benefits``/``actions``, and ``left``/``right``.
    """
    left_label, right_label, left_items, right_items = _derive_two_columns(raw_spec)
    return {
        "columns": [{"label": left_label}, {"label": right_label}],
        "left_items": left_items,
        "right_items": right_items,
    }


def _derive_two_columns(
    raw_spec: dict[str, Any],
) -> tuple[str, str, list[Any], list[Any]]:
    """Return (left_label, right_label, left_items, right_items) for a slide."""
    if "columns" in raw_spec:
        cols = raw_spec["columns"]
        if isinstance(cols, list) and len(cols) >= 2:
            return (
                _column_label(cols[0]),
                _column_label(cols[1]),
                _column_items(cols[0]),
                _column_items(cols[1]),
            )

    for left_key, right_key in _COLUMN_PAIRS[1:]:
        if left_key in raw_spec and right_key in raw_spec:
            return (
                _key_to_label(left_key),
                _key_to_label(right_key),
                _normalize_items(raw_spec[left_key]),
                _normalize_items(raw_spec[right_key]),
            )

    title = str(raw_spec.get("title", ""))
    return "Current", "Future", ([title] if title else []), ([title] if title else [])


def _column_label(column: Any) -> str:
    if isinstance(column, dict):
        return str(column.get("label") or column.get("title") or "")
    return str(column)


def _column_items(column: Any) -> list[Any]:
    if isinstance(column, dict):
        return _normalize_items(column.get("items", []))
    return []


def _normalize_items(items: Any) -> list[Any]:
    """Normalize a list of items so both text and bar renderers can consume them."""
    if not isinstance(items, list):
        return [str(items)] if items is not None else []

    result: list[Any] = []
    for item in items:
        if isinstance(item, dict):
            text = (
                item.get("text")
                or item.get("title")
                or item.get("name")
                or item.get("label")
                or ""
            )
            name = (
                item.get("name")
                or item.get("title")
                or item.get("label")
                or text
            )
            result.append({"text": str(text), "name": str(name)})
        else:
            result.append(str(item))
    return result


def _key_to_label(key: str) -> str:
    """Convert a snake_case key into a human-readable column label."""
    return key.replace("_", " ").title()


def _extract_card_source_items(raw_spec: dict[str, Any], source: str) -> list[Any]:
    """Return candidate items from ``raw_spec`` for a named card source."""
    if source == "cards":
        return raw_spec.get("cards", [])
    if source == "business_benefits":
        return raw_spec.get("business_benefits") or raw_spec.get("benefits") or []
    if source == "key_insights":
        return raw_spec.get("key_insights") or raw_spec.get("insights") or []
    if source == "recommendations":
        return raw_spec.get("recommendations", [])
    if source == "stages":
        return raw_spec.get("stages", [])
    if source == "kpis":
        return raw_spec.get("kpis", [])
    if source == "metrics":
        return raw_spec.get("metrics", [])
    return []


def _normalize_to_executive_card(item: Any, source: str) -> ExecutiveCardContent:
    """Normalize a source item into an ``ExecutiveCardContent`` instance."""
    if isinstance(item, ExecutiveCardContent):
        return item

    if isinstance(item, dict):
        title = item.get("title") or item.get("label") or item.get("name") or ""
        description = item.get("description", "")
        if source == "stages" and not description:
            activities = item.get("activities", [])
            description = ", ".join(str(activity) for activity in activities)
        return ExecutiveCardContent(title=str(title), description=str(description))

    return ExecutiveCardContent(title=str(item), description="")


def _normalize_to_kpi(item: Any) -> dict[str, str]:
    """Normalize a source item into a KPI card dict with ``label`` and ``value``."""
    if isinstance(item, dict):
        label = (
            item.get("label")
            or item.get("title")
            or item.get("name")
            or item.get("metric")
            or ""
        )
        value = (
            item.get("value")
            or item.get("target")
            or item.get("trend")
            or item.get("description")
            or ""
        )
        return {"label": str(label), "value": str(value)}
    return {"label": str(item), "value": ""}


# ── Infographic content builders ──────────────────────────────────────────────

# Required component counts per infographic layout.
_INFOGRAPHIC_COUNTS: dict[str, int] = {
    "IG-01": 4,   # 4 timeline event nodes
    "IG-02": 4,   # 4 roadmap phase bars
    "IG-03": 7,   # up to 7 process-flow step nodes
    "IG-04": 9,   # 3x3 matrix cells
    "IG-05": 4,   # 4 journey stages
    "IG-06": 4,   # 4 capability domains
}


def _build_infographic_content(
    raw_spec: dict[str, Any],
    pattern_id: str,
) -> dict[str, Any]:
    """
    Build a renderer-ready content dict for a production-enabled infographic
    visual pattern. Coordinates always come from the Layout Engine JSON.
    """
    if pattern_id == "IG-01":
        return {"events": _build_events(raw_spec)}
    if pattern_id == "IG-02":
        return {"phases": _build_phases(raw_spec)}
    if pattern_id == "IG-03":
        return {"steps": _build_steps(raw_spec)}
    if pattern_id == "IG-04":
        return {"cells": _build_cells(raw_spec)}
    if pattern_id == "IG-05":
        return {"stages": _build_journey_stages(raw_spec)}
    if pattern_id == "IG-06":
        return {"domains": _build_domains(raw_spec)}
    return {}


def _build_events(raw_spec: dict[str, Any]) -> list[dict[str, str]]:
    """IG-01: derive timeline events from structured content (no padding)."""
    sources = [
        raw_spec.get("events"),
        raw_spec.get("timeline"),
        raw_spec.get("milestones"),
        raw_spec.get("phases"),
        raw_spec.get("roadmap"),
        raw_spec.get("stages"),
    ]
    for source in sources:
        events = _normalize_nodes(source)
        if events:
            return events
    return []


def _build_phases(raw_spec: dict[str, Any]) -> list[dict[str, str]]:
    """IG-02: derive roadmap phases from structured content (no padding)."""
    sources = [
        raw_spec.get("phases"),
        raw_spec.get("roadmap"),
        raw_spec.get("stages"),
    ]
    for source in sources:
        phases = _normalize_bars(source)
        if phases:
            return phases
    return []


def _build_steps(raw_spec: dict[str, Any]) -> list[dict[str, str]]:
    """IG-03: derive process-flow steps from structured content (no padding)."""
    sources = [
        raw_spec.get("steps"),
        raw_spec.get("nodes"),
        raw_spec.get("process"),
        raw_spec.get("stages"),
    ]
    for source in sources:
        steps = _normalize_nodes(source)
        if steps:
            return steps
    return []


def _build_cells(raw_spec: dict[str, Any]) -> list[dict[str, str]]:
    """IG-04: derive matrix cells from rows, capabilities, or priorities (no padding)."""
    rows = raw_spec.get("rows")
    if isinstance(rows, list) and rows:
        cells: list[dict[str, str]] = []
        for row in rows:
            row_cells = row.get("cells") if isinstance(row, dict) else None
            if isinstance(row_cells, list):
                cells.extend(_normalize_cells(row_cells))
            elif isinstance(row, dict):
                cells.append(_normalize_cell(row))
            else:
                cells.append({"value": str(row)})
        if cells:
            return cells

    sources = [
        raw_spec.get("cells"),
        raw_spec.get("capabilities"),
        raw_spec.get("priorities"),
        raw_spec.get("matrix"),
    ]
    for source in sources:
        cells = _normalize_cells(source)
        if cells:
            return cells

    return []


def _build_journey_stages(raw_spec: dict[str, Any]) -> list[dict[str, str]]:
    """IG-05: derive journey stages from structured content (no padding)."""
    sources = [
        raw_spec.get("journey_stages"),
        raw_spec.get("stages"),
        raw_spec.get("touchpoints"),
    ]
    for source in sources:
        stages = _normalize_nodes(source)
        if stages:
            return stages
    return []


def _build_domains(raw_spec: dict[str, Any]) -> list[dict[str, Any]]:
    """IG-06: derive capability-map domains from structured content (no padding)."""
    if "domains" in raw_spec:
        domains = _normalize_domains(raw_spec["domains"])
        if domains:
            return domains

    sources = [
        raw_spec.get("capabilities"),
        raw_spec.get("functions"),
        raw_spec.get("business_areas"),
        raw_spec.get("stages"),
    ]
    for source in sources:
        domains = _normalize_domains(source)
        if domains:
            return domains

    return []


# ── Normalization helpers ─────────────────────────────────────────────────────


def _normalize_nodes(items: Any) -> list[dict[str, str]]:
    """Normalize items into node-shaped dicts (label/title/name)."""
    if not isinstance(items, list):
        return []
    result: list[dict[str, str]] = []
    for item in items:
        if isinstance(item, dict):
            label = (
                item.get("label")
                or item.get("title")
                or item.get("name")
                or item.get("date")
                or ""
            )
            description = item.get("description", "")
            result.append({"label": str(label), "description": str(description)})
        elif item:
            result.append({"label": str(item), "description": ""})
    return result


def _normalize_bars(items: Any) -> list[dict[str, str]]:
    """Normalize items into bar-shaped dicts (name + activities)."""
    if not isinstance(items, list):
        return []
    result: list[dict[str, str]] = []
    for item in items:
        if isinstance(item, dict):
            name = (
                item.get("name")
                or item.get("label")
                or item.get("title")
                or ""
            )
            activities = item.get("activities") or item.get("deliverables") or []
            if isinstance(activities, list):
                description = ", ".join(str(a) for a in activities)
            else:
                description = str(activities)
            result.append({"name": str(name), "description": description})
        elif item:
            result.append({"name": str(item), "description": ""})
    return result


def _normalize_cells(items: Any) -> list[dict[str, str]]:
    """Normalize items into matrix cell dicts (value)."""
    if not isinstance(items, list):
        return []
    result: list[dict[str, str]] = []
    for item in items:
        if isinstance(item, dict):
            value = (
                item.get("value")
                or item.get("name")
                or item.get("label")
                or item.get("title")
                or ""
            )
            result.append({"value": str(value)})
        elif item:
            result.append({"value": str(item)})
    return result


def _normalize_cell(item: Any) -> dict[str, str]:
    """Normalize a single matrix cell."""
    if isinstance(item, dict):
        value = (
            item.get("value")
            or item.get("name")
            or item.get("label")
            or item.get("title")
            or ""
        )
        return {"value": str(value)}
    return {"value": str(item)}


def _normalize_domains(items: Any) -> list[dict[str, Any]]:
    """Normalize items into capability-map domain dicts (name + capabilities)."""
    if not isinstance(items, list):
        return []
    result: list[dict[str, Any]] = []
    for item in items:
        if isinstance(item, dict):
            name = (
                item.get("name")
                or item.get("domain")
                or item.get("function")
                or item.get("area")
                or item.get("label")
                or item.get("title")
                or ""
            )
            capabilities = item.get("capabilities", [])
            if isinstance(capabilities, list):
                caps = [
                    {"name": str(c.get("name", c)) if isinstance(c, dict) else str(c)}
                    for c in capabilities
                ]
            else:
                caps = [{"name": str(capabilities)}]
            result.append({"name": str(name), "capabilities": caps})
        elif item:
            result.append({"name": str(item), "capabilities": []})
    return result


def _select_renderer(slide_type: str):
    """
    Select the appropriate Phase 1 renderer for the given slide type.

    This is the ONLY place in Phase 2 code that references renderer classes.
    All renderer logic remains inside ``ppt_renderer/`` untouched.

    Parameters
    ----------
    slide_type:
        Normalised slide type from ``SlideSpec.slide_type``.

    Returns
    -------
    Renderer instance with a ``.render(spec_dict, output_path)`` method.
    """
    if slide_type == "operating_model":
        return OperatingModelRenderer()
    # Default: ProcessFlowRenderer handles process_flow, comparison,
    # current_future, and unknown until dedicated renderers are built.
    return ProcessFlowRenderer()


def _remove_partial_slides(prs: Presentation, target_count: int) -> None:
    """Remove slides added after ``target_count`` (e.g. a failed render)."""
    while len(prs.slides) > target_count:
        slide = prs.slides[-1]
        slide._element.getparent().remove(slide._element)


def _render_text_fallback_slide(prs: Presentation, raw_spec: dict[str, Any]) -> None:
    """Draw a clean title + subtitle + bullet list when the renderer fails.

    Used as a last-resort fallback so a failed visual pipeline still produces a
    readable slide rather than aborting the whole deck. Up to six bullet lines
    are derived from the executive summary (split into sentences) and the
    description.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    title = str(raw_spec.get("title", "Slide"))
    subtitle = str(raw_spec.get("subtitle", "") or raw_spec.get("description", ""))

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.4), Inches(12.3), Inches(0.9)
    )
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(28)
    title_para.font.bold = True

    # Subtitle
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.6)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.word_wrap = True
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.text = subtitle
        subtitle_para.font.size = Pt(14)

    # Bullets from executive_summary / description
    bullets = _extract_fallback_bullets(raw_spec)
    if bullets:
        body_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(2.4), Inches(11.9), Inches(4.5)
        )
        body_frame = body_box.text_frame
        body_frame.word_wrap = True
        for index, bullet in enumerate(bullets):
            para = body_frame.paragraphs[0] if index == 0 else body_frame.add_paragraph()
            para.text = f"• {bullet}"
            para.font.size = Pt(14)
            if index > 0:
                para.space_before = Pt(8)


def _extract_fallback_bullets(raw_spec: dict[str, Any]) -> list[str]:
    """Return up to six concise bullet lines for the text fallback slide."""
    lines: list[str] = []
    summary = raw_spec.get("executive_summary") or raw_spec.get("description") or ""
    if isinstance(summary, str) and summary.strip():
        import re

        parts = re.split(r"(?<=[.!?])\s+", summary.strip())
        for part in parts:
            part = part.strip()
            if part:
                lines.append(part)
    return lines[:6]
