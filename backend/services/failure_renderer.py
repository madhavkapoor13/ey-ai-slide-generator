"""
backend/services/failure_renderer.py
=====================================
Failure rendering helper.

When the Deck Executor produces no successful slides, this module generates
a minimal one-slide placeholder deck so the frontend still receives a valid
.pptx file. It is intentionally separate from ``slide_service.py`` so that
failure presentation remains a distinct rendering concern.
"""

from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches, Pt

# Widescreen 16:9 dimensions shared across renderers.
_SLIDE_WIDTH_IN = 13.333
_SLIDE_HEIGHT_IN = 7.5


def render_placeholder_deck(output_path: str, errors: list[str]) -> str:
    """
    Render a single-slide placeholder deck explaining that generation failed.

    Parameters
    ----------
    output_path:
        File path where the placeholder .pptx will be saved.
    errors:
        Per-slide error messages collected from the failed deck execution.

    Returns
    -------
    str
        The ``output_path`` that was written.
    """
    prs = Presentation()
    prs.slide_width = Inches(_SLIDE_WIDTH_IN)
    prs.slide_height = Inches(_SLIDE_HEIGHT_IN)

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.75), Inches(2.0), Inches(11.8), Inches(1.0)
    )
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_para = title_frame.paragraphs[0]
    title_para.text = "Generation Failed"
    title_para.font.size = Pt(32)
    title_para.font.bold = True

    # Body
    body_box = slide.shapes.add_textbox(
        Inches(0.75), Inches(3.2), Inches(11.8), Inches(2.5)
    )
    body_frame = body_box.text_frame
    body_frame.word_wrap = True
    body_para = body_frame.paragraphs[0]
    body_para.text = (
        "We were unable to generate the requested slides. "
        "Please try again with a more detailed prompt."
    )
    body_para.font.size = Pt(16)

    if errors:
        detail_para = body_frame.add_paragraph()
        detail_para.text = f"Details: {'; '.join(errors[:5])}"
        detail_para.font.size = Pt(12)
        detail_para.space_before = Pt(12)

    prs.save(output_path)
    return output_path
